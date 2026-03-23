"""
HPP Product Framework — Plain Text Training Deck
Audience: Internal Hyundai Reps → Dealer Training
Font: Aptos throughout | No shapes or decorative formatting
Output: /Users/justin/Downloads/HPP Product Framework - Training Deck.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches as I

# ── Layout constants ────────────────────────────────────────────────────────
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
MARGIN_L = Inches(0.6)
MARGIN_T = Inches(0.5)
MARGIN_R = Inches(0.6)
CONTENT_W = SLIDE_W - MARGIN_L - MARGIN_R

F = "Aptos"

C_BLACK   = RGBColor(0x1A, 0x1A, 0x1A)
C_NAVY    = RGBColor(0x00, 0x20, 0x60)
C_MID     = RGBColor(0x00, 0x4C, 0x97)
C_TEAL    = RGBColor(0x00, 0x78, 0x8C)
C_PURPLE  = RGBColor(0x6A, 0x1A, 0x4D)
C_ORANGE  = RGBColor(0xE8, 0x6F, 0x1E)
C_GREY    = RGBColor(0x55, 0x55, 0x55)
C_WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
C_CORE    = RGBColor(0xC8, 0xA0, 0x32)
C_NONCORE = RGBColor(0x6D, 0x6D, 0x6D)

def new_prs():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs

def blank_slide(prs):
    layout = prs.slide_layouts[6]  # completely blank
    return prs.slides.add_slide(layout)

def txb(slide, text, left, top, width, height,
        size=11, bold=False, color=C_BLACK, align=PP_ALIGN.LEFT,
        italic=False, word_wrap=True):
    """Add a simple text box with a single run."""
    box = slide.shapes.add_textbox(I(left), I(top), I(width), I(height))
    tf = box.text_frame
    tf.word_wrap = word_wrap
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = F
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    return box

def txb_multi(slide, paras, left, top, width, height, word_wrap=True):
    """
    paras: list of dicts:
      { text, size, bold, color, align, italic, space_before, bullet }
    """
    box = slide.shapes.add_textbox(I(left), I(top), I(width), I(height))
    tf = box.text_frame
    tf.word_wrap = word_wrap
    first = True
    for pd in paras:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.alignment = pd.get("align", PP_ALIGN.LEFT)
        if pd.get("space_before"):
            p.space_before = Pt(pd["space_before"])
        if pd.get("space_after"):
            p.space_after = Pt(pd["space_after"])
        text = pd.get("text", "")
        if pd.get("bullet"):
            text = "  •  " + text
        r = p.add_run()
        r.text = text
        r.font.name = F
        r.font.size = Pt(pd.get("size", 11))
        r.font.bold = pd.get("bold", False)
        r.font.italic = pd.get("italic", False)
        r.font.color.rgb = pd.get("color", C_BLACK)
    return box

def hline(paras, text, size=11, bold=False, color=C_BLACK,
          space_before=4, space_after=0, bullet=False, italic=False):
    """Append a paragraph definition to paras list."""
    paras.append({
        "text": text, "size": size, "bold": bold, "color": color,
        "space_before": space_before, "space_after": space_after,
        "bullet": bullet, "italic": italic,
        "align": PP_ALIGN.LEFT,
    })

# ── Slide builders ──────────────────────────────────────────────────────────

def build_cover(prs):
    slide = blank_slide(prs)
    # Main title
    txb(slide, "Hyundai Protection Plan", 0.6, 2.2, 12.1, 1.0,
        size=40, bold=True, color=C_NAVY, align=PP_ALIGN.CENTER)
    txb(slide, "Product Framework", 0.6, 3.1, 12.1, 0.7,
        size=32, bold=False, color=C_NAVY, align=PP_ALIGN.CENTER)
    txb(slide, "Internal Training Reference  |  For Hyundai Reps & Dealer Partners",
        0.6, 4.0, 12.1, 0.5, size=13, bold=False, color=C_GREY, align=PP_ALIGN.CENTER)
    txb(slide, "CORE Products = HPP Branded  |  NON-CORE Products = Off-Make / White-Label",
        0.6, 4.65, 12.1, 0.4, size=11, bold=False, color=C_GREY, align=PP_ALIGN.CENTER,
        italic=True)
    txb(slide, "Confidential — For Internal Use Only",
        0.6, 6.8, 12.1, 0.4, size=9, bold=False, color=C_GREY, align=PP_ALIGN.CENTER,
        italic=True)

def build_toc(prs):
    slide = blank_slide(prs)
    txb(slide, "Table of Contents", 0.6, 0.3, 12.1, 0.6,
        size=22, bold=True, color=C_NAVY)
    txb(slide, "CORE = HPP Branded Product    NON-CORE = Off-Make / White-Label Product",
        0.6, 0.85, 12.1, 0.35, size=10, italic=True, color=C_GREY)

    sections = [
        ("MECHANICAL COVERAGE", C_NAVY, [
            ("Vehicle Service Protection (ICE)", "HFVI", True),
            ("Vehicle Service Protection — Competitive Makes (ICE)", "WFVI", False),
            ("Vehicle Service Protection Original Owner (ICE)", "HFOI", True),
            ("Certified Used Vehicle (CUV) Wrap", "HFCI", True),
            ("High Mileage Vehicle Service Protection (ICE)", "WFHM", False),
            ("EV Care Vehicle Service Protection", "HFVE", True),
            ("EV Care VSP — Competitive Makes", "WFVE", False),
            ("EV Care VSP — Original Owner", "HFOE", True),
            ("EV Care Certified Used Vehicle (CUV) Wrap", "HFCE", True),
            ("VSP Livery Coverage (ICE & EV)", "HCVL", True),
            ("Wear Protection (ICE & EV)", "HFWP", True),
        ]),
        ("MAINTENANCE PLANS", C_MID, [
            ("Maintenance Basic Wrap (ICE)", "HFBI", True),
            ("Maintenance Scheduled Wrap (ICE)", "HFSI", True),
            ("Pre-Paid Maintenance (ICE)", "HFPI", True),
            ("Pre-Paid Maintenance — Competitive Makes (ICE)", "WFPI", False),
            ("EV Care Maintenance", "HFEM", True),
        ]),
        ("DAMAGE CARE", C_TEAL, [
            ("Lease-End Protection", "HFEL", True),
            ("Multi-Coverage Protection (Tire & Wheel | Dent | Windshield | Key)", "HFMC", True),
            ("Appearance Protection", "HFAP", False),
        ]),
        ("LOSS & THEFT PROTECTION", C_PURPLE, [
            ("Guaranteed Asset Protection (GAP)", "HFGP / WFGP", True),
            ("Guaranteed Asset Protection Plus (GAP Plus)", "HFGS / WFGS", True),
            ("Theft Protection", "HFTP", False),
        ]),
    ]

    col_configs = [
        {"x": 0.6,  "sections": ["MECHANICAL COVERAGE"]},
        {"x": 7.0,  "sections": ["MAINTENANCE PLANS", "DAMAGE CARE", "LOSS & THEFT PROTECTION"]},
    ]

    col_data = {s[0]: s for s in sections}

    for col in col_configs:
        x = col["x"]
        y = 1.35
        for sec_name in col["sections"]:
            sec_title, sec_color, items = col_data[sec_name]
            paras = []
            hline(paras, sec_title, size=10, bold=True, color=sec_color, space_before=6)
            for (item_name, code, is_core) in items:
                core_tag = "[CORE]" if is_core else "[NON-CORE]"
                hline(paras, f"{item_name}  |  {code}  {core_tag}",
                      size=9.5, bold=False, color=C_BLACK, space_before=2, bullet=True)
            txb_multi(slide, paras, x, y, 6.1, 6.0)
            y += 0.28 + len(items) * 0.23 + 0.15

def build_section_divider(prs, title, subtitle, color):
    slide = blank_slide(prs)
    txb(slide, "HYUNDAI PROTECTION PLAN", 0.6, 2.5, 12.1, 0.5,
        size=13, bold=False, color=C_GREY, align=PP_ALIGN.CENTER)
    txb(slide, title, 0.6, 3.1, 12.1, 1.0,
        size=36, bold=True, color=color, align=PP_ALIGN.CENTER)
    txb(slide, subtitle, 0.6, 4.15, 12.1, 0.5,
        size=13, bold=False, color=C_GREY, align=PP_ALIGN.CENTER, italic=True)

def build_product_slide(prs, d):
    slide = blank_slide(prs)

    name   = d["name"]
    code   = d["code"]
    core   = d.get("core", True)
    ev     = d.get("ev_care", False)

    core_label = "CORE  |  HPP Branded" if core else "NON-CORE  |  Off-Make / White-Label"
    core_color = C_CORE if core else C_NONCORE
    ev_tag     = "  |  EV CARE" if ev else ""

    # ── Header line ──────────────────────────────────────────────────────────
    txb(slide, name, 0.6, 0.22, 10.5, 0.55,
        size=22, bold=True, color=C_NAVY)
    txb(slide, f"{code}{ev_tag}", 0.6, 0.72, 5.5, 0.32,
        size=10, bold=True, color=core_color)
    txb(slide, core_label, 6.5, 0.72, 6.2, 0.32,
        size=10, bold=False, color=core_color, align=PP_ALIGN.RIGHT)

    # Horizontal rule (thin text line as separator)
    txb(slide, "─" * 120, 0.6, 1.0, 12.1, 0.2, size=6, color=RGBColor(0xCC,0xCC,0xCC))

    # ── Overview ─────────────────────────────────────────────────────────────
    overview = d.get("overview", "")
    if overview:
        paras = []
        hline(paras, "PRODUCT OVERVIEW", size=9, bold=True, color=C_GREY, space_before=0)
        hline(paras, overview, size=10.5, bold=False, color=C_BLACK, space_before=2, italic=True)
        txb_multi(slide, paras, 0.6, 1.1, 12.1, 0.7)

    # ── Three-column layout ──────────────────────────────────────────────────
    # Col 1: Plans + Term/Coverage     x=0.6   w=4.3
    # Col 2: Eligibility + Reimb + Add x=5.1   w=4.0
    # Col 3: Exclusions + Transfer + Cancel + Circle  x=9.3  w=3.8

    C1X, C1W = 0.60, 4.30
    C2X, C2W = 5.10, 4.00
    C3X, C3W = 9.30, 3.83
    ROW_Y = 1.78

    # ── COL 1: Plans ─────────────────────────────────────────────────────────
    paras = []
    hline(paras, "AVAILABLE PLANS", size=9.5, bold=True, color=C_NAVY, space_before=0)
    for pl in d.get("plans", []):
        pl_name = pl["name"]
        if pl.get("sub"):
            pl_name += f"  {pl['sub']}"
        hline(paras, pl_name, size=10, bold=True, color=C_BLACK, space_before=5, bullet=True)
        if pl.get("desc"):
            hline(paras, pl["desc"], size=9, bold=False, color=C_GREY, space_before=1)

    hline(paras, "", size=4, space_before=4)
    hline(paras, "TERM & COVERAGE", size=9.5, bold=True, color=C_NAVY, space_before=2)
    for line in d.get("term_coverage", []):
        hline(paras, line, size=9.5, bold=False, color=C_BLACK, space_before=2, bullet=True)

    txb_multi(slide, paras, C1X, ROW_Y, C1W, 5.5)

    # ── COL 2: Eligibility + Reimbursement + Additional Benefits ─────────────
    paras = []
    elig_label = d.get("elig_label", "ELIGIBILITY")
    hline(paras, elig_label.upper(), size=9.5, bold=True, color=C_NAVY, space_before=0)
    for line in d.get("eligibility", []):
        hline(paras, line, size=9.5, bold=False, color=C_BLACK, space_before=2, bullet=True)

    reimb = d.get("reimbursement", [])
    if reimb:
        hline(paras, "", size=4, space_before=4)
        hline(paras, "CLAIM REIMBURSEMENT", size=9.5, bold=True, color=C_NAVY, space_before=2)
        for row in reimb:
            # row is (label, [(text, bold), ...])
            label, parts = row
            full_text = label + " ".join(t for t, _ in parts)
            hline(paras, full_text.strip(), size=9.5, bold=False, color=C_BLACK,
                  space_before=2, bullet=True)

    add_ben = d.get("additional_benefits", [])
    if add_ben:
        hline(paras, "", size=4, space_before=4)
        hline(paras, "ADDITIONAL BENEFITS", size=9.5, bold=True, color=C_NAVY, space_before=2)
        for row in add_ben:
            full_text = " ".join(t for t, _ in row)
            hline(paras, full_text.strip(), size=9.5, bold=False, color=C_BLACK,
                  space_before=2, bullet=True)

    txb_multi(slide, paras, C2X, ROW_Y, C2W, 5.5)

    # ── COL 3: Exclusions + Transfer + Cancellation + Circle ─────────────────
    paras = []
    excl = d.get("exclusions", [])
    if excl:
        hline(paras, "EXCLUSIONS", size=9.5, bold=True, color=C_NAVY, space_before=0)
        for line in excl:
            hline(paras, line, size=9, bold=False, color=C_BLACK, space_before=2, bullet=True)

    transfer = d.get("transfer", "")
    if transfer:
        hline(paras, "", size=4, space_before=4)
        hline(paras, "TRANSFER", size=9.5, bold=True, color=C_NAVY, space_before=2)
        hline(paras, transfer, size=9.5, bold=False, color=C_BLACK, space_before=2, bullet=True)

    cancel = d.get("cancel", [])
    if cancel:
        hline(paras, "", size=4, space_before=4)
        hline(paras, "CANCELLATION", size=9.5, bold=True, color=C_NAVY, space_before=2)
        for row in cancel:
            full_text = " ".join(t for t, _ in row)
            hline(paras, full_text.strip(), size=9.5, bold=False, color=C_BLACK,
                  space_before=2, bullet=True)

    circle = d.get("circle")
    if circle:
        hline(paras, "", size=4, space_before=6)
        hline(paras, "◆ CIRCLE PROGRAM", size=9.5, bold=True, color=C_ORANGE, space_before=2)
        hline(paras, f"Code: {circle.get('code', '')}",
              size=9.5, bold=True, color=C_BLACK, space_before=2)
        hline(paras, circle.get("benefit", ""),
              size=9.5, bold=False, color=C_BLACK, space_before=1)
        for r in circle.get("restrictions", []):
            hline(paras, r, size=9, bold=False, color=C_GREY, space_before=1, bullet=True)

    txb_multi(slide, paras, C3X, ROW_Y, C3W, 5.5)

# ── Product data (mirrors build_hpp_v2.py with all corrections) ─────────────

PRODUCTS = [

    # ════════════════ MECHANICAL COVERAGE ════════════════
    {
        "name": "Vehicle Service Protection (ICE)",
        "code": "HFVI", "core": True, "ev_care": False,
        "overview": "HPP-branded mechanical breakdown coverage for eligible Hyundai ICE vehicles. Parts, labor, and genuine Hyundai parts up to 10 years/150,000 miles. The only VSP product backed by Hyundai Motor America.",
        "plans": [
            {"name": "Powertrain", "sub": "(Stated Component)", "desc": "Engine, Transmission, Drive Axle incl. CV joints"},
            {"name": "Gold", "sub": "(Stated Component)", "desc": "Powertrain + front/rear suspension (incl. shocks), A/C, fuel system, electrical system"},
            {"name": "High Technology", "sub": "(Stated Component)", "desc": "Factory-installed ADAS, audio/visual systems, accessories (USB, 120V, wireless charging, WiFi/hotspot, mobile interfaces)"},
            {"name": "Platinum", "sub": "(Exclusionary)", "desc": "All covered parts except listed exclusions"},
            {"name": "Platinum Plus", "sub": "(Florida Only)", "desc": "Platinum + Term Protection items, headlamps, belts and hoses, electrical"},
        ],
        "term_coverage": [
            "Up to 10 years / 150,000 miles",
            "Deductible: $0 (at Hyundai dealer) or $100 (at independent shop)",
            "24-hr Emergency Roadside Assistance included",
            "Rental Car/Rideshare: up to $55/day, 10 days max",
            "Trip Interruption: up to $300/day, 5 days max ($1,500 max)",
            "Genuine Hyundai OEM parts used for all repairs",
        ],
        "elig_label": "Vehicle Eligibility",
        "eligibility": [
            "New and pre-owned Hyundai ICE vehicles",
            "Designed for purchase at time of sale; can be purchased later subject to eligibility",
            "Vehicles with salvage, junk, or Buy-Back titles excluded",
            "Non-U.S. spec models excluded",
            "Permitted Commercial Use allowed (single driver, rideshare, light service)",
            "Prohibited: hauling, livery, fleet/pool vehicles, daily rentals",
        ],
        "reimbursement": [
            ("", [("Repair Order submitted to HPP after repair is authorized and completed", False)]),
            ("", [("Genuine Hyundai parts: OEM pricing; Labor: published labor rate", False)]),
        ],
        "additional_benefits": [
            [("Diagnostic labor covered when repair is covered (Platinum & Gold)", False)],
            [("Necessary fluids/lubricants for covered repair covered (all plans)", False)],
            [("Compatibility: Platinum Plus CANNOT be sold with Wear Protection", False)],
        ],
        "exclusions": [
            "Maintenance services (oil changes, filters, brakes, etc.)",
            "Damage from accidents, collision, theft, vandalism",
            "Pre-existing conditions known at time of purchase",
            "Vehicles used for prohibited commercial purposes",
            "Wear items: brake pads/rotors, wiper blades, tires, belts (unless failure damages covered parts)",
        ],
        "transfer": "Transferable to subsequent owner: $75 fee",
        "cancel": [
            [("Within 30 days: full refund less claims paid", False)],
            [("After 30 days: pro-rata refund less claims paid and ", False), ("$75", True), (" cancellation fee", False)],
        ],
        "circle": {
            "code": "HFVI",
            "benefit": "Circle pricing available for Hyundai employees, family members, and affiliates. Discounted contract pricing on eligible VSP products.",
            "restrictions": [
                "Must qualify under Hyundai Circle program guidelines",
                "Verify eligibility at time of sale",
            ],
        },
    },
    {
        "name": "Vehicle Service Protection — Competitive Makes (ICE)",
        "code": "WFVI", "core": False, "ev_care": False,
        "overview": "PowerProtect mechanical coverage for pre-owned competitive make (non-Hyundai) ICE vehicles sold at Hyundai dealerships. Parts and labor up to 10 years/120,000 miles.",
        "plans": [
            {"name": "Powertrain", "sub": "(Stated Component)", "desc": "Engine, Transmission, Drive Axle incl. CV joints"},
            {"name": "Gold", "sub": "(Stated Component)", "desc": "Powertrain + front/rear suspension (incl. shocks), A/C, fuel system, electrical system"},
            {"name": "Platinum", "sub": "(Exclusionary)", "desc": "All covered parts except listed exclusions"},
            {"name": "Platinum Plus", "sub": "(Florida Only)", "desc": "Platinum + headlamps, belts and hoses, electrical"},
        ],
        "term_coverage": [
            "Up to 10 years / 120,000 miles",
            "Deductible: $0 (at selling dealer) or $100 (independent shop)",
            "24-hr Emergency Roadside Assistance included",
            "Rental Car/Rideshare: up to $55/day, 10 days max",
            "Trip Interruption: up to $300/day, 5 days max ($1,500 max)",
        ],
        "elig_label": "Vehicle Eligibility",
        "eligibility": [
            "Pre-owned competitive make (non-Hyundai/Kia/Genesis) ICE vehicles sold at Hyundai dealerships",
            "Vehicles with salvage, junk, or Buy-Back titles excluded",
            "Non-U.S. spec models excluded",
            "Permitted Commercial Use allowed (single driver, rideshare, light service)",
        ],
        "reimbursement": [
            ("", [("Repair Order submitted to administrator after authorization and repair", False)]),
        ],
        "additional_benefits": [
            [("Compatibility: Platinum Plus CANNOT be sold with Wear Protection", False)],
        ],
        "exclusions": [
            "Maintenance services",
            "Accidental damage, collision, vandalism, theft",
            "Pre-existing conditions known at time of purchase",
            "Prohibited commercial use vehicles",
        ],
        "transfer": "Transferable to subsequent owner: $75 fee",
        "cancel": [
            [("Within 30 days: full refund less claims paid", False)],
            [("After 30 days: pro-rata less claims paid and ", False), ("$75", True), (" cancellation fee", False)],
        ],
        "circle": None,
    },
    {
        "name": "Vehicle Service Protection Original Owner (ICE)",
        "code": "HFOI", "core": True, "ev_care": False,
        "overview": "HPP-branded VSP designed for vehicle lessees who are purchasing or planning to purchase their leased Hyundai vehicle at lease end. Platinum-level exclusionary coverage only.",
        "plans": [
            {"name": "Platinum", "sub": "(Exclusionary)", "desc": "All covered parts except listed exclusions — same breadth as standard VSP Platinum"},
            {"name": "Platinum Plus", "sub": "(Florida Only)", "desc": "Platinum + headlamps, belts and hoses, electrical coverage"},
        ],
        "term_coverage": [
            "Up to 10 years / 150,000 miles from original in-service date",
            "Deductible: $0 (Hyundai dealer) or $100 (independent shop)",
            "24-hr Emergency Roadside Assistance included",
            "Rental Car/Rideshare: up to $55/day, 10 days max",
            "Trip Interruption: up to $300/day, 5 days max ($1,500 max)",
        ],
        "elig_label": "Original Owner Eligibility",
        "eligibility": [
            "Vehicle must be currently owned by original owner or being purchased from lease by original lessee",
            "Purchased at time of vehicle purchase (new or lease buyout); post-sale subject to eligibility",
            "Vehicles with salvage/junk/Buy-Back titles excluded",
            "Non-U.S. spec models excluded",
            "Permitted Commercial Use allowed",
        ],
        "reimbursement": [
            ("", [("Repair Order submitted to HPP after authorization and repair", False)]),
            ("", [("Genuine Hyundai OEM parts; published labor rate", False)]),
        ],
        "additional_benefits": [
            [("Diagnostic labor covered when repair is covered (Platinum)", False)],
            [("Necessary fluids/lubricants for covered repair included", False)],
            [("Compatibility: Platinum Plus CANNOT be sold with Wear Protection", False)],
        ],
        "exclusions": [
            "Maintenance services",
            "Accidental/collision/vandalism/theft damage",
            "Pre-existing conditions",
            "Vehicles for prohibited commercial purposes",
        ],
        "transfer": "Transferable to subsequent owner: $75 fee",
        "cancel": [
            [("Within 30 days: full refund less claims paid", False)],
            [("After 30 days: pro-rata less claims paid and ", False), ("$75", True), (" cancellation fee", False)],
        ],
        "circle": None,
    },
    {
        "name": "Certified Used Vehicle (CUV) Wrap",
        "code": "HFCI", "core": True, "ev_care": False,
        "overview": "HPP-branded wrap coverage for Hyundai Certified Used ICE vehicles. Picks up where the remaining manufacturer warranty leaves off, extending coverage to match the original factory warranty period.",
        "plans": [
            {"name": "Gold", "sub": "(Stated Component)", "desc": "Front/rear suspension (incl. shocks), A/C, fuel system, electrical system"},
            {"name": "Platinum", "sub": "(Exclusionary)", "desc": "All covered parts except those listed as exclusions"},
            {"name": "Platinum Plus", "sub": "(Florida Only)", "desc": "Platinum + headlamps, belts and hoses, electrical"},
        ],
        "term_coverage": [
            "Up to 10 years / 150,000 miles from original in-service date",
            "Deductible: $0 (Hyundai dealer) or $100 (independent shop)",
            "24-hr Emergency Roadside Assistance included",
            "Rental Car/Rideshare: up to $55/day, 10 days max",
            "Trip Interruption: up to $300/day, 5 days max ($1,500 max)",
        ],
        "elig_label": "CUV Eligibility",
        "eligibility": [
            "Hyundai Certified Used ICE vehicles only",
            "Less than 6 years from original in-service date and less than 80,000 miles at purchase",
            "Must have at least 1 month and 1,000 miles of manufacturer's comprehensive warranty remaining",
            "Must be a Hyundai Certified Used vehicle (UCI inspection required)",
            "Vehicles with salvage/junk/Buy-Back titles excluded",
        ],
        "reimbursement": [
            ("", [("Genuine Hyundai OEM parts; published labor rate", False)]),
        ],
        "additional_benefits": [
            [("Compatibility: Platinum Plus CANNOT be sold with Wear Protection", False)],
        ],
        "exclusions": [
            "Non-certified used Hyundai vehicles",
            "Maintenance services",
            "Accidental/collision/vandalism/theft damage",
            "Vehicles for prohibited commercial purposes",
        ],
        "transfer": "Transferable to subsequent owner: $75 fee",
        "cancel": [
            [("Within 30 days: full refund less claims paid", False)],
            [("After 30 days: pro-rata less claims paid and ", False), ("$75", True), (" cancellation fee", False)],
        ],
        "circle": {
            "code": "HFCI",
            "benefit": "Circle pricing available for Hyundai employees, family members, and affiliates.",
            "restrictions": [
                "Must qualify under Hyundai Circle program guidelines",
                "Verify eligibility at time of sale",
            ],
        },
    },
    {
        "name": "High Mileage Vehicle Service Protection (ICE)",
        "code": "WFHM", "core": False, "ev_care": False,
        "overview": "PowerProtect Powertrain-only coverage for used high-mileage Hyundai and competitive make ICE vehicles. Designed for vehicles between 60,001 and 175,000 miles. Available at time of vehicle sale only.",
        "plans": [
            {"name": "Powertrain", "sub": "(Stated Component Only)", "desc": "Engine, Transmission, Drive Axle incl. CV joints. No Gold or Platinum tiers available for this product."},
        ],
        "term_coverage": [
            "Term options: 1 year/12,000 miles up to 3 years/36,000 miles",
            "Vehicle odometer: 60,001 – 175,000 miles at time of purchase",
            "Deductible: $100 per repair visit",
            "24-hr Emergency Roadside Assistance included",
        ],
        "elig_label": "High Mileage Eligibility",
        "eligibility": [
            "Used Hyundai and competitive make ICE vehicles with 60,001 – 175,000 miles at time of sale",
            "Available at time of vehicle sale only (no post-sale purchases)",
            "Vehicles with salvage/junk/Buy-Back titles excluded",
            "Non-U.S. spec models excluded",
            "No Circle pricing available for this product",
        ],
        "reimbursement": [
            ("", [("Repair Order submitted to administrator after authorization and repair", False)]),
        ],
        "additional_benefits": [
            [("No Circle program pricing available", False)],
        ],
        "exclusions": [
            "Gold and Platinum coverage levels not available",
            "Vehicles under 60,000 miles or over 175,000 miles at time of sale",
            "Maintenance services, wear items, accidental damage",
            "Prohibited commercial use",
        ],
        "transfer": "Transferable to subsequent owner: $75 fee",
        "cancel": [
            [("Within 30 days: full refund less claims paid", False)],
            [("After 30 days: pro-rata less claims paid and ", False), ("$75", True), (" cancellation fee", False)],
        ],
        "circle": None,
    },
    {
        "name": "EV Care Vehicle Service Protection",
        "code": "HFVE", "core": True, "ev_care": True,
        "overview": "HPP-branded EV Care mechanical breakdown coverage for eligible Hyundai Electric and Hybrid vehicles. Parts, labor, and genuine Hyundai parts up to 12 years/200,000 miles.",
        "plans": [
            {"name": "Battery", "sub": "(Stated Component)", "desc": "Lithium-Ion Battery Pack, Battery Module, BMS, Battery Degradation (<70%), HV Pre Charger, Traction Motor (incl. housing), Onboard Charger, Inverter, Converter"},
            {"name": "High Technology", "sub": "(Stated Component)", "desc": "Factory-installed ADAS, audio/visual systems, accessories (USB, 120V, wireless charging, WiFi/hotspot, mobile interfaces)"},
            {"name": "Platinum", "sub": "(Exclusionary)", "desc": "All covered parts except listed exclusions"},
            {"name": "Platinum Plus", "sub": "(Florida Only)", "desc": "Platinum + Term Protection items, headlamps, belts and hoses, electrical"},
        ],
        "term_coverage": [
            "Up to 12 years / 200,000 miles",
            "Deductible: $0 (Hyundai dealer) or $100 (independent shop)",
            "24-hr Emergency Roadside Assistance included",
            "Rental Car/Rideshare: up to $55/day, 10 days max",
            "Trip Interruption: up to $300/day, 5 days max ($1,500 max)",
            "Genuine Hyundai OEM parts used for all repairs",
        ],
        "elig_label": "EV Vehicle Eligibility",
        "eligibility": [
            "New and pre-owned Hyundai Electric and Hybrid vehicles",
            "Designed for purchase at time of sale; can be purchased later subject to eligibility",
            "Vehicles with salvage/junk/Buy-Back titles excluded",
            "Non-U.S. spec models excluded",
            "Permitted Commercial Use allowed (single driver, rideshare, light service)",
        ],
        "reimbursement": [
            ("", [("Repair Order submitted to HPP after authorization and repair", False)]),
            ("", [("Genuine Hyundai parts; published labor rate", False)]),
        ],
        "additional_benefits": [
            [("Compatibility: Platinum Plus CANNOT be sold with Wear Protection", False)],
            [("Diagnostic labor covered when repair is covered (Platinum)", False)],
        ],
        "exclusions": [
            "Maintenance services",
            "Accidental/collision/vandalism/theft damage",
            "Battery degradation if above 70% capacity",
            "Prohibited commercial use vehicles",
        ],
        "transfer": "Transferable to subsequent owner: $75 fee",
        "cancel": [
            [("Within 30 days: full refund less claims paid", False)],
            [("After 30 days: pro-rata less claims paid and ", False), ("$75", True), (" cancellation fee", False)],
        ],
        "circle": {
            "code": "HFVE",
            "benefit": "Circle pricing available for Hyundai employees, family members, and affiliates on EV Care VSP products.",
            "restrictions": [
                "Must qualify under Hyundai Circle program guidelines",
                "Verify eligibility at time of sale",
            ],
        },
    },
    {
        "name": "EV Care VSP — Competitive Makes",
        "code": "WFVE", "core": False, "ev_care": True,
        "overview": "PowerProtect EV mechanical coverage for pre-owned competitive make Electric vehicles sold at Hyundai dealerships. Parts and labor up to 12 years/120,000 miles.",
        "plans": [
            {"name": "Platinum", "sub": "(Exclusionary)", "desc": "All covered parts except listed exclusions"},
            {"name": "Platinum Plus", "sub": "(Florida Only)", "desc": "Platinum + headlamps, belts and hoses, electrical"},
        ],
        "term_coverage": [
            "Up to 12 years / 120,000 miles",
            "Deductible: $0 (selling dealer) or $100 (independent shop)",
            "24-hr Emergency Roadside Assistance included",
            "Rental Car/Rideshare: up to $55/day, 10 days max",
        ],
        "elig_label": "Competitive EV Eligibility",
        "eligibility": [
            "Pre-owned competitive make (non-Hyundai) Electric vehicles sold at Hyundai dealerships",
            "Vehicles with salvage/junk/Buy-Back titles excluded",
            "Non-U.S. spec models excluded",
            "Permitted Commercial Use allowed",
        ],
        "reimbursement": [
            ("", [("Repair Order submitted to administrator after authorization and repair", False)]),
        ],
        "additional_benefits": [
            [("Compatibility: Platinum Plus CANNOT be sold with Wear Protection", False)],
        ],
        "exclusions": [
            "Maintenance services",
            "Accidental/collision/vandalism/theft damage",
            "Prohibited commercial use",
        ],
        "transfer": "Transferable to subsequent owner: $75 fee",
        "cancel": [
            [("Within 30 days: full refund less claims paid", False)],
            [("After 30 days: pro-rata less claims paid and ", False), ("$75", True), (" cancellation fee", False)],
        ],
        "circle": None,
    },
    {
        "name": "EV Care VSP — Original Owner",
        "code": "HFOE", "core": True, "ev_care": True,
        "overview": "HPP-branded EV Care VSP for original owners/lessees purchasing their Hyundai EV at lease end. Battery and Platinum exclusionary coverage options.",
        "plans": [
            {"name": "Battery", "sub": "(Stated Component)", "desc": "Lithium-Ion Battery Pack, Battery Module, BMS, Battery Degradation (<70%), HV Pre Charger, Traction Motor (incl. housing), Onboard Charger, Inverter, Converter"},
            {"name": "Platinum", "sub": "(Exclusionary)", "desc": "All covered parts except listed exclusions"},
            {"name": "Platinum Plus", "sub": "(Florida Only)", "desc": "Platinum + Term Protection items, headlamps, belts and hoses, electrical"},
        ],
        "term_coverage": [
            "Up to 12 years / 200,000 miles from original in-service date",
            "Deductible: $0 (Hyundai dealer) or $100 (independent shop)",
            "24-hr Emergency Roadside Assistance included",
            "Rental Car/Rideshare: up to $55/day, 10 days max",
        ],
        "elig_label": "EV Original Owner Eligibility",
        "eligibility": [
            "Vehicle must be currently owned by original owner or being purchased from lease by original lessee",
            "New Hyundai Electric or Hybrid vehicles",
            "Vehicles with salvage/junk/Buy-Back titles excluded",
            "Non-U.S. spec models excluded",
        ],
        "reimbursement": [
            ("", [("Genuine Hyundai OEM parts; published labor rate", False)]),
        ],
        "additional_benefits": [
            [("Compatibility: Platinum Plus CANNOT be sold with Wear Protection", False)],
        ],
        "exclusions": [
            "Maintenance services",
            "Accidental damage, collision, vandalism, theft",
            "Battery degradation above 70% capacity threshold",
            "Prohibited commercial use",
        ],
        "transfer": "Transferable to subsequent owner: $75 fee",
        "cancel": [
            [("Within 30 days: full refund less claims paid", False)],
            [("After 30 days: pro-rata less claims paid and ", False), ("$75", True), (" cancellation fee", False)],
        ],
        "circle": None,
    },
    {
        "name": "EV Care Certified Used Vehicle (CUV) Wrap",
        "code": "HFCE", "core": True, "ev_care": True,
        "overview": "HPP-branded EV Care wrap for Hyundai Certified Used Electric and Hybrid vehicles. Extends coverage beyond remaining manufacturer warranty.",
        "plans": [
            {"name": "Battery", "sub": "(Stated Component)", "desc": "Lithium-Ion Battery Pack, BMS, Battery Degradation (<70%), HV Pre Charger, Traction Motor (incl. housing), Onboard Charger, Inverter, Converter"},
            {"name": "Gold", "sub": "(Stated Component)", "desc": "Front/rear suspension (incl. shocks), A/C, fuel system, electrical system"},
            {"name": "Platinum", "sub": "(Exclusionary)", "desc": "All covered parts except listed exclusions"},
            {"name": "Platinum Plus", "sub": "(Florida Only)", "desc": "Platinum + headlamps, belts and hoses, electrical"},
        ],
        "term_coverage": [
            "Up to 12 years / 200,000 miles from original in-service date",
            "Deductible: $0 (Hyundai dealer) or $100 (independent shop)",
            "24-hr Emergency Roadside Assistance included",
            "Rental Car/Rideshare: up to $55/day, 10 days max",
        ],
        "elig_label": "EV CUV Eligibility",
        "eligibility": [
            "Hyundai Certified Used Electric and Hybrid vehicles only",
            "Less than 6 years from original in-service date and less than 80,000 miles at purchase",
            "Must have at least 1 month and 1,000 miles of manufacturer's comprehensive warranty remaining",
            "UCI (Used Car Inspection) required",
            "Vehicles with salvage/junk/Buy-Back titles excluded",
        ],
        "reimbursement": [
            ("", [("Genuine Hyundai OEM parts; published labor rate", False)]),
        ],
        "additional_benefits": [
            [("Compatibility: Platinum Plus CANNOT be sold with Wear Protection", False)],
        ],
        "exclusions": [
            "Non-certified used Hyundai EV/Hybrid vehicles",
            "Maintenance services",
            "Battery degradation above 70%",
            "Prohibited commercial use",
        ],
        "transfer": "Transferable to subsequent owner: $75 fee",
        "cancel": [
            [("Within 30 days: full refund less claims paid", False)],
            [("After 30 days: pro-rata less claims paid and ", False), ("$75", True), (" cancellation fee", False)],
        ],
        "circle": {
            "code": "HFCE",
            "benefit": "Circle pricing available for Hyundai employees, family members, and affiliates.",
            "restrictions": [
                "Must qualify under Hyundai Circle program guidelines",
            ],
        },
    },
    {
        "name": "VSP Livery Coverage (ICE & EV)",
        "code": "HCVL", "core": True, "ev_care": False,
        "overview": "HPP-branded VSP for new Hyundai ICE and EV vehicles used for legal Livery purposes (limousine, chauffeured, private hire, shuttle). Available only through select dealers with District Manager approval. Coverage from 5 years/100,000 to 5 years/300,000 miles.",
        "plans": [
            {"name": "ICE — Powertrain", "sub": "(Stated Component)", "desc": "Engine, Transmission, Drive Axle incl. CV joints"},
            {"name": "ICE — Gold", "sub": "(Stated Component)", "desc": "Powertrain + front/rear suspension (incl. shocks), A/C, fuel system, electrical system"},
            {"name": "ICE — Platinum", "sub": "(Exclusionary)", "desc": "All covered parts except listed exclusions"},
            {"name": "EV — Battery", "sub": "(Stated Component)", "desc": "Lithium-Ion Battery Pack, BMS, Battery Degradation (<70%), HV Pre Charger, Traction Motor, Onboard Charger, Inverter, Converter"},
            {"name": "EV — Platinum", "sub": "(Exclusionary)", "desc": "All covered parts except listed exclusions"},
        ],
        "term_coverage": [
            "5 years / 100,000 miles up to 5 years / 300,000 miles",
            "Deductible: $0",
            "24-hr Emergency Roadside Assistance included",
            "Rental Car/Rideshare: up to $55/day, 10 days max",
            "Trip Interruption: up to $300/day, 5 days max ($1,500 max)",
        ],
        "elig_label": "Livery Eligibility",
        "eligibility": [
            "New ICE and EV Hyundai vehicles only",
            "Odometer must be less than 10,000 miles at time of purchase",
            "Available at time of vehicle sale only",
            "Vehicle must be used for legal, licensed Livery purposes only",
            "Livery defined as: limousine, chauffeured vehicle, private hire vehicle, or shuttle service",
            "Contact District Manager of Insurance for availability",
        ],
        "reimbursement": [
            ("", [("Genuine Hyundai OEM parts; published labor rate", False)]),
        ],
        "additional_benefits": [
            [("Available only through select dealers — contact DM of Insurance", False)],
            [("Diagnostic labor covered when repair is covered (Platinum & Gold)", False)],
        ],
        "exclusions": [
            "Non-Livery commercial use (hauling, delivery, fleet pool, towing, government/military)",
            "Non-U.S. spec models",
            "Salvage/junk/Buy-Back vehicles",
            "Maintenance services, wear items",
        ],
        "transfer": "Transferable to subsequent owner: $75 fee",
        "cancel": [
            [("Within 30 days: full refund less claims paid", False)],
            [("After 30 days: pro-rata less claims paid and ", False), ("$75", True), (" cancellation fee", False)],
        ],
        "circle": None,
    },
    {
        "name": "Wear Protection (ICE & EV)",
        "code": "HFWP", "core": True, "ev_care": False,
        "overview": "HPP-branded short-term wear coverage designed for lessees and short-term owners. Covers scheduled wear items that are not covered by the manufacturer warranty, helping customers avoid unexpected out-of-pocket costs.",
        "plans": [
            {"name": "Standard", "sub": "", "desc": "One set front & rear brake pads/shoes, one 12V battery replacement, one wheel alignment, unlimited headlamps (non-impact), unlimited engine belts & hoses, one set wiper blades, unlimited fuses & light bulbs (non-impact)"},
            {"name": "+ Brake Rotor Add-On", "sub": "(Optional)", "desc": "Adds one replacement set of front and rear brake rotors during the term"},
        ],
        "term_coverage": [
            "Up to 4 years / 60,000 miles",
            "Deductible: $0",
            "Designed for lease-length terms",
        ],
        "elig_label": "Wear Protection Eligibility",
        "eligibility": [
            "New Hyundai ICE and EV vehicles",
            "Available at time of vehicle sale only",
            "NOT compatible with Platinum Plus VSP coverage",
            "Vehicles for prohibited commercial use excluded",
        ],
        "reimbursement": [
            ("", [("Published Hyundai labor rate; OEM parts pricing", False)]),
        ],
        "additional_benefits": [
            [("Designed to complement lease products — helps protect against lease-end charges", False)],
            [("CANNOT be sold with Platinum Plus VSP coverage", False)],
        ],
        "exclusions": [
            "Mechanical breakdown coverage (use VSP for that)",
            "Impact damage (headlamps, bulbs damaged by road impact)",
            "Tires",
            "Pre-existing conditions",
            "Prohibited commercial use vehicles",
        ],
        "transfer": "Non-transferable",
        "cancel": [
            [("Within 30 days: full refund less claims paid", False)],
            [("After 30 days: pro-rata less claims paid and ", False), ("$75", True), (" cancellation fee", False)],
        ],
        "circle": {
            "code": "HFWP",
            "benefit": "Circle pricing available for Hyundai employees, family members, and affiliates.",
            "restrictions": [
                "Must qualify under Hyundai Circle program guidelines",
            ],
        },
    },

    # ════════════════ MAINTENANCE PLANS ════════════════
    {
        "name": "Maintenance Basic Wrap (ICE)",
        "code": "HFBI", "core": True, "ev_care": False,
        "overview": "Extends and wraps Hyundai Complimentary Maintenance (HCM) on MY2025 and older ICE vehicles. Covers the same services as HCM: oil & filter, tire rotation, and multi-point inspection. Keeps customers returning to the dealership for service.",
        "plans": [
            {"name": "Basic Wrap", "sub": "(HFBI)", "desc": "Oil & oil filter change, tire rotation, multi-point inspection per interval. Severe Usage upgrade available. Extends beyond standard HCM (3 yrs/36,000 mi)."},
        ],
        "term_coverage": [
            "Up to 8 years / 96,000 miles",
            "MY2025 and older Hyundai ICE vehicles with HCM only",
            "Deductible: $0",
            "Reimbursed per Maintenance Interval schedule",
        ],
        "elig_label": "Basic Wrap Eligibility",
        "eligibility": [
            "MY2025 and older ICE Hyundai vehicles WITH active HCM only",
            "Available at time of sale or before 6,000 miles on odometer",
            "NOT eligible on any vehicle that already has Pre-Paid Maintenance (PPM)",
            "Electric vehicles excluded — see EV Care Maintenance",
        ],
        "reimbursement": [
            ("", [("Reimbursed per Maintenance Interval reimbursement sheets", False)]),
            ("", [("Services must be performed within 5,000 miles or 6 months of scheduled interval", False)]),
        ],
        "additional_benefits": [
            [("Non-Transferable — contract stays with original buyer", False)],
            [("Severe Usage interval upgrade available at purchase", False)],
        ],
        "exclusions": [
            "Electric vehicles (EV) — use EV Care Maintenance instead",
            "MY2026+ vehicles — use Pre-Paid Maintenance (HFPI)",
            "Vehicles without HCM",
            "Services performed outside the allowed interval window (5,000 mi / 6 months)",
            "Mechanical breakdown repairs",
        ],
        "transfer": "Non-Transferable",
        "cancel": [
            [("Within 30 days: full refund if no covered services provided", False)],
            [("Next 35 months: full refund less ", False), ("$50", True), (" fee if no services provided", False)],
            [("After initial look period or if service provided: pro-rata less cost of services + ", False), ("$50", True), (" fee", False)],
        ],
        "circle": None,
    },
    {
        "name": "Maintenance Scheduled Wrap (ICE)",
        "code": "HFSI", "core": True, "ev_care": False,
        "overview": "Extends and wraps Hyundai Complimentary Maintenance (HCM) on MY2025 and older ICE vehicles. Includes all Basic Wrap services plus all owner's manual-scheduled services: cabin filter, engine air filter, brake fluid, spark plugs, transmission fluid, and other fluids.",
        "plans": [
            {"name": "Scheduled Wrap", "sub": "(HFSI)", "desc": "All Basic Wrap services (oil & filter, tire rotation, multi-point inspection) plus cabin air filter, engine air filter, brake fluid, spark plugs, transmission fluid, differential/case oil per owner's manual schedule. Severe Usage upgrade available."},
        ],
        "term_coverage": [
            "Up to 8 years / 96,000 miles",
            "MY2025 and older Hyundai ICE vehicles with HCM only",
            "Deductible: $0",
            "Reimbursed per Maintenance Interval schedule",
        ],
        "elig_label": "Scheduled Wrap Eligibility",
        "eligibility": [
            "MY2025 and older ICE Hyundai vehicles WITH active HCM only",
            "Available at time of sale or before 6,000 miles on odometer",
            "NOT eligible on any vehicle that already has Pre-Paid Maintenance (PPM)",
            "Electric vehicles excluded — see EV Care Maintenance",
        ],
        "reimbursement": [
            ("", [("Reimbursed per Maintenance Interval reimbursement sheets", False)]),
            ("", [("Services must be performed within 5,000 miles or 6 months of scheduled interval", False)]),
        ],
        "additional_benefits": [
            [("Non-Transferable — contract stays with original buyer", False)],
            [("Severe Usage interval upgrade available at purchase", False)],
            [("Broader coverage than Basic Wrap — includes all scheduled service items", False)],
        ],
        "exclusions": [
            "Electric vehicles (EV) — use EV Care Maintenance instead",
            "MY2026+ vehicles — use Pre-Paid Maintenance Scheduled (HFSI/PPM)",
            "Vehicles without HCM",
            "Services performed outside the allowed interval window (5,000 mi / 6 months)",
            "Mechanical breakdown repairs",
        ],
        "transfer": "Non-Transferable",
        "cancel": [
            [("Within 30 days: full refund if no covered services provided", False)],
            [("Next 35 months: full refund less ", False), ("$50", True), (" fee if no services provided", False)],
            [("After initial look period or if service provided: pro-rata less cost of services + ", False), ("$50", True), (" fee", False)],
        ],
        "circle": None,
    },
    {
        "name": "Pre-Paid Maintenance (ICE)",
        "code": "HFPI", "core": True, "ev_care": False,
        "overview": "Prepaid maintenance program for Hyundai ICE vehicles. Basic (HFPI) available on all model years at any time. Scheduled (HFSI code used for PPM) available on new MY2026+ vehicles at time of sale only. Reimbursed per interval schedule.",
        "plans": [
            {"name": "Basic Maintenance", "sub": "(HFPI)", "desc": "Oil & oil filter services at selected intervals; tire rotations and multi-point inspection. Normal or Severe usage options. All Hyundai ICE model years."},
            {"name": "Scheduled Maintenance", "sub": "(MY2026+ New Only)", "desc": "Includes Basic services plus cabin filter, engine air filter, brake fluid, spark plugs, transmission fluid, differential/case oil per owner's manual. New MY2026+ at time of sale only; ≤6,000 miles."},
        ],
        "term_coverage": [
            "Up to 8 years (96 months) / 96,000 miles",
            "Deductible: $0",
            "Basic: available anytime (no time-of-sale restriction)",
            "Scheduled: new MY2026+ only, at time of sale, ≤6,000 miles",
        ],
        "elig_label": "PPM Eligibility",
        "eligibility": [
            "All Hyundai ICE vehicles (Basic); new MY2026+ only (Scheduled)",
            "NOT eligible on vehicles with active HCM — use Maintenance Wrap instead",
            "Electric vehicles excluded — see EV Care Maintenance",
        ],
        "reimbursement": [
            ("", [("Reimbursed per Pre-Paid Maintenance Interval reimbursement sheets", False)]),
            ("", [("Services must be performed within interval guidelines", False)]),
        ],
        "additional_benefits": [
            [("Keeps customers returning to the dealership on a scheduled cadence", False)],
            [("Severe Usage interval option available", False)],
        ],
        "exclusions": [
            "Vehicles with active HCM (use Maintenance Wrap)",
            "Electric vehicles",
            "Mechanical breakdown repairs not covered",
            "Services performed outside interval window",
        ],
        "transfer": "Non-Transferable",
        "cancel": [
            [("Within 30 days: full refund if no services provided", False)],
            [("After 30 days: pro-rata less cost of services provided + ", False), ("$50", True), (" cancellation fee", False)],
        ],
        "circle": {
            "code": "HFEP",
            "benefit": "Circle pricing available for Hyundai employees, family members, and affiliates on Pre-Paid Maintenance contracts.",
            "restrictions": [
                "Must qualify under Hyundai Circle program guidelines",
                "Verify eligibility at time of sale",
            ],
        },
    },
    {
        "name": "Pre-Paid Maintenance — Competitive Makes (ICE)",
        "code": "WFPI", "core": False, "ev_care": False,
        "overview": "Prepaid maintenance for competitive make (non-Hyundai) ICE vehicles sold at Hyundai dealerships. Three simple plan options. Can be sold at any time — no time-of-sale restriction.",
        "plans": [
            {"name": "Basic", "sub": "", "desc": "Conventional oil & oil filter changes at selected intervals"},
            {"name": "Plus", "sub": "", "desc": "Conventional oil changes + tire rotations for covered vehicle"},
            {"name": "Synthetic Plus", "sub": "", "desc": "Synthetic oil changes + tire rotations for covered vehicle"},
        ],
        "term_coverage": [
            "Up to 7 years / 105,000 miles",
            "Deductible: $0",
            "Can be sold at any time — no time-of-sale restriction",
        ],
        "elig_label": "Competitive Makes Eligibility",
        "eligibility": [
            "Competitive make (non-Hyundai) ICE vehicles sold at Hyundai dealerships",
            "Available at any time — not restricted to time of vehicle sale",
        ],
        "reimbursement": [
            ("", [("Reimbursed per Maintenance Interval reimbursement sheets", False)]),
        ],
        "additional_benefits": [
            [("Flexible selling point — can be offered post-sale", False)],
        ],
        "exclusions": [
            "Hyundai-branded vehicles — use HFPI instead",
            "Electric vehicles",
            "Services beyond selected plan scope",
        ],
        "transfer": "Non-Transferable",
        "cancel": [
            [("Within 30 days: full refund if no services provided", False)],
            [("After 30 days: pro-rata less cost of services + ", False), ("$50", True), (" cancellation fee", False)],
        ],
        "circle": None,
    },
    {
        "name": "EV Care Maintenance",
        "code": "HFEM", "core": True, "ev_care": True,
        "overview": "Prepaid EV maintenance program for Hyundai Electric and Hybrid vehicles. Basic available on new and used EVs at any time. Maintenance Plus available on new MY2026+ EVs at time of sale only. Reimbursed per interval schedule.",
        "plans": [
            {"name": "Basic", "sub": "(New + Used)", "desc": "Tire rotations and multi-point maintenance inspection every interval; cabin air filter replacement every other interval"},
            {"name": "Maintenance Plus", "sub": "(New MY2026+ Only)", "desc": "All Basic items + all owner's manual scheduled services + wiper blades (every other interval), wheel alignment at 32K & 72K mi, 12V battery at 72K mi, brake pads at 72K mi, brake fluid at 48K & 96K mi"},
        ],
        "term_coverage": [
            "Up to 8 years (96 months) / 96,000 miles",
            "Deductible: $0",
            "Basic: sold at any time on new and used EVs",
            "Maintenance Plus: new MY2026+ EV only at time of sale, ≤6,000 miles",
        ],
        "elig_label": "EV Maintenance Eligibility",
        "eligibility": [
            "New and used Hyundai Electric and Hybrid vehicles (Basic)",
            "New MY2026+ Hyundai EV only at time of sale with ≤6,000 miles (Maintenance Plus)",
            "ICE vehicles excluded — use HFPI or Maintenance Wrap instead",
        ],
        "reimbursement": [
            ("", [("Reimbursed per EV Care Maintenance Interval reimbursement sheets", False)]),
        ],
        "additional_benefits": [
            [("Keeps EV customers returning to Hyundai dealer service department", False)],
            [("Maintenance Plus is the most comprehensive EV maintenance offering", False)],
        ],
        "exclusions": [
            "ICE vehicles — use Pre-Paid Maintenance instead",
            "Mechanical breakdown repairs not covered",
            "Services performed outside interval window",
        ],
        "transfer": "Non-Transferable",
        "cancel": [
            [("Within 30 days: full refund if no services provided", False)],
            [("After 30 days: pro-rata less cost of services + ", False), ("$50", True), (" cancellation fee", False)],
        ],
        "circle": {
            "code": "HFEM",
            "benefit": "Circle pricing available for Hyundai employees, family members, and affiliates on EV Care Maintenance contracts.",
            "restrictions": [
                "Must qualify under Hyundai Circle program guidelines",
            ],
        },
    },

    # ════════════════ DAMAGE CARE ════════════════
    {
        "name": "Lease-End Protection",
        "code": "HFEL", "core": True, "ev_care": False,
        "overview": "HPP-branded product that waives eligible excess wear and use charges at lease end. Covers billable charges assessed by Hyundai Motor Finance/Hyundai Lease Titling Trust, up to $1,000 per single event. Available on Hyundai Motor Finance leases only.",
        "plans": [
            {"name": "Lease-End Protection", "sub": "", "desc": "Waives eligible excess wear and use charges assessed at lease return. Single event limit of $1,000. Covers up to 200 excess miles at $0.20/mile. Covers items including paint damage, dents, glass, carpet, trim, tires, chrome, audio, mirrors, and more."},
        ],
        "term_coverage": [
            "Minimum 12 months / Maximum 72 months",
            "Maximum eligible odometer: 150,000 miles",
            "Single event limit: $1,000",
            "Deductible: $0",
            "Vehicle must be turned in within 1 year of original scheduled termination date",
            "Not available in NY or MD",
        ],
        "elig_label": "Lease-End Eligibility",
        "eligibility": [
            "Hyundai Motor Finance/Hyundai Lease Titling Trust leases only",
            "Available on Hyundai ICE, EV, hybrid, and hydrogen fuel cell leased vehicles",
            "Available at time of lease only",
            "Commercial and fleet vehicles excluded; professional use (realtor, rideshare) allowed",
            "NOT available in New York (NY) or Maryland (MD)",
        ],
        "reimbursement": [
            ("", [("Benefit paid directly to Hyundai Motor Finance/Hyundai Lease Titling Trust", False)]),
            ("", [("No benefit if vehicle is purchased by customer or dealer at lease end", False)]),
        ],
        "additional_benefits": [
            [("Covers items including: paint damage, exterior dents, glass, carpet, trim, mirrors, audio, tires", False)],
            [("Covers up to 200 excess miles at $0.20/mile waived", False)],
        ],
        "exclusions": [
            "Alteration charges and repairs completed prior to lease turn-in",
            "Any single charge exceeding $1,000 (excluded amount is customer's responsibility)",
            "Vehicles purchased by customer or dealer at lease end",
            "Not available in NY or MD",
        ],
        "transfer": "Transferable to subsequent lessee: $75 fee",
        "cancel": [
            [("Within 30 days: full refund", False)],
            [("After 30 days: pro-rata refund subject to ", False), ("$75", True), (" cancellation fee", False)],
            [("If waiver benefit has been provided: deemed fully earned — no cancellation", False)],
        ],
        "circle": {
            "code": "HFEL",
            "benefit": "Circle pricing available for Hyundai employees, family members, and affiliates on Lease-End Protection contracts.",
            "restrictions": [
                "Must qualify under Hyundai Circle program guidelines",
                "Verify eligibility at time of lease",
            ],
        },
    },
    {
        "name": "Multi-Coverage Protection",
        "code": "HFMC", "core": True, "ev_care": False,
        "overview": "Umbrella damage care product that bundles Tire & Wheel Protection, Dent Protection, Windshield Protection, and Key Protection into a single contract. Available in four bundle packages (Platinum, Gold, Silver, Bronze) or as individual stand-alone components on the Multi-Coverage form.",
        "plans": [
            {"name": "Platinum Package", "sub": "(Complete Bundle)", "desc": "Tire & Wheel Protection + Dent Protection + Windshield Protection + Key Protection"},
            {"name": "Gold Package", "sub": "(No Key)", "desc": "Tire & Wheel Protection + Dent Protection + Windshield Protection"},
            {"name": "Silver Package", "sub": "(FL Bundle Only)", "desc": "Tire & Wheel Protection + Dent Protection + Key Protection"},
            {"name": "Bronze Package", "sub": "(Entry Bundle)", "desc": "Tire & Wheel Protection + Dent Protection"},
        ],
        "term_coverage": [
            "Up to 7 years / Unlimited miles",
            "No deductible",
            "No limit on occurrences (except: max 8 cosmetic wheel repairs; max 1 windshield replacement; hail damage up to $1,000/occurrence)",
            "Windshield Protection NOT available in Florida (stand-alone or bundle)",
            "Tire must have min. 3/32 inch tread depth at time of damage",
        ],
        "elig_label": "Multi-Coverage Eligibility",
        "eligibility": [
            "Hyundai vehicles and competitive makes (incl. Kia or Genesis) sold at Hyundai dealerships",
            "Available at time of vehicle sale or lease only",
            "Commercial and fleet vehicles excluded; professional use (realtor, rideshare) allowed",
            "Windshield Protection: NOT available in Florida",
            "Key Protection: Motor Club license required in Massachusetts",
        ],
        "reimbursement": [
            ("Tire & Wheel:", [("Tires: 10% markup (not to exceed MSRP); Wheels: 20% markup (not to exceed MSRP); Labor: up to $50 repair / $30 replacement", False)]),
            ("Dent:", [("Up to $65 per single dent; up to $100 per panel; hail up to $1,000/occurrence", False)]),
            ("Windshield:", [("Chip repair labor; one (1) windshield replacement including installation and calibration", False)]),
            ("Key:", [("Vehicle key/fob replacement up to $800/occurrence; lockout up to $100; towing up to $100", False)]),
        ],
        "additional_benefits": [
            [("Products can be bundled or sold as individual stand-alone on the Multi-Coverage form", False)],
            [("Towing reimbursement up to $100 for road hazard tire/wheel damage", False)],
            [("Dent repairs completed at dealership or customer-preferred location", False)],
        ],
        "exclusions": [
            "Windshield: NOT available in FL; repair-only (no replacement) in AZ, GA, ME, NY, TX",
            "Dent: dents/dings over 4 inches in diameter not covered",
            "Tire: tires with less than 3/32 inch tread depth at time of damage",
            "Key: specialty or aftermarket keys not meeting manufacturer specs",
            "Intentional damage, pre-existing conditions",
        ],
        "transfer": "Transferable to subsequent owner or lessee: $50 fee",
        "cancel": [
            [("Within 30 days: full refund less claims paid (60 days in Utah)", False)],
            [("After 30 days: pro-rata less claims paid and ", False), ("$50", True), (" cancellation fee", False)],
            [("A cancellation request cancels ALL products purchased on the contract", False)],
        ],
        "circle": None,
    },
    {
        "name": "Appearance Protection",
        "code": "HFAP", "core": False, "ev_care": False,
        "overview": "Permanently applied chemical protection system for vehicle exterior paint, interior fabric, and interior vinyl/leather surfaces. Guards against stains, finish damage, UV fading, and environmental conditions. NON-CANCELLABLE.",
        "plans": [
            {"name": "Exterior Paint Protection", "sub": "", "desc": "Protects factory painted surfaces from UV oxidation/fading, corrosion (sea/road salt, acid rain, hard water, insects, de-icing), bird droppings, tree sap, spray paint/overspray"},
            {"name": "Interior Fabric Protection", "sub": "", "desc": "Covers fabric seating and interior surfaces from stains: food/drink (soda, coffee, gum, baby food), pet waste, makeup, dye transfer, ink, crayons"},
            {"name": "Interior Vinyl & Leather Protection", "sub": "", "desc": "Covers stains or damage to vinyl/leather surfaces — benefit is professional cleaning, repair, dye, or replacement of affected area"},
        ],
        "term_coverage": [
            "120 months (10 years) only / Unlimited miles",
            "Deductible: $0",
            "Permanently applied at time of vehicle sale",
            "NON-CANCELLABLE",
        ],
        "elig_label": "Appearance Protection Eligibility",
        "eligibility": [
            "New and pre-owned Hyundai vehicles (ICE, EV, hybrid, hydrogen)",
            "Competitive makes sold at Hyundai dealerships",
            "Pre-owned vehicles: must be 5 model years old or newer",
            "Available at time of vehicle sale only",
            "Commercial and fleet vehicles excluded; professional use allowed",
        ],
        "reimbursement": [
            ("", [("Exterior: repair, repaint, or refinish of affected area per agreement limits", False)]),
            ("", [("Interior: professional cleaning, repair, dye, or replacement of affected area", False)]),
        ],
        "additional_benefits": [
            [("Rental Car Benefit: up to $50/day, 10 days max while vehicle is in for covered repair", False)],
            [("Chemical products can be ordered through HPP marketing site", False)],
        ],
        "exclusions": [
            "Pre-owned vehicles older than 5 model years",
            "Dye transfer limited to professional stain removal techniques only",
            "Pre-existing damage on used vehicles (limited to professional reconditioning)",
            "Commercial and fleet vehicles",
        ],
        "transfer": "Transferable to subsequent owner: $50 fee",
        "cancel": [[("NON-CANCELLABLE", False)]],
        "circle": None,
    },

    # ════════════════ LOSS & THEFT PROTECTION ════════════════
    {
        "name": "Guaranteed Asset Protection (GAP)",
        "code": "HFGP (Hyundai) / WFGP (Off-Make)", "core": True, "ev_care": False,
        "overview": "Waives the difference between the vehicle's actual cash value (per primary insurer) and the outstanding loan balance in the event of a total loss. Finance purchases only. Also waives up to $1,000 of the primary insurance deductible.",
        "plans": [
            {"name": "GAP", "sub": "(HFGP — Hyundai | WFGP — Off-Make)", "desc": "Waives the gap between vehicle cash value and remaining finance balance. Also waives up to $1,000 of customer's primary insurance deductible (not available in all states). No mileage limitations. No LTV limit."},
        ],
        "term_coverage": [
            "Term matches finance agreement term",
            "No mileage limitations",
            "No LTV (Loan-to-Value) limit",
            "Deductible waiver: up to $1,000 (not available in all states)",
        ],
        "elig_label": "GAP Eligibility",
        "eligibility": [
            "Finance purchases only (not available on lease)",
            "New and pre-owned Hyundai vehicles (HFGP); competitive makes (WFGP)",
            "Available at time of vehicle sale only",
            "Primary auto insurance required",
        ],
        "reimbursement": [
            ("", [("Benefit paid to lienholder; deductible waiver paid to customer", False)]),
        ],
        "additional_benefits": [
            [("No LTV limit — covers any loan amount regardless of vehicle value", False)],
            [("Waives up to $1,000 of primary insurance deductible (not available in all states)", False)],
        ],
        "exclusions": [
            "Lease agreements — use Lease-End Protection instead",
            "Total loss due to intentional acts",
            "Overdue payments and delinquent charges not covered",
        ],
        "transfer": "Non-Transferable",
        "cancel": [
            [("Within 30 days: full refund", False)],
            [("After 30 days: pro-rata refund based on time and/or mileage, less ", False), ("$75", True), (" cancellation fee", False)],
        ],
        "circle": {
            "code": "HFEG",
            "benefit": "Circle pricing available for Hyundai employees, family members, and affiliates on GAP contracts.",
            "restrictions": [
                "Must qualify under Hyundai Circle program guidelines",
            ],
        },
    },
    {
        "name": "Guaranteed Asset Protection Plus (GAP Plus)",
        "code": "HFGS (Hyundai) / WFGS (Off-Make)", "core": True, "ev_care": False,
        "overview": "All benefits of standard GAP, plus an additional $2,000 credit toward the purchase or lease of a replacement vehicle at the same dealership. Finance purchases only.",
        "plans": [
            {"name": "GAP Plus", "sub": "(HFGS — Hyundai | WFGS — Off-Make)", "desc": "All GAP benefits + $2,000 replacement vehicle credit at selling dealer (or nearest Hyundai dealer if customer relocates >50 miles). Credit applied toward purchase or lease of replacement vehicle."},
        ],
        "term_coverage": [
            "Term matches finance agreement term",
            "No mileage limitations",
            "No LTV (Loan-to-Value) limit",
            "$2,000 replacement vehicle credit at selling dealer",
        ],
        "elig_label": "GAP Plus Eligibility",
        "eligibility": [
            "Finance purchases only",
            "New and pre-owned Hyundai vehicles (HFGS); competitive makes (WFGS)",
            "Available at time of vehicle sale only",
            "Replacement credit must be used at selling dealer unless customer moves >50 miles",
        ],
        "reimbursement": [
            ("", [("GAP benefit paid to lienholder; $2,000 credit applied at dealer", False)]),
        ],
        "additional_benefits": [
            [("$2,000 credit drives repeat vehicle purchase at selling dealership", False)],
            [("Strong dealer retention tool — rewards loyal customers", False)],
        ],
        "exclusions": [
            "Lease agreements",
            "Replacement credit cannot be used if customer remains within 50 miles of selling dealer and dealer is operational",
            "Total loss due to intentional acts",
        ],
        "transfer": "Non-Transferable",
        "cancel": [
            [("Within 30 days: full refund", False)],
            [("After 30 days: pro-rata refund less ", False), ("$75", True), (" cancellation fee", False)],
        ],
        "circle": {
            "code": "HFEG",
            "benefit": "Circle pricing available for Hyundai employees, family members, and affiliates on GAP Plus contracts.",
            "restrictions": [
                "Must qualify under Hyundai Circle program guidelines",
            ],
        },
    },
    {
        "name": "Theft Protection",
        "code": "HFTP", "core": False, "ev_care": False,
        "overview": "Permanently marks the vehicle with VIN identification codes on windows and body panels, plus anti-theft decals. Provides a financial warranty benefit up to $5,000 if the vehicle is stolen and declared a total loss or not recovered. NON-CANCELLABLE.",
        "plans": [
            {"name": "Theft Protection", "sub": "", "desc": "Permanent VIN markings on windows and body panels + anti-theft deterrent decals. No activation required. Warranty benefit up to $5,000 if vehicle is stolen and not recovered or is declared a total loss."},
        ],
        "term_coverage": [
            "Lifetime of vehicle ownership",
            "No mileage limitation",
            "Warranty benefit: up to $5,000",
            "Deductible: $0",
            "NON-CANCELLABLE",
        ],
        "elig_label": "Theft Protection Eligibility",
        "eligibility": [
            "New and pre-owned Hyundai vehicles",
            "Competitive makes sold at Hyundai dealerships",
            "Available at time of vehicle sale",
        ],
        "reimbursement": [
            ("", [("Warranty benefit up to $5,000 paid to customer if vehicle stolen and total loss declared", False)]),
        ],
        "additional_benefits": [
            [("Permanent deterrent — VIN markings visible to thieves on windows and panels", False)],
            [("No activation required — protection begins immediately upon application", False)],
        ],
        "exclusions": [
            "Vehicle recovered without total loss declaration — no benefit payable",
            "Intentional acts by owner or insured",
        ],
        "transfer": "Transferable to subsequent owner: no transfer fee",
        "cancel": [[("NON-CANCELLABLE", False)]],
        "circle": None,
    },
]

# ── Main ─────────────────────────────────────────────────────────────────────

def main():
    prs = new_prs()

    build_cover(prs)
    build_toc(prs)

    # MECHANICAL COVERAGE
    build_section_divider(prs, "MECHANICAL COVERAGE",
        "VSP ICE  |  Competitive Makes  |  Original Owner  |  CUV Wrap  |  High Mileage  |  EV Care  |  Livery  |  Wear Protection",
        color=C_NAVY)
    mech = [
        "Vehicle Service Protection (ICE)",
        "Vehicle Service Protection — Competitive Makes (ICE)",
        "Vehicle Service Protection Original Owner (ICE)",
        "Certified Used Vehicle (CUV) Wrap",
        "High Mileage Vehicle Service Protection (ICE)",
        "EV Care Vehicle Service Protection",
        "EV Care VSP — Competitive Makes",
        "EV Care VSP — Original Owner",
        "EV Care Certified Used Vehicle (CUV) Wrap",
        "VSP Livery Coverage (ICE & EV)",
        "Wear Protection (ICE & EV)",
    ]
    for d in PRODUCTS:
        if d["name"] in mech:
            build_product_slide(prs, d)

    # MAINTENANCE PLANS
    build_section_divider(prs, "MAINTENANCE PLANS",
        "Basic Wrap  |  Scheduled Wrap  |  Pre-Paid Maintenance  |  Competitive Makes  |  EV Care Maintenance",
        color=C_MID)
    maint = [
        "Maintenance Basic Wrap (ICE)",
        "Maintenance Scheduled Wrap (ICE)",
        "Pre-Paid Maintenance (ICE)",
        "Pre-Paid Maintenance — Competitive Makes (ICE)",
        "EV Care Maintenance",
    ]
    for d in PRODUCTS:
        if d["name"] in maint:
            build_product_slide(prs, d)

    # DAMAGE CARE
    build_section_divider(prs, "DAMAGE CARE",
        "Lease-End Protection  |  Multi-Coverage Protection  |  Appearance Protection",
        color=C_TEAL)
    damage = [
        "Lease-End Protection",
        "Multi-Coverage Protection",
        "Appearance Protection",
    ]
    for d in PRODUCTS:
        if d["name"] in damage:
            build_product_slide(prs, d)

    # LOSS & THEFT
    build_section_divider(prs, "LOSS & THEFT PROTECTION",
        "GAP  |  GAP Plus  |  Theft Protection",
        color=C_PURPLE)
    loss = [
        "Guaranteed Asset Protection (GAP)",
        "Guaranteed Asset Protection Plus (GAP Plus)",
        "Theft Protection",
    ]
    for d in PRODUCTS:
        if d["name"] in loss:
            build_product_slide(prs, d)

    out = "/Users/justin/Downloads/HPP Product Framework - Training Deck.pptx"
    prs.save(out)
    print(f"Saved: {out}")
    print(f"Total slides: {len(prs.slides)}")

if __name__ == "__main__":
    main()
