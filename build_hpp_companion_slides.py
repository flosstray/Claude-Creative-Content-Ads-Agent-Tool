"""
HPP Companion Slides — Off-Make / Competitive Make Versions (Option B)
Products: VSP Competitive Makes (WFVI) + EV Care VSP Competitive Makes (WFVE)
Format: Plain text, Aptos font — matches Training Deck style
Output: /Users/justin/Downloads/HPP Companion Slides - Off-Make.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches as I

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

F       = "Aptos"
C_BLACK  = RGBColor(0x1A, 0x1A, 0x1A)
C_NAVY   = RGBColor(0x00, 0x20, 0x60)
C_GREY   = RGBColor(0x55, 0x55, 0x55)
C_MID    = RGBColor(0x00, 0x4C, 0x97)
C_ORANGE = RGBColor(0xE8, 0x6F, 0x1E)
C_WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
C_NONCORE = RGBColor(0x6D, 0x6D, 0x6D)
C_AMBER  = RGBColor(0xC8, 0x6A, 0x00)   # off-make callout accent


def new_prs():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = Inches(7.5)
    return prs


def blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def txb(slide, text, left, top, width, height,
        size=11, bold=False, color=C_BLACK,
        align=PP_ALIGN.LEFT, italic=False, word_wrap=True):
    box = slide.shapes.add_textbox(I(left), I(top), I(width), I(height))
    tf  = box.text_frame
    tf.word_wrap = word_wrap
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name  = F
    r.font.size  = Pt(size)
    r.font.bold  = bold
    r.font.italic = italic
    r.font.color.rgb = color
    return box


def txb_multi(slide, paras, left, top, width, height, word_wrap=True):
    box = slide.shapes.add_textbox(I(left), I(top), I(width), I(height))
    tf  = box.text_frame
    tf.word_wrap = word_wrap
    first = True
    for pd in paras:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = pd.get("align", PP_ALIGN.LEFT)
        if pd.get("space_before"):
            p.space_before = Pt(pd["space_before"])
        if pd.get("space_after"):
            p.space_after  = Pt(pd["space_after"])
        text = pd.get("text", "")
        if pd.get("bullet"):
            text = "  •  " + text
        r = p.add_run()
        r.text = text
        r.font.name   = F
        r.font.size   = Pt(pd.get("size", 11))
        r.font.bold   = pd.get("bold", False)
        r.font.italic = pd.get("italic", False)
        r.font.color.rgb = pd.get("color", C_BLACK)
    return box


def hl(paras, text, size=9.5, bold=False, color=C_BLACK,
        space_before=3, space_after=0, bullet=False, italic=False):
    paras.append({
        "text": text, "size": size, "bold": bold, "color": color,
        "space_before": space_before, "space_after": space_after,
        "bullet": bullet, "italic": italic, "align": PP_ALIGN.LEFT,
    })


def build_companion_slide(prs, d):
    """
    Companion slide layout (same 3-col structure as training deck):
      Header: product name + code + NON-CORE badge
      Reference banner: 'Companion to [parent product]'
      Col 1: Available Plans + Key Differences vs. parent + Terms
      Col 2: Eligibility (with off-make-specific rules called out)
      Col 3: Exclusions + Transfer + Cancellation
    """
    slide = blank_slide(prs)

    name   = d["name"]
    code   = d["code"]
    parent = d.get("parent_product", "")

    # ── Header ───────────────────────────────────────────────────────────────
    txb(slide, name, 0.6, 0.22, 10.5, 0.55,
        size=22, bold=True, color=C_NAVY)
    txb(slide, f"{code}", 0.6, 0.72, 5.5, 0.32,
        size=10, bold=True, color=C_NONCORE)
    txb(slide, "NON-CORE  |  Off-Make / Competitive Makes", 6.5, 0.72, 6.2, 0.32,
        size=10, bold=False, color=C_NONCORE, align=PP_ALIGN.RIGHT)

    # ── Companion reference banner ────────────────────────────────────────────
    txb(slide, f"─" * 120, 0.6, 1.0, 12.1, 0.2, size=6, color=RGBColor(0xCC,0xCC,0xCC))
    txb(slide,
        f"COMPANION SLIDE  |  Off-Make version of:  {parent}",
        0.6, 1.08, 12.1, 0.32,
        size=9.5, bold=False, italic=True, color=C_AMBER)

    # ── Overview ─────────────────────────────────────────────────────────────
    overview = d.get("overview", "")
    if overview:
        paras = []
        hl(paras, "PRODUCT OVERVIEW", size=9, bold=True, color=C_GREY, space_before=0)
        hl(paras, overview, size=10.5, bold=False, italic=True, color=C_BLACK, space_before=2)
        txb_multi(slide, paras, 0.6, 1.36, 12.1, 0.65)

    # ── Three-column layout ──────────────────────────────────────────────────
    C1X, C1W = 0.60, 4.30
    C2X, C2W = 5.10, 4.00
    C3X, C3W = 9.30, 3.83
    ROW_Y = 1.98

    # ── COL 1: Plans + Key Differences + Terms ────────────────────────────────
    paras = []
    hl(paras, "AVAILABLE PLANS", size=9.5, bold=True, color=C_NAVY, space_before=0)
    for pl in d.get("plans", []):
        pl_name = pl["name"]
        if pl.get("sub"):
            pl_name += f"  {pl['sub']}"
        hl(paras, pl_name, size=10, bold=True, color=C_BLACK, space_before=5, bullet=True)
        if pl.get("desc"):
            hl(paras, pl["desc"], size=9, bold=False, color=C_GREY, space_before=1)

    # Key differences vs. parent
    diffs = d.get("key_differences", [])
    if diffs:
        hl(paras, "", size=4, space_before=4)
        hl(paras, "KEY DIFFERENCES vs. HPP CORE VERSION", size=9.5, bold=True,
           color=C_AMBER, space_before=2)
        for diff in diffs:
            hl(paras, diff, size=9.5, bold=False, color=C_BLACK, space_before=2, bullet=True)

    hl(paras, "", size=4, space_before=4)
    hl(paras, "TERM & COVERAGE", size=9.5, bold=True, color=C_NAVY, space_before=2)
    for line in d.get("term_coverage", []):
        hl(paras, line, size=9.5, bold=False, color=C_BLACK, space_before=2, bullet=True)

    txb_multi(slide, paras, C1X, ROW_Y, C1W, 5.3)

    # ── COL 2: Eligibility + Reimbursement ───────────────────────────────────
    paras = []
    hl(paras, "VEHICLE ELIGIBILITY", size=9.5, bold=True, color=C_NAVY, space_before=0)
    for line in d.get("eligibility", []):
        hl(paras, line, size=9.5, bold=False, color=C_BLACK, space_before=2, bullet=True)

    reimb = d.get("reimbursement", [])
    if reimb:
        hl(paras, "", size=4, space_before=4)
        hl(paras, "CLAIM REIMBURSEMENT", size=9.5, bold=True, color=C_NAVY, space_before=2)
        for row in reimb:
            label, parts = row
            full_text = label + " ".join(t for t, _ in parts)
            hl(paras, full_text.strip(), size=9.5, bold=False, color=C_BLACK,
               space_before=2, bullet=True)

    add_ben = d.get("additional_benefits", [])
    if add_ben:
        hl(paras, "", size=4, space_before=4)
        hl(paras, "ADDITIONAL BENEFITS", size=9.5, bold=True, color=C_NAVY, space_before=2)
        for row in add_ben:
            full_text = " ".join(t for t, _ in row)
            hl(paras, full_text.strip(), size=9.5, bold=False, color=C_BLACK,
               space_before=2, bullet=True)

    txb_multi(slide, paras, C2X, ROW_Y, C2W, 5.3)

    # ── COL 3: Exclusions + Transfer + Cancellation ───────────────────────────
    paras = []
    excl = d.get("exclusions", [])
    if excl:
        hl(paras, "EXCLUSIONS", size=9.5, bold=True, color=C_NAVY, space_before=0)
        for line in excl:
            hl(paras, line, size=9, bold=False, color=C_BLACK, space_before=2, bullet=True)

    transfer = d.get("transfer", "")
    if transfer:
        hl(paras, "", size=4, space_before=4)
        hl(paras, "TRANSFER", size=9.5, bold=True, color=C_NAVY, space_before=2)
        hl(paras, transfer, size=9.5, bold=False, color=C_BLACK, space_before=2, bullet=True)

    cancel = d.get("cancel", [])
    if cancel:
        hl(paras, "", size=4, space_before=4)
        hl(paras, "CANCELLATION", size=9.5, bold=True, color=C_NAVY, space_before=2)
        for row in cancel:
            full_text = " ".join(t for t, _ in row)
            hl(paras, full_text.strip(), size=9.5, bold=False, color=C_BLACK,
               space_before=2, bullet=True)

    txb_multi(slide, paras, C3X, ROW_Y, C3W, 5.3)


# ── Companion slide data ─────────────────────────────────────────────────────

COMPANIONS = [
    {
        "name": "VSP — Competitive Makes (ICE)",
        "code": "WFVI",
        "parent_product": "Vehicle Service Protection (ICE)  |  HFVI",
        "overview": (
            "PowerProtect mechanical breakdown coverage for pre-owned competitive make (non-Hyundai) "
            "ICE vehicles sold at Hyundai dealerships. Same core structure as HPP VSP — key differences "
            "are vehicle eligibility, plan availability, and additive (not original in-service) term structure."
        ),
        "plans": [
            {
                "name": "Powertrain",
                "sub": "(Stated Component)",
                "desc": "Engine, Transmission, Drive Axle incl. CV joints",
            },
            {
                "name": "Gold",
                "sub": "(Stated Component)",
                "desc": (
                    "Powertrain + front/rear suspension (incl. shocks), A/C, fuel system, "
                    "electrical system"
                ),
            },
            {
                "name": "Platinum",
                "sub": "(Exclusionary)",
                "desc": "All covered parts except listed exclusions",
            },
            {
                "name": "Platinum Plus",
                "sub": "(Florida Only)",
                "desc": "Platinum + headlamps, belts and hoses, electrical",
            },
        ],
        "key_differences": [
            "HIGH TECHNOLOGY plan NOT available (competitive makes only)",
            "Term is ADDITIVE — starts from contract sale date/mileage, not original in-service date",
            "Max term: 10 years / 120,000 miles (vs. 150,000 mi on HFVI)",
            "Pre-owned vehicles only — not available on new competitive make vehicles",
            "Vehicle must be 9 model years old or less with less than 120,000 miles at purchase",
            "UCI (Used Car Inspection) required at time of sale",
            "No Circle program pricing available",
            "Genuine Hyundai OEM parts NOT used — like kind and quality parts",
        ],
        "term_coverage": [
            "Up to 10 years / 120,000 miles (additive from contract sale date)",
            "Deductible: $0, $100, or $250 per repair visit",
            "24-hr Emergency Roadside Assistance included",
            "Rental Car/Rideshare: up to $55/day, 10 days max",
            "Trip Interruption: up to $300/day, 5 days max ($1,500 max)",
        ],
        "eligibility": [
            "Pre-owned competitive make (non-Hyundai) ICE vehicles sold at Hyundai dealerships",
            "9 model years old or less at time of contract purchase",
            "Less than 120,000 miles on odometer at time of contract purchase",
            "January 1st is the anniversary date for determining vehicle age",
            "If purchased post-sale: must have at least 1 month and 1,000 miles of manufacturer warranty remaining",
            "UCI (Used Car Inspection) required for all pre-owned VSP sales",
            "Salvage, junk, or Buy-Back titled vehicles not eligible",
            "Non-U.S. spec models not eligible",
            "Permitted Commercial Use allowed (single driver, rideshare, light service)",
            "Prohibited: hauling, livery, fleet/pool, daily rentals, towing, government/military",
        ],
        "reimbursement": [
            ("", [("Like kind and quality parts; published labor rate (not OEM Hyundai parts)", False)]),
            ("", [("Repair Order submitted to administrator after authorization and repair", False)]),
        ],
        "additional_benefits": [
            [("Diagnostic labor covered when repair is covered", False)],
            [("Necessary fluids/lubricants for covered repair included", False)],
            [("Platinum Plus CANNOT be sold with Wear Protection", False)],
        ],
        "exclusions": [
            "New competitive make vehicles — pre-owned only",
            "High Technology plan not available",
            "Hyundai, Kia, or Genesis vehicles — use HFVI instead",
            "Vehicles over 9 model years old or over 120,000 miles at purchase",
            "Maintenance services, wear items",
            "Accidental damage, collision, vandalism, theft",
            "Prohibited commercial use",
        ],
        "transfer": "Transferable to subsequent owner: $75 fee",
        "cancel": [
            [("Within 30 days: full refund less claims paid", False)],
            [("After 30 days: pro-rata less claims paid and $75 cancellation fee", False)],
        ],
    },
    {
        "name": "EV Care VSP — Competitive Makes",
        "code": "WFVE",
        "parent_product": "EV Care Vehicle Service Protection  |  HFVE",
        "overview": (
            "PowerProtect EV mechanical breakdown coverage for pre-owned competitive make Electric vehicles "
            "sold at Hyundai dealerships. Platinum exclusionary plan only. Key differences from HPP EV Care VSP: "
            "additive term, lower max mileage, pre-owned only, and broader Platinum coverage scope including "
            "competitive-make-specific EV components."
        ),
        "plans": [
            {
                "name": "Platinum",
                "sub": "(Exclusionary — Only Plan Available)",
                "desc": (
                    "All covered parts except listed exclusions. Covers: Battery & High Technology components, "
                    "Electronic Air Compressor, Electric Power Control Unit, Transmission, Powertrain components, "
                    "any components originally manufactured/installed by the vehicle OEM, A/C Refrigerant Charge, "
                    "12V Battery (defective within first 3 years), Climate Control, Shocks/Suspension, "
                    "Steering, EV Regenerative Brakes"
                ),
            },
            {
                "name": "Platinum Plus",
                "sub": "(Florida Only)",
                "desc": "Platinum + headlamps, belts and hoses, electrical",
            },
        ],
        "key_differences": [
            "PLATINUM ONLY — Battery and High Technology standalone plans NOT available",
            "Term is ADDITIVE — starts from contract sale date/mileage, not original in-service date",
            "Max term: 12 years / 120,000 miles (vs. 200,000 mi on HFVE)",
            "Pre-owned vehicles only — 9 model years old or less",
            "Max odometer at purchase: 108,001 miles (vs. no stated limit on HFVE)",
            "UCI (Used Car Inspection) required at time of sale",
            "No Circle program pricing available",
            "Platinum scope covers competitive EV components (not just Hyundai-spec parts)",
        ],
        "term_coverage": [
            "Up to 12 years / 120,000 miles (additive from contract sale date)",
            "Deductible: $0, $100, or $250 per repair visit",
            "24-hr Emergency Roadside Assistance included",
            "Rental Car/Rideshare: up to $55/day, 10 days max",
            "Trip Interruption: up to $300/day, 5 days max ($1,500 max)",
        ],
        "eligibility": [
            "Pre-owned competitive make (non-Hyundai) Electric vehicles sold at Hyundai dealerships",
            "9 model years old or less at time of contract purchase",
            "Less than 108,001 miles on odometer at time of contract purchase",
            "January 1st is the anniversary date for determining vehicle age",
            "If purchased post-sale: must have at least 1 month and 1,000 miles of manufacturer warranty remaining",
            "UCI (Used Car Inspection) required for all pre-owned EV VSP sales",
            "Salvage, junk, or Buy-Back titled vehicles not eligible",
            "Non-U.S. spec models not eligible",
            "Permitted Commercial Use allowed (single driver, rideshare, light service)",
            "Light Duty Commercial Use NOT available with High Technology Coverage Plan",
        ],
        "reimbursement": [
            ("", [("Like kind and quality parts; published labor rate", False)]),
            ("", [("Repair Order submitted to administrator after authorization and repair", False)]),
        ],
        "additional_benefits": [
            [("Diagnostic labor covered when repair is covered", False)],
            [("Necessary fluids/lubricants for covered repair included", False)],
            [("Platinum Plus CANNOT be sold with Wear Protection", False)],
        ],
        "exclusions": [
            "New competitive make EVs — pre-owned only",
            "Battery and High Technology standalone plans not available",
            "Hyundai, Kia, or Genesis EVs — use HFVE instead",
            "Vehicles over 9 model years old or over 108,000 miles at purchase",
            "Maintenance services, wear items",
            "Accidental damage, collision, vandalism, theft",
            "Prohibited commercial use",
        ],
        "transfer": "Transferable to subsequent owner: $75 fee",
        "cancel": [
            [("Within 30 days: full refund less claims paid", False)],
            [("After 30 days: pro-rata less claims paid and $75 cancellation fee", False)],
        ],
    },
]


def main():
    prs = new_prs()
    for d in COMPANIONS:
        build_companion_slide(prs, d)
    out = "/Users/justin/Downloads/HPP Companion Slides - Off-Make.pptx"
    prs.save(out)
    print(f"Saved: {out}")
    print(f"Total slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
