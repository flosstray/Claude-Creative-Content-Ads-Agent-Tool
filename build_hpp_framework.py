"""
HPP Product Framework PowerPoint Builder
Recreates the HPP product layout following the Dealer Operations Manual structure,
using the Product Knowledge 360 / HPP deck as formatting inspiration.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Brand colours ──────────────────────────────────────────────────────────────
HYUNDAI_BLUE   = RGBColor(0x00, 0x2C, 0x5F)   # deep navy
HYUNDAI_SKY    = RGBColor(0x00, 0x7A, 0xC1)   # Hyundai accent blue
EV_TEAL        = RGBColor(0x00, 0x98, 0x91)   # teal for EV Care products
CORE_GOLD      = RGBColor(0xC8, 0xA0, 0x32)   # gold badge for CORE
NON_CORE_GREY  = RGBColor(0x6D, 0x6D, 0x6D)   # grey badge for NON-CORE
WHITE          = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GREY_BG  = RGBColor(0xF4, 0xF6, 0xF9)
DARK_TEXT      = RGBColor(0x1A, 0x1A, 0x2E)
MED_GREY       = RGBColor(0x9E, 0x9E, 0x9E)
SECTION_LINE   = RGBColor(0xDD, 0xE3, 0xEA)
CIRCLE_ORANGE  = RGBColor(0xE8, 0x6F, 0x1E)


def add_rect(slide, left, top, width, height, fill_rgb=None, line_rgb=None, line_width=None):
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    shape = slide.shapes.add_shape(1, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.line.fill.background()
    if fill_rgb:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_rgb
    else:
        shape.fill.background()
    if line_rgb:
        shape.line.color.rgb = line_rgb
        if line_width:
            shape.line.width = Pt(line_width)
    else:
        shape.line.fill.background()
    return shape


def add_text_box(slide, text, left, top, width, height,
                 font_size=10, bold=False, color=DARK_TEXT,
                 italic=False, align=PP_ALIGN.LEFT, wrap=True):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    txBox.word_wrap = wrap
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox


def set_slide_bg(slide, rgb):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = rgb


def add_multi_para_textbox(slide, paras, left, top, width, height):
    """paras: list of (text, font_size, bold, color, indent_level)"""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    txBox.word_wrap = True
    tf = txBox.text_frame
    tf.word_wrap = True
    first = True
    for (text, font_size, bold, color, indent) in paras:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.level = indent
        run = p.add_run()
        run.text = text
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.color.rgb = color
    return txBox


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE BUILDERS
# ─────────────────────────────────────────────────────────────────────────────

def build_cover(prs):
    slide_layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(slide_layout)
    set_slide_bg(slide, HYUNDAI_BLUE)

    # Full-width gradient bar at top (simulate with rectangle)
    add_rect(slide, 0, 0, 13.33, 0.08, fill_rgb=HYUNDAI_SKY)

    # White logo text area
    add_text_box(slide, "HYUNDAI PROTECTION PLAN",
                 0.5, 1.8, 12.3, 0.7,
                 font_size=36, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    add_text_box(slide, "Product Framework Guide",
                 0.5, 2.6, 12.3, 0.5,
                 font_size=20, bold=False, color=RGBColor(0xB0, 0xC8, 0xE8), align=PP_ALIGN.CENTER)

    # Divider
    add_rect(slide, 2.5, 3.25, 8.33, 0.04, fill_rgb=HYUNDAI_SKY)

    add_text_box(slide, "Core & Non-Core Products  |  F&I Dealer Reference",
                 0.5, 3.45, 12.3, 0.4,
                 font_size=12, bold=False, color=RGBColor(0xB0, 0xC8, 0xE8), align=PP_ALIGN.CENTER)

    # Categories listed
    cats = [
        ("MECHANICAL COVERAGE", 4.3),
        ("MAINTENANCE PLANS", 5.0),
        ("DAMAGE CARE", 5.7),
        ("LOSS & THEFT PROTECTION", 6.4),
    ]
    for label, y in cats:
        add_rect(slide, 1.5, y, 10.33, 0.45, fill_rgb=RGBColor(0x00, 0x44, 0x8A))
        add_text_box(slide, label, 1.5, y + 0.05, 10.33, 0.35,
                     font_size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # Footer
    add_text_box(slide, "For Internal Use Only  |  Dealer Operations Manual Reference",
                 0.5, 7.1, 12.3, 0.3,
                 font_size=9, color=MED_GREY, align=PP_ALIGN.CENTER)

    # CORE legend
    add_rect(slide, 0.5, 6.6, 1.1, 0.3, fill_rgb=CORE_GOLD)
    add_text_box(slide, "● CORE", 0.5, 6.6, 1.1, 0.3,
                 font_size=9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_rect(slide, 1.8, 6.6, 1.5, 0.3, fill_rgb=NON_CORE_GREY)
    add_text_box(slide, "● NON-CORE", 1.8, 6.6, 1.5, 0.3,
                 font_size=9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    add_rect(slide, 3.5, 6.6, 2.0, 0.3, fill_rgb=CIRCLE_ORANGE)
    add_text_box(slide, "● CIRCLE PROGRAM", 3.5, 6.6, 2.0, 0.3,
                 font_size=9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)


def build_section_divider(prs, section_title, subtitle="", bg=HYUNDAI_BLUE):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_bg(slide, bg)

    add_rect(slide, 0, 0, 13.33, 0.08, fill_rgb=HYUNDAI_SKY)

    add_text_box(slide, "HYUNDAI PROTECTION PLAN",
                 0.5, 2.5, 12.3, 0.5,
                 font_size=14, bold=False, color=RGBColor(0xB0, 0xC8, 0xE8),
                 align=PP_ALIGN.CENTER)

    add_text_box(slide, section_title,
                 0.5, 3.1, 12.3, 0.9,
                 font_size=40, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    if subtitle:
        add_rect(slide, 2.5, 4.15, 8.33, 0.04, fill_rgb=HYUNDAI_SKY)
        add_text_box(slide, subtitle,
                     0.5, 4.3, 12.3, 0.4,
                     font_size=13, color=RGBColor(0xB0, 0xC8, 0xE8),
                     align=PP_ALIGN.CENTER)

    add_text_box(slide, "For Internal Use Only",
                 0.5, 7.1, 12.3, 0.3,
                 font_size=9, color=MED_GREY, align=PP_ALIGN.CENTER)


# Helper to add badge
def add_badge(slide, label, left, top, width=1.0, color=CORE_GOLD):
    add_rect(slide, left, top, width, 0.22, fill_rgb=color)
    add_text_box(slide, label, left, top + 0.01, width, 0.2,
                 font_size=8, bold=True, color=WHITE, align=PP_ALIGN.CENTER)


def add_circle_badge(slide, left, top):
    add_rect(slide, left, top, 1.3, 0.22, fill_rgb=CIRCLE_ORANGE)
    add_text_box(slide, "● CIRCLE PROGRAM", left, top + 0.01, 1.3, 0.2,
                 font_size=7.5, bold=True, color=WHITE, align=PP_ALIGN.CENTER)


# ─────────────────────────────────────────────────────────────────────────────
# PRODUCT CARD SLIDE  (standard two-column layout)
# ─────────────────────────────────────────────────────────────────────────────

def build_product_slide(prs, product_data):
    """
    product_data keys:
      name, subtitle, code, core (bool), ev_care (bool), circle (bool),
      overview,
      plans: list of {name, desc}
      features: list of str
      components: list of str  (optional)
      terms_mileage: str
      eligibility: list of str
      exclusions: list of str
      transfer_fee: str
      cancel_fee: str
      deductible: str
      additional_benefits: list of str
    """
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_bg(slide, WHITE)

    d = product_data
    is_ev = d.get("ev_care", False)
    header_color = EV_TEAL if is_ev else HYUNDAI_BLUE

    # ── Top header bar ──────────────────────────────────────────────────────
    add_rect(slide, 0, 0, 13.33, 1.0, fill_rgb=header_color)
    add_text_box(slide, d["name"], 0.25, 0.08, 9.5, 0.55,
                 font_size=20, bold=True, color=WHITE)
    if d.get("subtitle"):
        add_text_box(slide, d["subtitle"], 0.25, 0.58, 9.0, 0.35,
                     font_size=11, color=RGBColor(0xB0, 0xD8, 0xFF))

    # Code badge (top right)
    code_x = 10.5
    if d.get("code"):
        add_rect(slide, code_x, 0.08, 2.6, 0.38, fill_rgb=RGBColor(0x00, 0x1A, 0x40))
        add_text_box(slide, f"Code: {d['code']}", code_x + 0.1, 0.1, 2.4, 0.35,
                     font_size=9, bold=True, color=RGBColor(0x90, 0xC0, 0xFF))

    # CORE / NON-CORE badge
    badge_color = CORE_GOLD if d.get("core", True) else NON_CORE_GREY
    badge_label = "CORE" if d.get("core", True) else "NON-CORE"
    add_badge(slide, badge_label, code_x, 0.55, width=1.2, color=badge_color)

    if d.get("circle"):
        add_circle_badge(slide, code_x + 1.3, 0.55)

    # ── Overview strip ──────────────────────────────────────────────────────
    add_rect(slide, 0, 1.0, 13.33, 0.38, fill_rgb=LIGHT_GREY_BG)
    add_text_box(slide, d.get("overview", ""),
                 0.25, 1.03, 13.0, 0.35,
                 font_size=9.5, italic=True, color=DARK_TEXT)

    # ── Two column layout starting at y=1.45 ───────────────────────────────
    # Left column: Plans + Features & Benefits + Components
    # Right column: Terms/Mileage | Eligibility | Exclusions | Fees

    COL1_X = 0.2
    COL1_W = 6.3
    COL2_X = 6.85
    COL2_W = 6.2
    COL_TOP = 1.48
    ROW_H   = 0.22

    y_l = COL_TOP
    y_r = COL_TOP

    # ── LEFT: Plans ─────────────────────────────────────────────────────────
    if d.get("plans"):
        add_rect(slide, COL1_X, y_l, COL1_W, 0.26, fill_rgb=header_color)
        add_text_box(slide, "PLANS & COVERAGE LEVELS",
                     COL1_X + 0.1, y_l + 0.03, COL1_W - 0.2, 0.22,
                     font_size=9, bold=True, color=WHITE)
        y_l += 0.28

        for plan in d["plans"]:
            add_rect(slide, COL1_X, y_l, COL1_W, 0.04, fill_rgb=SECTION_LINE)
            y_l += 0.06
            # Plan name
            add_text_box(slide, f"▸ {plan['name']}",
                         COL1_X + 0.05, y_l, COL1_W * 0.35, ROW_H,
                         font_size=8.5, bold=True, color=header_color)
            # Plan desc
            add_text_box(slide, plan.get("desc", ""),
                         COL1_X + COL1_W * 0.35, y_l, COL1_W * 0.65, ROW_H * 2,
                         font_size=8, color=DARK_TEXT)
            y_l += ROW_H * 1.5

    # ── LEFT: Features & Benefits ───────────────────────────────────────────
    if d.get("features"):
        y_l += 0.05
        add_rect(slide, COL1_X, y_l, COL1_W, 0.26, fill_rgb=header_color)
        add_text_box(slide, "FEATURES & BENEFITS",
                     COL1_X + 0.1, y_l + 0.03, COL1_W - 0.2, 0.22,
                     font_size=9, bold=True, color=WHITE)
        y_l += 0.28

        for feat in d["features"]:
            add_text_box(slide, f"• {feat}",
                         COL1_X + 0.1, y_l, COL1_W - 0.15, ROW_H * 1.6,
                         font_size=8, color=DARK_TEXT)
            y_l += ROW_H * 1.1

    # ── LEFT: Components (if any) ────────────────────────────────────────────
    if d.get("components"):
        y_l += 0.05
        add_rect(slide, COL1_X, y_l, COL1_W, 0.26, fill_rgb=RGBColor(0x1A, 0x4A, 0x7A))
        add_text_box(slide, "KEY COMPONENTS COVERED",
                     COL1_X + 0.1, y_l + 0.03, COL1_W - 0.2, 0.22,
                     font_size=9, bold=True, color=WHITE)
        y_l += 0.28
        # Two column within column
        half = len(d["components"]) // 2 + len(d["components"]) % 2
        for i, comp in enumerate(d["components"]):
            cx = COL1_X + 0.1 if i < half else COL1_X + COL1_W / 2 + 0.05
            cy = COL_TOP + 0.28 + (i if i < half else i - half) * (ROW_H * 0.95)
            # Overwrite y_l only for the last of left column
            if i == 0:
                cy = y_l
            elif i < half:
                cy = y_l + (i) * (ROW_H * 0.95)
            else:
                cy = y_l + (i - half) * (ROW_H * 0.95)
            add_text_box(slide, f"• {comp}", cx, cy,
                         COL1_W / 2 - 0.15, ROW_H * 1.4,
                         font_size=7.5, color=DARK_TEXT)
        y_l += max(half, len(d["components"]) - half) * (ROW_H * 0.95)

    # ── RIGHT: Terms & Mileage ──────────────────────────────────────────────
    if d.get("terms_mileage"):
        add_rect(slide, COL2_X, y_r, COL2_W, 0.26, fill_rgb=header_color)
        add_text_box(slide, "TERMS & MILEAGE",
                     COL2_X + 0.1, y_r + 0.03, COL2_W - 0.2, 0.22,
                     font_size=9, bold=True, color=WHITE)
        y_r += 0.28
        add_rect(slide, COL2_X, y_r, COL2_W, 0.04, fill_rgb=SECTION_LINE)
        y_r += 0.06
        add_text_box(slide, d["terms_mileage"],
                     COL2_X + 0.1, y_r, COL2_W - 0.2, 0.6,
                     font_size=8.5, bold=True, color=DARK_TEXT)
        y_r += 0.68

    if d.get("deductible"):
        add_text_box(slide, f"Deductible: {d['deductible']}",
                     COL2_X + 0.1, y_r, COL2_W - 0.2, ROW_H,
                     font_size=8.5, color=DARK_TEXT)
        y_r += ROW_H + 0.04

    # ── RIGHT: Eligibility ──────────────────────────────────────────────────
    if d.get("eligibility"):
        y_r += 0.04
        add_rect(slide, COL2_X, y_r, COL2_W, 0.26, fill_rgb=RGBColor(0x1A, 0x5C, 0x3A))
        add_text_box(slide, "ELIGIBILITY",
                     COL2_X + 0.1, y_r + 0.03, COL2_W - 0.2, 0.22,
                     font_size=9, bold=True, color=WHITE)
        y_r += 0.28
        for elig in d["eligibility"]:
            add_text_box(slide, f"• {elig}",
                         COL2_X + 0.1, y_r, COL2_W - 0.2, ROW_H * 1.7,
                         font_size=8, color=DARK_TEXT)
            y_r += ROW_H * 1.15

    # ── RIGHT: Exclusions ───────────────────────────────────────────────────
    if d.get("exclusions"):
        y_r += 0.06
        add_rect(slide, COL2_X, y_r, COL2_W, 0.26, fill_rgb=RGBColor(0x8B, 0x1A, 0x1A))
        add_text_box(slide, "KEY EXCLUSIONS",
                     COL2_X + 0.1, y_r + 0.03, COL2_W - 0.2, 0.22,
                     font_size=9, bold=True, color=WHITE)
        y_r += 0.28
        for excl in d["exclusions"]:
            add_text_box(slide, f"✕ {excl}",
                         COL2_X + 0.1, y_r, COL2_W - 0.2, ROW_H * 1.7,
                         font_size=8, color=DARK_TEXT)
            y_r += ROW_H * 1.1

    # ── RIGHT: Transfer & Cancellation ─────────────────────────────────────
    y_r += 0.1
    add_rect(slide, COL2_X, y_r, COL2_W, 0.26, fill_rgb=RGBColor(0x4A, 0x3A, 0x00))
    add_text_box(slide, "TRANSFER / CANCELLATION",
                 COL2_X + 0.1, y_r + 0.03, COL2_W - 0.2, 0.22,
                 font_size=9, bold=True, color=WHITE)
    y_r += 0.28
    if d.get("transfer_fee"):
        add_text_box(slide, f"Transfer Fee: {d['transfer_fee']}",
                     COL2_X + 0.1, y_r, COL2_W - 0.2, ROW_H,
                     font_size=8.5, color=DARK_TEXT)
        y_r += ROW_H
    if d.get("cancel_fee"):
        add_text_box(slide, f"Cancellation: {d['cancel_fee']}",
                     COL2_X + 0.1, y_r, COL2_W - 0.2, ROW_H * 2,
                     font_size=8.5, color=DARK_TEXT)
        y_r += ROW_H * 1.2

    # ── RIGHT: Additional Benefits ──────────────────────────────────────────
    if d.get("additional_benefits"):
        y_r += 0.06
        add_rect(slide, COL2_X, y_r, COL2_W, 0.26, fill_rgb=RGBColor(0x1A, 0x3A, 0x6A))
        add_text_box(slide, "ADDITIONAL BENEFITS",
                     COL2_X + 0.1, y_r + 0.03, COL2_W - 0.2, 0.22,
                     font_size=9, bold=True, color=WHITE)
        y_r += 0.28
        for ben in d["additional_benefits"]:
            add_text_box(slide, f"✓ {ben}",
                         COL2_X + 0.1, y_r, COL2_W - 0.2, ROW_H * 1.7,
                         font_size=8, color=DARK_TEXT)
            y_r += ROW_H * 1.1

    # ── Footer ───────────────────────────────────────────────────────────────
    add_rect(slide, 0, 7.3, 13.33, 0.2, fill_rgb=LIGHT_GREY_BG)
    add_text_box(slide,
                 "Some terms vary by state and/or contract revision date. For illustrative purposes only. See applicable product agreement for details.",
                 0.2, 7.32, 13.0, 0.18,
                 font_size=6.5, color=MED_GREY, italic=True)


# ─────────────────────────────────────────────────────────────────────────────
# PRODUCT DATA
# ─────────────────────────────────────────────────────────────────────────────

ADDITIONAL_BENEFITS_MECH = [
    "24-Hour Emergency Roadside Assistance (towing, jump start, flat tire, lockout, fuel delivery)",
    "Rental Car Reimbursement: up to $55/day, 10 days max",
    "Trip Interruption Benefit: up to $300/day, 5 days max ($1,500 total; not available in NY)",
    "Diagnostic labor covered when repair is covered",
    "Replacement fluids covered in conjunction with a covered repair",
]

ADDITIONAL_BENEFITS_EV = [
    "24-Hour Emergency Roadside Assistance (towing, jump start, flat tire, lockout, fuel delivery)",
    "Rental Car / Ride Share Reimbursement: up to $55/day, 10 days max",
    "Trip Interruption Benefit: up to $300/day, 5 days max ($1,500 total; not available in NY)",
    "Diagnostic labor covered when repair is covered",
    "Replacement fluids covered in conjunction with a covered repair",
]

PRODUCTS = [
    # ════════════════════════════════════════════════════════════════════
    # MECHANICAL — ICE
    # ════════════════════════════════════════════════════════════════════
    {
        "name": "Vehicle Service Protection [ICE]",
        "subtitle": "Mechanical breakdown coverage — parts & labor — up to 10 years / 150,000 miles",
        "code": "HFVI",
        "core": True,
        "ev_care": False,
        "circle": True,
        "overview": "Four-tier stated/exclusionary plan protecting ICE Hyundai vehicles against unanticipated mechanical repairs. Coverage is additive, beginning the day the contract is sold.",
        "plans": [
            {"name": "Powertrain", "desc": "Engine, Transmission, Drive Axle incl. CV joints (Stated Component)"},
            {"name": "Gold", "desc": "Powertrain + Front/Rear Suspension (incl. shocks), A/C, Fuel System, Electrical System (Stated Component)"},
            {"name": "High Technology", "desc": "Factory-installed ADAS, Audio/Visual, Accessories (USB, 120V, wireless charging, WiFi/hotspot, phone interfaces) (Stated Component)"},
            {"name": "Platinum", "desc": "All Covered Parts except listed exclusions; includes Powertrain, Gold & Hi-Tech + Brakes (excl. wear), Power Steering, CV Boots (Exclusionary)"},
            {"name": "Platinum Plus (FL Only)", "desc": "Platinum + headlamps, belts & hoses, electrical — requires dealer enrollment & updated DAMA setup"},
        ],
        "features": [
            "Additive term coverage — starts on contract sale date",
            "Deductible options: $0 or $100",
            "Optional Light Duty Commercial Use Coverage (surcharge applies)",
            "High Technology coverage not available with Light Duty Commercial",
        ],
        "terms_mileage": "Up to 10 years / 150,000 miles (additive from contract sale date)",
        "deductible": "$0 or $100",
        "eligibility": [
            "New or pre-owned ICE Hyundai vehicles",
            "≤9 model years old; <140,001 miles at purchase",
            "If purchased after vehicle sale: ≥1 month & 1,000 miles manufacturer warranty remaining",
            "UCI required for pre-owned vehicles with 50,001+ miles or >4.5 years old",
            "Commercial/fleet excluded unless Permitted Commercial Purpose",
        ],
        "exclusions": [
            "Total loss, salvaged, junk, or Buy-Back vehicles",
            "Non-U.S. specification vehicles (e.g., Canada-import)",
            "Prohibited Commercial Use (livery, hauling, delivery, rentals, etc.)",
            "EVs and Hybrids — use EV Care VSP",
            "High Technology plan incompatible with Light Duty Commercial Use",
        ],
        "transfer_fee": "$75 (private party or lease assumption only; not transferable if dealer is party to resale)",
        "cancel_fee": "Within 30 days: full refund less claims. After 30 days: pro-rata less claims + $75 fee",
        "additional_benefits": ADDITIONAL_BENEFITS_MECH,
    },
    {
        "name": "Vehicle Service Protection [ICE] — Competitive Makes",
        "subtitle": "PowerProtect mechanical coverage for non-Hyundai pre-owned vehicles — up to 10 years / 120,000 miles",
        "code": "HPVI",
        "core": False,
        "ev_care": False,
        "circle": False,
        "overview": "White-label PowerProtect VSP for pre-owned ICE competitive makes sold at Hyundai dealerships. Three plan tiers. High Technology coverage not available.",
        "plans": [
            {"name": "Powertrain", "desc": "Engine, Transmission, Drive Axle incl. CV joints"},
            {"name": "Gold", "desc": "Powertrain + Front/Rear Suspension, A/C, Fuel System, Electrical System"},
            {"name": "Platinum", "desc": "All Covered Parts except listed exclusions; includes Powertrain & Gold + Brakes (excl. wear), Power Steering, CV Boots"},
            {"name": "Platinum Plus (FL Only)", "desc": "Platinum + headlamps, belts & hoses, electrical"},
        ],
        "features": [
            "Additive term coverage — starts on contract sale date",
            "Deductible options: $0, $100, or $250",
            "High Technology coverage NOT available for Competitive Makes",
        ],
        "terms_mileage": "Up to 10 years / 120,000 miles (additive from contract sale date)",
        "deductible": "$0, $100, or $250",
        "eligibility": [
            "Pre-owned ICE competitive make (non-Hyundai) vehicles",
            "≤9 model years; <120,000 miles at purchase",
            "UCI required for pre-owned",
            "Commercial/fleet excluded unless Permitted Commercial Purpose",
        ],
        "exclusions": [
            "Hyundai vehicles — use standard HPP VSP",
            "Total loss, salvaged, junk, or Buy-Back vehicles",
            "Non-U.S. specification vehicles",
            "Prohibited Commercial Use",
        ],
        "transfer_fee": "$75 (private party or lease assumption only)",
        "cancel_fee": "Within 30 days: full refund less claims. After 30 days: pro-rata less claims + $75 fee",
        "additional_benefits": ADDITIONAL_BENEFITS_MECH,
    },
    {
        "name": "Vehicle Service Protection Original Owner [ICE]",
        "subtitle": "Platinum exclusionary coverage measured from the original in-service date — up to 10 yrs / 150,000 mi",
        "code": "HFOI",
        "core": True,
        "ev_care": False,
        "circle": False,
        "overview": "Built for original owners and lessees buying out their vehicle. Coverage measures from the original in-service date. Platinum plan only; under 59,000 miles at purchase.",
        "plans": [
            {"name": "Platinum", "desc": "All Covered Parts except listed exclusions; includes Brakes (excl. wear), Shocks, Suspension, Fuel, Electrical, Power Steering, Climate Control, Nav, Audio, Hi-Tech (BlueLink, Bluetooth, HomeLink)"},
            {"name": "Platinum Plus (FL Only)", "desc": "Platinum + headlamps, belts & hoses, electrical"},
        ],
        "features": [
            "Coverage from original in-service date and zero miles",
            "Original owner or original lessee buy-out required",
            "Deductible options: $0 or $100",
        ],
        "terms_mileage": "Up to 10 years / 150,000 miles from original in-service date (zero-mile start)",
        "deductible": "$0 or $100",
        "eligibility": [
            "ICE Hyundai vehicles: >10,000 miles but <59,000 miles",
            "Must be currently owned by original owner or purchased by original lessee",
            "UCI required for pre-owned",
            "If post-sale: ≥1 month & 1,000 miles of manufacturer warranty remaining",
        ],
        "exclusions": [
            "Vehicles ≥59,000 miles at purchase — use standard VSP",
            "Non-original owners (except original lessee buyout)",
            "Total loss, salvaged, junk, Buy-Back vehicles",
            "Non-U.S. specification vehicles",
        ],
        "transfer_fee": "$75 (private party or lease assumption only)",
        "cancel_fee": "Within 30 days: full refund less claims. After 30 days: pro-rata less claims + $75 fee",
        "additional_benefits": ADDITIONAL_BENEFITS_MECH,
    },
    {
        "name": "Certified Used Vehicle (CUV) Wrap [ICE]",
        "subtitle": "Enhances HPP CUV Powertrain warranty — up to 10 years / 150,000 miles from original in-service date",
        "code": "HFCI",
        "core": True,
        "ev_care": False,
        "circle": True,
        "overview": "Wrap coverage layered on top of the Hyundai CUV powertrain limited warranty. Adds Gold or Platinum mechanical coverage for certified used ICE Hyundai vehicles.",
        "plans": [
            {"name": "Gold", "desc": "Front/Rear Suspension (incl. shocks), A/C, Fuel System, Electrical System (Stated Component)"},
            {"name": "Platinum", "desc": "All Covered Parts except listed exclusions; includes Gold + Brakes (excl. wear), Power Steering, Factory Nav, Audio, Hi-Tech (BlueLink, Bluetooth, HomeLink)"},
            {"name": "Platinum Plus (FL Only)", "desc": "Gold + Platinum + headlamps, belts & hoses, electrical"},
        ],
        "features": [
            "Coverage from original in-service date and zero miles",
            "Deductible options: $0, $50 (disappearing at selling dealer), or $100",
            "CUV Powertrain deductible will not exceed selected deductible on CUV Wrap",
            "Circle: Requires Certificate; max markup $100 over dealer cost",
        ],
        "terms_mileage": "Up to 10 years / 150,000 miles from original in-service date (zero-mile start)",
        "deductible": "$0, $50 (disappearing), or $100",
        "eligibility": [
            "Hyundai Certified Used (CUV) ICE vehicles only",
            "<6 years from original in-service date AND <80,000 miles",
            "Must have active CPO/CUV certification",
            "UCI required",
            "If post-sale: ≥1 month & 1,000 miles of manufacturer warranty remaining",
        ],
        "exclusions": [
            "Vehicle exceeds 80,000 miles",
            "Non-CUV/CPO certified vehicles",
            "Total loss, salvaged, junk, Buy-Back vehicles",
            "Platinum Plus incompatible with Wear Protection",
        ],
        "transfer_fee": "$75 (private party or lease assumption only)",
        "cancel_fee": "Within 30 days: full refund less claims. After 30 days: pro-rata less claims + $75 fee",
        "additional_benefits": ADDITIONAL_BENEFITS_MECH,
    },
    {
        "name": "High Mileage Vehicle Service Protection [ICE]",
        "subtitle": "PowerProtect Powertrain-only coverage for high-mileage used vehicles — up to 5 yrs / 60,000 mi",
        "code": "HPHI",
        "core": False,
        "ev_care": False,
        "circle": False,
        "overview": "Designed for used vehicles with 60,001–175,000 miles. Powertrain-only stated component plan for Hyundai and competitive makes. Available at time of vehicle sale only.",
        "plans": [
            {"name": "Powertrain", "desc": "Engine, Transmission, Drive Axle incl. CV joints (Stated Component)"},
        ],
        "features": [
            "Additive term — starts on contract sale date",
            "Deductible options: $250 or $500",
            "Available only at time of vehicle sale (not post-sale)",
        ],
        "terms_mileage": "Up to 5 years / 60,000 miles (additive from contract sale date)",
        "deductible": "$250 or $500",
        "eligibility": [
            "Used ICE Hyundai vehicles AND competitive makes sold at Hyundai dealerships",
            "≤9 model years; odometer 60,001–175,000 miles",
            "Available at time of vehicle sale ONLY",
        ],
        "exclusions": [
            "Vehicles <60,001 or >175,000 miles — use standard VSP or no coverage",
            "Post-sale purchase not permitted",
            "Rideshare and commercial fleets (except permitted professional use)",
        ],
        "transfer_fee": "$75 (private party or lease assumption only)",
        "cancel_fee": "Within 30 days: full refund less claims. After 30 days: pro-rata less claims + $75 fee",
        "additional_benefits": ADDITIONAL_BENEFITS_MECH,
    },
    {
        "name": "EV Care Vehicle Service Protection",
        "subtitle": "Hyundai EV & Hybrid mechanical coverage — up to 12 years / 200,000 miles",
        "code": "HFVE",
        "core": True,
        "ev_care": True,
        "circle": True,
        "overview": "Comprehensive mechanical repair protection for Hyundai EVs, HEVs, PHEVs, and Nexo vehicles. Three plan tiers including Battery, High Technology, and full exclusionary Platinum.",
        "plans": [
            {"name": "Battery", "desc": "Lithium-Ion Battery Pack, Battery Module, BMS, Battery Degradation (<70%), HV Pre-Charger, Traction Motor (incl. housing), Onboard Charger, Inverter, Converter"},
            {"name": "High Technology", "desc": "Factory-installed ADAS, Audio/Visual, Accessories (USB, 120V, wireless charging, WiFi/hotspot)"},
            {"name": "Platinum", "desc": "All Covered Parts except listed exclusions — includes Battery + Hi-Tech + Electronic Air Compressor, EPCU, Transmission, Powertrain, A/C Refrigerant, 12V Battery (first 3 yrs), Climate Control, Shocks/Suspension, Steering, EV Regenerative Brakes"},
            {"name": "Platinum Plus (FL Only)", "desc": "Platinum + headlamps, belts & hoses, electrical"},
        ],
        "features": [
            "Additive term — starts on contract sale date",
            "Deductible options: $0 or $100",
            "Circle: Post-sale purchase requires ≥1 month warranty + 1,000 miles; max markup $100 over dealer cost; not available in FL",
        ],
        "terms_mileage": "Up to 12 years / 200,000 miles (additive from contract sale date)",
        "deductible": "$0 or $100",
        "eligibility": [
            "New or pre-owned Hyundai EVs/HEVs/PHEVs/Nexo vehicles only (not ICE)",
            "≤9 model years; <140,001 miles at purchase",
            "In-warranty vehicles eligible at time of vehicle purchase",
            "Post-sale: ≥1 month & 1,000 miles manufacturer warranty remaining",
            "UCI required for pre-owned",
        ],
        "exclusions": [
            "ICE vehicles — use standard VSP",
            "Total loss, salvaged, junk, Buy-Back vehicles",
            "Non-U.S. specification vehicles",
            "Prohibited Commercial Use",
            "High Technology plan incompatible with Light Duty Commercial Use",
        ],
        "transfer_fee": "$75 (private party or lease assumption only)",
        "cancel_fee": "Within 30 days: full refund less claims. After 30 days: pro-rata less claims + $75 fee",
        "additional_benefits": ADDITIONAL_BENEFITS_EV,
    },
    {
        "name": "EV Care VSP — Competitive Makes",
        "subtitle": "PowerProtect EV coverage for pre-owned competitive make EVs — up to 12 years / 120,000 miles",
        "code": "HPVE",
        "core": False,
        "ev_care": True,
        "circle": False,
        "overview": "White-label EV VSP for pre-owned non-Hyundai EV competitive makes. Platinum exclusionary plan only. Eligibility cap: <108,001 miles.",
        "plans": [
            {"name": "Platinum", "desc": "All Covered Parts except listed exclusions — includes all Battery & Hi-Tech components, EPCU, Transmission, Powertrain, A/C Refrigerant, 12V Battery, Climate Control, Shocks/Suspension, Steering, EV Regenerative Brakes"},
            {"name": "Platinum Plus (FL Only)", "desc": "Platinum + headlamps, belts & hoses, electrical"},
        ],
        "features": [
            "Additive term — starts on contract sale date",
            "Deductible options: $0, $100, or $250",
        ],
        "terms_mileage": "Up to 12 years / 120,000 miles (additive from contract sale date)",
        "deductible": "$0, $100, or $250",
        "eligibility": [
            "Pre-owned competitive make EVs only",
            "≤9 model years; <108,001 miles at purchase",
            "UCI required",
        ],
        "exclusions": [
            "Hyundai EVs — use EV Care VSP (HFVE)",
            "Total loss, salvaged, junk, Buy-Back vehicles",
            "Non-U.S. specification vehicles",
        ],
        "transfer_fee": "$75 (private party or lease assumption only)",
        "cancel_fee": "Within 30 days: full refund less claims. After 30 days: pro-rata less claims + $75 fee",
        "additional_benefits": ADDITIONAL_BENEFITS_EV,
    },
    {
        "name": "EV Care VSP — Original Owner",
        "subtitle": "EV Platinum exclusionary coverage from original in-service date — up to 12 yrs / 200,000 mi",
        "code": "HFOE",
        "core": True,
        "ev_care": True,
        "circle": False,
        "overview": "Built for EV/HEV lessees buying out their vehicle or original Hyundai EV owners with <59,000 miles. Coverage measured from original in-service date.",
        "plans": [
            {"name": "Battery", "desc": "Lithium-Ion Battery Pack, BMS, Battery Degradation (<70%), HV Pre-Charger, Traction Motor, Onboard Charger, Inverter, Converter"},
            {"name": "Platinum", "desc": "All Covered Parts except listed exclusions — Battery + EPCU, Transmission, Powertrain, A/C Refrigerant, Radio/Sound, 12V Battery, Climate Control, Shocks, Steering, EV Regen Brakes, Blind Spot Monitor, Nav/Audio/Tech"},
            {"name": "Platinum Plus (FL Only)", "desc": "Platinum + headlamps, belts & hoses, electrical"},
        ],
        "features": [
            "Coverage from original in-service date and zero miles",
            "Original EV owner or original lessee buy-out required",
            "Deductible options: $0 or $100",
        ],
        "terms_mileage": "Up to 12 years / 200,000 miles from original in-service date",
        "deductible": "$0 or $100",
        "eligibility": [
            "Hyundai EVs/HEVs/PHEVs/Nexo: >10,000 miles but <59,000 miles",
            "Must be currently owned by original owner or purchased by original lessee",
            "UCI required for pre-owned",
        ],
        "exclusions": [
            "Vehicles ≥59,000 miles at purchase",
            "Non-original owners (except original lessee buyout)",
            "Total loss, salvaged, junk, Buy-Back vehicles",
        ],
        "transfer_fee": "$75 (private party or lease assumption only)",
        "cancel_fee": "Within 30 days: full refund less claims. After 30 days: pro-rata less claims + $75 fee",
        "additional_benefits": ADDITIONAL_BENEFITS_EV,
    },
    {
        "name": "EV Care Certified Used Vehicle (CUV) Wrap",
        "subtitle": "Enhances Hyundai EV CUV powertrain warranty — up to 12 years / 200,000 miles from in-service date",
        "code": "HFCE",
        "core": True,
        "ev_care": True,
        "circle": True,
        "overview": "Wrap coverage for Hyundai EV/Hybrid certified used vehicles. Three plan tiers: Battery, Gold, Platinum — layered on top of CUV HEV/PHEV/EV Battery warranty.",
        "plans": [
            {"name": "Battery", "desc": "Lithium-Ion Battery Pack, BMS, Battery Degradation (<70%), HV Pre-Charger, Traction Motor, Onboard Charger, Inverter, Converter"},
            {"name": "Gold", "desc": "Front/Rear Suspension (incl. shocks), A/C, Fuel System, Electrical System"},
            {"name": "Platinum", "desc": "All Covered Parts except listed exclusions; includes Gold + Brakes (excl. wear), Power Steering, Factory Nav, Audio, Hi-Tech (BlueLink, Bluetooth, HomeLink)"},
            {"name": "Platinum Plus (FL Only)", "desc": "Battery + Gold + Platinum + headlamps, belts & hoses, electrical"},
        ],
        "features": [
            "Coverage from original in-service date and zero miles",
            "Deductible options: $0, $50 (disappearing at selling dealer), or $100",
            "Circle: Requires Certificate; max markup $100 over dealer cost; active CPO/CUV cert required",
        ],
        "terms_mileage": "Up to 12 years / 200,000 miles from original in-service date (zero-mile start)",
        "deductible": "$0, $50 (disappearing), or $100",
        "eligibility": [
            "Hyundai EV/Hybrid Certified Used Vehicles (CUV) only",
            "<6 years from original in-service date AND <80,000 miles",
            "Must have active CPO/CUV certification",
            "UCI required",
        ],
        "exclusions": [
            "Vehicle exceeds 80,000 miles",
            "Non-CUV/CPO certified vehicles",
            "Total loss, salvaged, junk, Buy-Back vehicles",
            "Platinum Plus incompatible with Wear Protection",
        ],
        "transfer_fee": "$75 (private party or lease assumption only)",
        "cancel_fee": "Within 30 days: full refund less claims. After 30 days: pro-rata less claims + $75 fee",
        "additional_benefits": ADDITIONAL_BENEFITS_EV,
    },
    {
        "name": "Vehicle Service Protection — Livery Coverage [ICE & EV]",
        "subtitle": "Mechanical coverage for legal Livery use vehicles — 5 years / 100,000 mi to 5 yrs / 300,000 mi",
        "code": "HCVL",
        "core": True,
        "ev_care": False,
        "circle": False,
        "overview": "Specialized VSP for limousines, chauffeured vehicles, private hire, and shuttles. Available only through select dealers — contact your District Manager of Insurance.",
        "plans": [
            {"name": "ICE — Powertrain", "desc": "Engine, Transmission, Drive Axle incl. CV joints"},
            {"name": "ICE — Gold", "desc": "Powertrain + Front/Rear Suspension, A/C, Fuel System, Electrical System"},
            {"name": "ICE — Platinum", "desc": "All Covered Parts except listed exclusions (incl. Powertrain & Gold)"},
            {"name": "EV — Battery", "desc": "Lithium-Ion Battery Pack, BMS, Battery Degradation (<70%), HV Pre-Charger, Traction Motor, Onboard Charger, Inverter, Converter"},
            {"name": "EV — Platinum", "desc": "All Covered Parts except listed exclusions (incl. Battery components)"},
        ],
        "features": [
            "No deductible on any plan",
            "Uses genuine Hyundai Parts for all repairs",
            "Only available to certain dealers — contact District Manager of Insurance",
            "New vehicles only; ≤10,000 miles at purchase",
        ],
        "terms_mileage": "5 years / 100,000 miles  →  5 years / 300,000 miles",
        "deductible": "$0",
        "eligibility": [
            "New ICE and EV Hyundai vehicles only",
            "<10,000 miles at time of purchase",
            "Vehicle must be licensed for legal livery use (limousine, chauffeured, private hire, shuttle)",
            "Available at time of vehicle sale ONLY",
            "Available through select dealers only",
        ],
        "exclusions": [
            "Pre-owned vehicles",
            "Vehicles >10,000 miles at purchase",
            "Non-livery commercial use",
            "Total loss, salvaged, junk, Buy-Back vehicles",
        ],
        "transfer_fee": "$75 (private party or lease assumption only)",
        "cancel_fee": "Within 30 days: full refund less claims. After 30 days: pro-rata less claims + $75 fee",
        "additional_benefits": [
            "24-Hour Emergency Roadside Assistance",
            "Rental Car / Ride Share Reimbursement: up to $55/day, 10 days max",
            "Trip Interruption Benefit: up to $300/day, 5 days max (not available in NY)",
        ],
    },
    {
        "name": "Wear Protection [ICE & EV]",
        "subtitle": "Coverage for common wear items on new vehicles — up to 60 months / 60,000 miles",
        "code": "HFWP / HFEW (Circle)",
        "core": True,
        "ev_care": False,
        "circle": True,
        "overview": "Covers out-of-pocket costs for typical wear items including brakes, battery, alignment, headlamps, belts, hoses, wipers, fuses, and bulbs. New vehicles only; at time of sale.",
        "plans": [
            {"name": "Standard", "desc": "Brake pads/shoes (1 set front & rear), 12V Battery (1), Wheel Alignment (1), Headlamps (unlimited, non-impact), Engine Belts & Hoses (unlimited), Wiper Blades (1 set), Fuses & Light Bulbs (unlimited, non-impact)"},
            {"name": "+ Brake Rotor (Optional)", "desc": "Adds one replacement set of front & rear brake rotors during the term"},
        ],
        "features": [
            "Additive term — starts on contract sale date",
            "$0 deductible",
            "Cannot be sold on vehicles with Hyundai Maintenance Wrap Plus (HCM)",
            "Available at time of vehicle sale ONLY",
            "Circle: Requires Certificate; max markup $100 over dealer cost (Code: HFEW)",
        ],
        "terms_mileage": "Up to 60 months / 60,000 miles (additive from contract sale date)",
        "deductible": "$0",
        "eligibility": [
            "New Hyundai ICE, EV, Hybrid, and Hydrogen Fuel Cell vehicles",
            "Pre-owned ICE and EV competitive makes (≤10,000 miles at purchase)",
            "Vehicle mileage cannot exceed 10,000 at time of purchase",
            "Available at time of vehicle sale ONLY",
        ],
        "exclusions": [
            "Vehicles with Hyundai Maintenance Wrap Plus (not compatible)",
            "Hybrid & EV batteries excluded from 12V battery benefit",
            "Impact-damaged headlamps, bulbs not covered",
            "Competitive makes on new vehicles (only pre-owned competitive makes eligible)",
            "Platinum Plus incompatible with Wear Protection",
        ],
        "transfer_fee": "$75 (private party or lease assumption only)",
        "cancel_fee": "Within 30 days: full refund less claims. After 30 days: pro-rata less claims + $75 fee",
        "additional_benefits": ADDITIONAL_BENEFITS_MECH,
    },

    # ════════════════════════════════════════════════════════════════════
    # MAINTENANCE PLANS
    # ════════════════════════════════════════════════════════════════════
    {
        "name": "Maintenance Wrap [ICE]",
        "subtitle": "Extends / wraps Hyundai Complimentary Maintenance (HCM) — up to 8 years / 96,000 miles",
        "code": "HFMW",
        "core": True,
        "ev_care": False,
        "circle": False,
        "overview": "Enhances and extends HCM coverage on MY2025 and older ICE vehicles. Two plans: Basic (mirrors HCM) and Scheduled (HCM + additional owner's manual services). Not available on EVs.",
        "plans": [
            {"name": "Basic Maintenance Wrap", "desc": "Mirrors HCM — oil & oil filter, tire rotation, multi-point inspection. Extension beyond 3 yrs/36K mi. Option to upgrade to Severe Usage intervals."},
            {"name": "Scheduled Maintenance Wrap", "desc": "Basic + additional owner's manual items: air filters, brake fluid, spark plugs, fluids, etc. Extension and Severe Usage option available."},
        ],
        "features": [
            "Only for MY2025 and older vehicles that include HCM",
            "Can be sold after vehicle sale but only before 6,000 miles on odometer",
            "Covered services must be performed within 5,000 miles or 6 months of scheduled interval",
            "Synthetic oil only (effective Oct 15, 2023)",
            "Non-Transferable",
        ],
        "terms_mileage": "Up to 8 years / 96,000 miles (coverage beyond HCM term of 3 yrs/36K mi)",
        "deductible": "$0",
        "eligibility": [
            "MY2025 and older ICE Hyundai vehicles WITH HCM only",
            "Available at time of sale or before 6,000 miles",
            "Pre-Paid Maintenance NOT eligible on any vehicle with HCM",
        ],
        "exclusions": [
            "EV vehicles",
            "Vehicles without HCM — use Pre-Paid Maintenance",
            "MY2026+ vehicles — use Pre-Paid Maintenance (PPM)",
            "Mechanical breakdown repairs not covered",
            "Services not performed within 5,000 miles or 6 months of interval",
        ],
        "transfer_fee": "Non-Transferable",
        "cancel_fee": "Within 30 days: full refund if no covered services. Next 35 months: full refund less $50 fee if no covered services. After: pro-rata less services rendered + $50 fee",
        "additional_benefits": [],
    },
    {
        "name": "Pre-Paid Maintenance [ICE]",
        "subtitle": "Prepaid factory-scheduled maintenance for ICE & Hybrid Hyundai — up to 8 years / 96,000 miles",
        "code": "HFPI (Basic) / HFSI (Scheduled)",
        "core": True,
        "ev_care": False,
        "circle": True,
        "overview": "Locks in today's prices for oil changes, tire rotations, inspections, and scheduled services. Basic available anytime; Scheduled only on new MY2026+ vehicles at time of sale.",
        "plans": [
            {"name": "Basic", "desc": "Oil & oil filter, tire rotation, multi-point inspection at selected intervals. Normal or Severe usage. Available on all Hyundai ICE & Hybrid vehicles without HCM."},
            {"name": "Scheduled", "desc": "Basic + climate control air filter, engine air filter, brake fluid, spark plugs, transmission oil, differential/case oil per interval sheet. New MY2026+ only; at time of sale; <6,000 miles."},
        ],
        "features": [
            "Synthetic Oil Only",
            "Normal or Severe usage options",
            "Basic: available at any time; Scheduled: new vehicles at time of sale only",
            "Reimbursed per maintenance interval sheets",
            "Circle (Friends & Family): Requires Certificate; max markup $100 over dealer cost (Code: HFEP)",
            "Helps ensure compliance with factory warranty requirements",
        ],
        "terms_mileage": "Up to 8 years / 96,000 miles",
        "deductible": "$0",
        "eligibility": [
            "ICE & Hybrid Hyundai vehicles (not EV)",
            "Basic: any vehicle without HCM, any time",
            "Scheduled: new MY2026+ vehicles only; ≤6,000 miles; at time of sale",
            "NOT available on any vehicle with HCM",
        ],
        "exclusions": [
            "Vehicles with HCM (use Maintenance Wrap instead)",
            "EV vehicles — use EV Care Maintenance",
            "Mechanical breakdown repairs not covered",
            "Services not completed within interval requirements",
        ],
        "transfer_fee": "$50 (private party or lease assumption only, within 30 days of transfer)",
        "cancel_fee": "Within 30 days: full refund less claims. After 30 days: pro-rata less services + $75 fee",
        "additional_benefits": [],
    },
    {
        "name": "Pre-Paid Maintenance — Competitive Makes [ICE]",
        "subtitle": "Prepaid oil changes & tire rotations for non-Hyundai ICE vehicles — up to 7 years / 105,000 miles",
        "code": "HPCM",
        "core": False,
        "ev_care": False,
        "circle": False,
        "overview": "White-label PPM for competitive make ICE vehicles sold at Hyundai dealerships. Three plans: Basic (oil only), Plus (oil + rotation), Synthetic Plus (synthetic oil + rotation). Sold at any time.",
        "plans": [
            {"name": "Basic", "desc": "Conventional oil & oil filter changes at selected intervals"},
            {"name": "Plus", "desc": "Conventional oil changes + tire rotations"},
            {"name": "Synthetic Plus", "desc": "Synthetic oil changes + tire rotations"},
        ],
        "features": [
            "Multiple term/interval options",
            "Can be sold at any time (no time-of-sale restriction)",
            "Locks in today's service prices",
        ],
        "terms_mileage": "Up to 7 years / 105,000 miles",
        "deductible": "$0",
        "eligibility": [
            "ICE competitive make (non-Hyundai) vehicles sold at Hyundai dealerships",
            "Can be sold at any time",
            "Commercial/fleet excluded unless permitted professional use",
        ],
        "exclusions": [
            "Hyundai vehicles — use HPP Pre-Paid Maintenance",
            "EV / Hybrid competitive makes",
        ],
        "transfer_fee": "$50 (private party or lease assumption only)",
        "cancel_fee": "Within 30 days: full refund less claims. After 30 days: pro-rata less services + $75 fee",
        "additional_benefits": [],
    },
    {
        "name": "EV Care Maintenance",
        "subtitle": "Prepaid EV-specific maintenance — up to 8 years / 96,000 miles",
        "code": "HFEM (Circle)",
        "core": True,
        "ev_care": True,
        "circle": True,
        "overview": "Designed for Hyundai EV and Fuel Cell vehicles. Two plans: Basic (tire rotation + inspection + cabin filter) and Maintenance Plus (all scheduled EV services). Basic anytime; Plus at time of sale only.",
        "plans": [
            {"name": "Basic (New + Used)", "desc": "Tire rotation + multi-point maintenance inspection every interval; Cabin air filter replacement every other interval"},
            {"name": "Maintenance Plus (New MY2026+)", "desc": "All Basic items + all owner's manual scheduled services + wiper blades (every other interval), wheel alignment at 32K/72K mi, 12V battery at 72K mi, brake pads at 72K mi, brake fluid at 48K/96K mi"},
        ],
        "features": [
            "Basic: available on all Hyundai EV and Fuel Cell vehicles; sold at any time",
            "Maintenance Plus: new MY2026+ Hyundai EV only; ≤3 years old; ≤6,000 miles; at time of sale",
            "Reimbursed per maintenance interval sheets",
            "Circle (Friends & Family): Requires Certificate; max markup $100 over dealer cost (Code: HFEM)",
        ],
        "terms_mileage": "Up to 8 years / 96,000 miles",
        "deductible": "$0",
        "eligibility": [
            "Hyundai EV and Fuel Cell vehicles (not ICE, not PHEV/Plug-in Hybrid)",
            "Basic: all model years; any mileage; sold at any time",
            "Maintenance Plus: new MY2026+ EV; ≤3 years old; ≤6,000 miles; at time of sale",
            "Maintenance Plus: vehicle must not have HCM",
        ],
        "exclusions": [
            "Plug-in Hybrid vehicles not eligible",
            "ICE vehicles — use Pre-Paid Maintenance",
            "Maintenance Plus not for vehicles with HCM",
        ],
        "transfer_fee": "$50 (private party or lease assumption only, within 30 days of transfer)",
        "cancel_fee": "Within 30 days: full refund less claims. After 30 days: pro-rata less services + $75 fee",
        "additional_benefits": [],
    },

    # ════════════════════════════════════════════════════════════════════
    # DAMAGE CARE
    # ════════════════════════════════════════════════════════════════════
    {
        "name": "Lease-End Protection",
        "subtitle": "Waives billable excess wear & use charges at lease turn-in — up to 72 months",
        "code": "HFLP",
        "core": True,
        "ev_care": False,
        "circle": True,
        "overview": "Available only on leases assigned to Hyundai Motor Finance / Hyundai Lease Titling Trust. Waives covered excess wear and use charges with an unlimited benefit and $1,000 single-event limit.",
        "plans": [
            {"name": "Single Plan", "desc": "Covers billable excess wear & use charges waived by Hyundai Lease Titling Trust/HMF. Vehicle can be turned in any time prior to or within 1 year after the original scheduled termination date."},
        ],
        "features": [
            "Unlimited waiver benefit; single event limit $1,000",
            "$0 deductible",
            "Up to 200 excess miles covered at $0.20/mile",
            "Missing parts: up to $150 total ($100 with Circle Certificate)",
            "Lease term must match protection term",
            "Circle (Friends & Family): Requires Certificate; max markup $100 over dealer cost (Code: HFEL)",
            "Not available in NY or MD",
        ],
        "components": [
            "Paint damage & exterior dents", "Front/rear door & hatch rubber seals",
            "Tail, turn signal & brake light bulbs", "Muffler(s) and tail pipe",
            "Front & rear windshield, side glass chips/breakage", "Convertible & vinyl tops",
            "Carpet rips, tears & stains", "Door, trunk & hatch handles",
            "Alloy & steel wheels & wheel covers", "Factory-installed audio",
            "Scratches, chips, pits or cracks", "Headlamps, running lights, fog lights",
            "Bumper covers & surrounding trim", "Interior upholstery rips/tears/stains",
            "Chrome, bright metal, moldings & trim", "Side & rear-view mirrors",
            "Hood ornaments & body insignia", "Worn tires",
        ],
        "terms_mileage": "12 to 72 months; up to 150,000 miles (must match lease term)",
        "deductible": "$0",
        "eligibility": [
            "Hyundai leases with HMF / Hyundai Lease Titling Trust only",
            "New Hyundai vehicles (ICE, EV, Hybrid, Hydrogen)",
            "Purchased at time of lease only",
        ],
        "exclusions": [
            "Not available in NY or MD",
            "Alteration charges or pre-turn-in repairs not covered",
            "No benefit if vehicle is purchased by customer or dealer",
            "Single event charges >$1,000 excluded",
            "Excess mileage >200 miles not covered",
        ],
        "transfer_fee": "$75 (to subsequent lessee only)",
        "cancel_fee": "Within 30 days: full refund. After 30 days: pro-rata + $75 fee. Once waiver paid: fully earned, no cancellation",
        "additional_benefits": [],
    },
    {
        "name": "Multi-Coverage Protection",
        "subtitle": "Bundled cosmetic & ancillary protection — Tire & Wheel, Dent, Windshield, Key — up to 7 years",
        "code": "HFMC",
        "core": False,
        "ev_care": False,
        "circle": False,
        "overview": "Four bundle packages on a single contract. Products can be sold as a bundle or as stand-alone on the Multi-Coverage form. Available on Hyundai and competitive makes including Kia & Genesis.",
        "plans": [
            {"name": "Platinum Bundle", "desc": "Tire & Wheel + Dent + Key + Windshield (Complete Package)"},
            {"name": "Gold Bundle", "desc": "Tire & Wheel + Dent + Windshield (No Key)"},
            {"name": "Silver Bundle", "desc": "Tire & Wheel + Dent + Key (FL Bundle Only)"},
            {"name": "Bronze Bundle", "desc": "Tire & Wheel + Dent (Entry Package)"},
        ],
        "features": [
            "Products can be sold as bundle OR as stand-alone on the MC form",
            "Windshield not available stand-alone in FL; not in AZ, GA, ME, NY, TX for replacement",
            "No deductible on any product",
            "No limit on occurrences (except hail damage & cosmetic wheel repairs)",
            "No mileage limitations",
            "Available at time of vehicle purchase only",
        ],
        "terms_mileage": "Up to 7 years (84 months) / Unlimited miles",
        "deductible": "$0",
        "eligibility": [
            "New, used, and CUV Hyundai vehicles",
            "Used competitive makes including Kia & Genesis",
            "No mileage or year limitations",
            "Available at time of vehicle purchase only",
        ],
        "exclusions": [
            "Individual product exclusions apply (see underlying product agreements)",
            "Windshield stand-alone not available in FL",
            "Cancellation of one product cancels ALL products on the contract",
        ],
        "transfer_fee": "$50 (private party or lease assumption only, within 30 days)",
        "cancel_fee": "Within 30 days: full refund less claims. After 30 days: pro-rata less claims + $50 fee. Cancellation serves to cancel ALL bundled products.",
        "additional_benefits": [],
    },
    {
        "name": "Tire & Wheel Protection",
        "subtitle": "Road hazard coverage for tire repair/replacement and wheel replacement — up to 7 years / unlimited miles",
        "code": "HFTW / HFTC (with Cosmetic)",
        "core": False,
        "ev_care": False,
        "circle": False,
        "overview": "Covers tire repair or replacement and structural wheel replacement from road hazard damage. Optional Cosmetic Wheel Coverage (up to 8 repairs). Run-flat and aftermarket tires covered.",
        "plans": [
            {"name": "Standard T&W", "desc": "Road hazard tire repair/replacement + structural wheel replacement + TPMS replacement (from road hazard)"},
            {"name": "+ Cosmetic Wheel (Optional)", "desc": "Adds cosmetic wheel repair for alloy/aluminum wheels via sanding/painting/refinishing — up to 8 individual wheel repairs during term"},
        ],
        "features": [
            "Covers glass, metal, potholes, debris, nails, and other road hazards",
            "Run-flat tires covered; aftermarket tires & wheels covered if meet manufacturer specs",
            "Up to $100 towing reimbursement per road hazard event",
            "10% markup on tires (not to exceed MSRP); 20% markup on wheels (not to exceed MSRP)",
            "Labor: up to $50 for repair; up to $30 for replacement",
            "Cosmetic wheel repair: up to $150 labor per wheel (covered)",
            "No prior authorization for tire repair at Hyundai dealer (reimbursement within 30 days)",
        ],
        "terms_mileage": "Up to 7 years (84 months) / Unlimited miles",
        "deductible": "$0",
        "eligibility": [
            "All Hyundai vehicles and competitive makes including Kia & Genesis sold at Hyundai dealerships",
            "Tire must have ≥3/32 inch tread depth at time of damage",
            "Available at time of vehicle purchase or lease only",
        ],
        "exclusions": [
            "Tires with <3/32 inch tread depth at time of damage",
            "Cosmetic wheel damage requiring material addition, bending, or straightening",
            "Commercial/fleet vehicles (permitted professional use allowed)",
            "Replacement wheels are not remanufactured",
        ],
        "transfer_fee": "$50 (within 30 days, private party or lease assumption only)",
        "cancel_fee": "Full refund within 30 days (60 days in UT) if no claim. After: pro-rata less claims + $50 fee",
        "additional_benefits": [],
    },
    {
        "name": "Dent Protection",
        "subtitle": "Paintless dent repair (PDR) for minor dents & dings — up to 7 years / unlimited miles",
        "code": "HFDP",
        "core": False,
        "ev_care": False,
        "circle": False,
        "overview": "Covers paintless dent repair process for minor dents on horizontal and vertical steel or aluminum body panels. Repairs may be completed at the dealership or customer's preferred location.",
        "plans": [
            {"name": "PDR Coverage", "desc": "Covers dents/dings ≤4 inches in diameter on accessible horizontal or vertical body panels. No limit on occurrences."},
            {"name": "Hail Damage (Optional)", "desc": "Up to $1,000 per occurrence (not available in all states) — lesser of comprehensive deductible or $1,000"},
        ],
        "features": [
            "No harm to factory finish",
            "Permanent repair of minor dents and dings",
            "Repairs at dealership or customer's preferred location",
            "Up to $65 labor for a single dent; up to $100 per panel",
            "No limit on occurrences",
        ],
        "terms_mileage": "Up to 7 years (84 months) / Unlimited miles",
        "deductible": "$0",
        "eligibility": [
            "New, used, and CUV Hyundai vehicles (no mileage or year limitation)",
            "Used competitive makes including Kia & Genesis when on branded consumer agreement",
            "Available at time of vehicle purchase/lease only",
        ],
        "exclusions": [
            "Dents/dings >4 inches in diameter",
            "Damage requiring repainting",
            "Creases exceeding PDR size limits",
            "Rust, corrosion, or pre-existing damage",
            "Intentional damage",
        ],
        "transfer_fee": "$50 (within 30 days, private party or lease assumption only)",
        "cancel_fee": "Full refund within 30 days (60 days in UT) if no claim. After: pro-rata less claims + $50 fee",
        "additional_benefits": [],
    },
    {
        "name": "Windshield Protection",
        "subtitle": "Chip repair & one windshield replacement — up to 7 years / unlimited miles",
        "code": "HFWS",
        "core": False,
        "ev_care": False,
        "circle": False,
        "overview": "Covers chip/crack repair and one front windshield replacement (including installation and calibration) caused by propelled rocks or road debris. Not available in FL. Repair-only in AZ, GA, ME, NY, TX.",
        "plans": [
            {"name": "Standard", "desc": "Unlimited chip/crack repairs + one (1) front windshield replacement (incl. installation & calibration). Repair-only states: AZ, GA, ME, NY, TX."},
        ],
        "features": [
            "No limit on repair occurrences",
            "One replacement includes installation and calibration",
            "Convenient repairs at dealership or customer's preferred location",
            "Up to $50 labor for single chip; up to $100 for multiple chips",
            "NOT available in Florida",
        ],
        "terms_mileage": "Up to 7 years (84 months) / Unlimited miles",
        "deductible": "$0",
        "eligibility": [
            "New, used, and CUV Hyundai vehicles (no mileage or year limitation)",
            "Used competitive makes including Kia & Genesis when on branded consumer agreement",
            "Available at time of vehicle purchase/lease only",
        ],
        "exclusions": [
            "NOT available in Florida (stand-alone or bundle)",
            "Repair-only (NO replacement) in AZ, GA, ME, NY, TX",
            "Stress cracks not covered",
            "Cracks over 6 inches not covered",
            "Rear and side glass not covered (front windshield only)",
        ],
        "transfer_fee": "$50 (within 30 days, private party or lease assumption only)",
        "cancel_fee": "Full refund within 30 days (60 days in UT) if no claim. After: pro-rata less claims + $50 fee",
        "additional_benefits": [],
    },
    {
        "name": "Appearance Protection",
        "subtitle": "Exterior paint & interior fabric/leather/vinyl protection — up to 120 months / unlimited miles",
        "code": "HFAP",
        "core": False,
        "ev_care": False,
        "circle": False,
        "overview": "Permanently applied interior and exterior chemical protection system. Covers paint from UV, corrosion, bird droppings, tree sap; interior fabric/leather/vinyl from stains and damage. Non-cancellable.",
        "plans": [
            {"name": "Exterior Paint Protection", "desc": "Protects factory painted surfaces from oxidation/UV fading, corrosion (sea/road salt, acid rain), bird droppings, tree sap, spray paint"},
            {"name": "Interior Fabric Protection", "desc": "Protects fabric seating/interior surfaces from stains (food/drink, pet waste, makeup, ink, dye transfer)"},
            {"name": "Interior Vinyl & Leather Protection", "desc": "Protects vinyl/leather surfaces from stains, dye transfer — benefit: cleaning, repair, dye, or replacement"},
        ],
        "features": [
            "Permanently applied at time of sale",
            "Rental Car Benefit: up to $50/day, 10 days max if vehicle in shop for covered repair",
            "Available on new and pre-owned Hyundai (incl. EV, Hybrid, Hydrogen) and competitive makes",
            "NON-CANCELLABLE",
        ],
        "terms_mileage": "120 months (10 years) / Unlimited miles",
        "deductible": "$0",
        "eligibility": [
            "New and pre-owned Hyundai vehicles (ICE, EV, Hybrid, Hydrogen)",
            "Competitive makes sold at HPP dealership",
            "Pre-owned vehicles: ≤5 model years old",
            "Available at time of sale only",
        ],
        "exclusions": [
            "Pre-owned vehicles older than 5 model years",
            "Commercial/fleet vehicles (permitted professional use allowed)",
        ],
        "transfer_fee": "$50 (private party; within 30 days with required documents)",
        "cancel_fee": "NON-CANCELLABLE",
        "additional_benefits": [
            "Rental Car Benefit: up to $50/day, 10 days max while vehicle is in for a covered repair",
        ],
    },

    # ════════════════════════════════════════════════════════════════════
    # LOSS & THEFT PROTECTION
    # ════════════════════════════════════════════════════════════════════
    {
        "name": "Guaranteed Asset Protection (GAP)",
        "subtitle": "Waives the gap between insurance settlement and remaining loan balance on a total loss",
        "code": "HFGP (Hyundai) / WFGP (Off-Make)",
        "core": True,
        "ev_care": False,
        "circle": True,
        "overview": "In a total loss, HPP GAP waives the difference between the vehicle's cash value (per primary insurer) and the remaining balance on the customer's finance agreement. Finance purchases only; not available on leases.",
        "plans": [
            {"name": "GAP", "desc": "Waives the difference between vehicle cash value (primary insurer) and remaining finance balance. Also waives up to $1,000 of customer's primary insurance deductible (not in all states)."},
        ],
        "features": [
            "No MSRP limitations",
            "No LTV limit",
            "No mileage limitations",
            "$0 deductible",
            "Circle: Requires Certificate; customer cost fixed at $100 over dealer cost (Code: HFEG)",
            "Off-Make (non-Hyundai) products: NON-CORE (Code: WFGP)",
            "Not available in CA, NY, TX",
        ],
        "terms_mileage": "Up to 96 months",
        "deductible": "$0",
        "eligibility": [
            "Finance purchases only (not available on leases)",
            "New, used, CUV Hyundai vehicles and used competitive makes",
            "Available at time of vehicle sale only",
        ],
        "exclusions": [
            "Leases not eligible",
            "Amounts owed due to missed payments or extraordinary events not waived",
            "NOT available in CA, NY, TX",
            "Deductible waiver not available in all states",
        ],
        "transfer_fee": "$75 (only upon assumption of retail installment contract)",
        "cancel_fee": "Within 30 days: full refund. After 30 days: pro-rata + $75 fee. Once waiver paid: fully earned, no cancellation",
        "additional_benefits": [],
    },
    {
        "name": "Guaranteed Asset Protection Plus (GAP Plus)",
        "subtitle": "GAP + $2,000 credit toward a replacement vehicle purchase at the selling dealer",
        "code": "HFGS (Hyundai) / WFGS (Off-Make)",
        "core": True,
        "ev_care": False,
        "circle": True,
        "overview": "All GAP benefits plus a $2,000 replacement vehicle credit applicable only at the original selling Hyundai dealer (or nearest if customer moves >50 miles). Not available in AK, KS, NE, NY, TN, TX, WV, CA.",
        "plans": [
            {"name": "GAP Plus", "desc": "All GAP benefits + $2,000 credit toward purchase/lease of replacement vehicle at original selling dealer (or nearest Hyundai dealer if >50 miles away)."},
        ],
        "features": [
            "Includes all standard GAP benefits (gap waiver + up to $1,000 deductible waiver)",
            "$2,000 replacement vehicle credit — must be used at original selling dealer or nearest Hyundai dealer",
            "No MSRP or LTV limitations",
            "Circle: Requires Certificate; customer cost fixed at $100 over dealer cost",
            "Off-Make NON-CORE (Code: WFGS)",
            "NOT available in AK, KS, NE, NY, TN, TX, WV, CA",
        ],
        "terms_mileage": "Up to 96 months",
        "deductible": "$0",
        "eligibility": [
            "Finance purchases only (not available on leases)",
            "New, used, CUV Hyundai vehicles and used competitive makes",
            "Available at time of vehicle sale only",
        ],
        "exclusions": [
            "NOT available in AK, KS, NE, NY, TN, TX, WV, CA",
            "Leases not eligible",
            "GAP Plus credit only at Hyundai dealership",
        ],
        "transfer_fee": "$75 (only upon assumption of retail installment contract)",
        "cancel_fee": "Within 30 days: full refund. After 30 days: pro-rata + $75 fee. Once waiver paid: fully earned, no cancellation",
        "additional_benefits": [],
    },
    {
        "name": "Key Protection",
        "subtitle": "Replacement coverage for lost, stolen, or damaged keys & fobs — up to 7 years",
        "code": "HFKP",
        "core": False,
        "ev_care": False,
        "circle": False,
        "overview": "Comprehensive key and fob replacement including programming, lockout assistance, home lockout, towing reimbursement, and alternate transportation. Motor Club license required in MA.",
        "plans": [
            {"name": "Key Protection", "desc": "Vehicle key/transponder replacement + programming + 24-hr lockout assistance + home lockout + towing reimbursement + alternate transportation + additional non-programmable keys"},
        ],
        "features": [
            "Vehicle Key Replacement: up to $800 per occurrence",
            "24-hr vehicle lockout assistance: locksmith up to $100",
            "24-hr home lockout assistance: up to $100",
            "Additional keys on keyring (if vehicle key lost/damaged): up to $250",
            "Alternate transportation: up to $75 (taxi, rental, public transit)",
            "Towing for key replacement/programming: up to $100",
            "Emergency Message Relay (not in all states)",
            "Motor Club license required in Massachusetts",
        ],
        "terms_mileage": "Up to 7 years / Unlimited miles",
        "deductible": "$0",
        "eligibility": [
            "New, used, CUV Hyundai and competitive makes including Kia & Genesis sold at Hyundai dealerships",
            "Pre-owned vehicles must receive at least 2 keys at time of sale",
            "Prior authorization required for key replacement benefit",
            "Customer must return to selling dealer if within 50 miles or nearest Hyundai dealer",
            "Available at time of vehicle purchase/lease only",
        ],
        "exclusions": [
            "Coverage for enrolled vehicle keys/fobs only",
            "Intentional acts, abuse, or unauthorized modification",
            "Specialty or aftermarket keys not meeting manufacturer specs",
            "CA: limited to vehicle key/fob, vehicle towing, vehicle lockout, taxi/rental only — no home lockout, concierge, hotel/entertainment discounts, or emergency relay",
        ],
        "transfer_fee": "$50 (within 30 days, private party or lease assumption only)",
        "cancel_fee": "Within 30 days: full refund less claims. After 30 days: pro-rata less claims + $50 fee",
        "additional_benefits": [],
    },
    {
        "name": "Theft Protection",
        "subtitle": "Theft deterrent system + financial warranty benefit up to $5,000 if vehicle is stolen total loss",
        "code": "HFTP",
        "core": False,
        "ev_care": False,
        "circle": False,
        "overview": "Vehicle is permanently marked with anti-theft identification codes and warning decals. If the vehicle is stolen and declared a total loss (unrecovered), provides a financial warranty benefit up to $5,000. NON-CANCELLABLE.",
        "plans": [
            {"name": "Theft Protection", "desc": "Permanent VIN markings + anti-theft decals. Warranty benefit up to $5,000 if vehicle is stolen and declared a total loss or not recovered. No activation required."},
        ],
        "features": [
            "No deductible",
            "No mileage limitations",
            "No activation required",
            "Benefit checks payable to customer, sent to dealership",
            "Benefit helps cover: insurance deductible, sales tax, tag & title fees, negative equity, down payment for replacement",
            "NON-CANCELLABLE",
            "FL: flat warranty benefit capped at $5,000",
            "CT: dealer must have written consent before VIN etching",
        ],
        "terms_mileage": "60 months (5 years) / Unlimited miles",
        "deductible": "$0",
        "eligibility": [
            "New & pre-owned Hyundai vehicles (ICE, EV, Hybrid, Hydrogen)",
            "Competitive makes including Kia & Genesis sold at Hyundai dealerships",
            "Available at time of vehicle sale or lease only",
        ],
        "exclusions": [
            "Commercial/fleet vehicles (permitted professional use allowed)",
        ],
        "transfer_fee": "$50 (private party; within 30 days with required documents)",
        "cancel_fee": "NON-CANCELLABLE",
        "additional_benefits": [],
    },
]

# ─────────────────────────────────────────────────────────────────────────────
# TOC SLIDE
# ─────────────────────────────────────────────────────────────────────────────

def build_toc(prs):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_bg(slide, WHITE)

    add_rect(slide, 0, 0, 13.33, 1.0, fill_rgb=HYUNDAI_BLUE)
    add_text_box(slide, "TABLE OF CONTENTS",
                 0.3, 0.3, 12.7, 0.6,
                 font_size=26, bold=True, color=WHITE)

    sections = [
        ("MECHANICAL COVERAGE", HYUNDAI_BLUE, [
            "Vehicle Service Protection [ICE]                                  CORE",
            "VSP Competitive Makes                                              NON-CORE",
            "VSP Original Owner [ICE]                                          CORE",
            "CUV Wrap [ICE]                                                    CORE",
            "High Mileage VSP [ICE]                                            NON-CORE",
            "EV Care VSP                                                       CORE",
            "EV Care VSP — Competitive Makes                                   NON-CORE",
            "EV Care VSP — Original Owner                                      CORE",
            "EV Care CUV Wrap                                                  CORE",
            "VSP Livery Coverage [ICE & EV]                                    CORE",
            "Wear Protection [ICE & EV]                                        CORE",
        ]),
        ("MAINTENANCE PLANS", HYUNDAI_SKY, [
            "Maintenance Wrap [ICE]                                            CORE",
            "Pre-Paid Maintenance [ICE]                                        CORE",
            "Pre-Paid Maintenance — Competitive Makes                          NON-CORE",
            "EV Care Maintenance                                               CORE",
        ]),
        ("DAMAGE CARE", EV_TEAL, [
            "Lease-End Protection                                              CORE",
            "Multi-Coverage Protection                                         NON-CORE",
            "Tire & Wheel Protection                                           NON-CORE",
            "Dent Protection                                                   NON-CORE",
            "Windshield Protection                                             NON-CORE",
            "Appearance Protection                                             NON-CORE",
        ]),
        ("LOSS & THEFT PROTECTION", RGBColor(0x6A, 0x1A, 0x4D), [
            "Guaranteed Asset Protection (GAP)                                 CORE",
            "GAP Plus                                                          CORE",
            "Key Protection                                                    NON-CORE",
            "Theft Protection                                                  NON-CORE",
        ]),
    ]

    col_x = [0.25, 7.0]
    col_idx = 0
    y_positions = [1.15, 1.15]

    for sec_title, sec_color, items in sections:
        cx = col_x[col_idx]
        y = y_positions[col_idx]
        cw = 6.5

        add_rect(slide, cx, y, cw, 0.28, fill_rgb=sec_color)
        add_text_box(slide, sec_title, cx + 0.1, y + 0.04, cw - 0.2, 0.22,
                     font_size=9.5, bold=True, color=WHITE)
        y += 0.3

        for item in items:
            parts = item.rsplit("CORE", 1)
            if "NON-CORE" in item:
                badge_txt = "NON-CORE"
                badge_clr = NON_CORE_GREY
                item_name = item.replace("                                   NON-CORE", "").strip()
            elif "CORE" in item:
                badge_txt = "CORE"
                badge_clr = CORE_GOLD
                item_name = item.replace("                                  CORE", "").strip()
            else:
                badge_txt = ""
                badge_clr = None
                item_name = item.strip()

            add_text_box(slide, f"  {item_name}", cx, y, cw - 1.0, 0.21,
                         font_size=8, color=DARK_TEXT)
            if badge_clr:
                add_rect(slide, cx + cw - 0.9, y + 0.01, 0.85, 0.19, fill_rgb=badge_clr)
                add_text_box(slide, badge_txt, cx + cw - 0.9, y + 0.02, 0.85, 0.17,
                             font_size=7, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
            y += 0.21

        y_positions[col_idx] = y + 0.15
        col_idx = 1 - col_idx

    add_text_box(slide,
                 "CORE = HPP Branded Products   |   NON-CORE = White-Label / Non-Branded Products   |   ● = Circle Program Available",
                 0.3, 7.3, 12.7, 0.22,
                 font_size=8, color=MED_GREY, italic=True, align=PP_ALIGN.CENTER)


# ─────────────────────────────────────────────────────────────────────────────
# MAIN
# ─────────────────────────────────────────────────────────────────────────────

def main():
    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)

    # Cover
    build_cover(prs)

    # TOC
    build_toc(prs)

    # Section: Mechanical
    build_section_divider(prs, "MECHANICAL COVERAGE",
                          "Vehicle Service Protection  |  EV Care  |  Wear Protection")

    mech_names = [
        "Vehicle Service Protection [ICE]",
        "Vehicle Service Protection [ICE] — Competitive Makes",
        "Vehicle Service Protection Original Owner [ICE]",
        "Certified Used Vehicle (CUV) Wrap [ICE]",
        "High Mileage Vehicle Service Protection [ICE]",
        "EV Care Vehicle Service Protection",
        "EV Care VSP — Competitive Makes",
        "EV Care VSP — Original Owner",
        "EV Care Certified Used Vehicle (CUV) Wrap",
        "Vehicle Service Protection — Livery Coverage [ICE & EV]",
        "Wear Protection [ICE & EV]",
    ]
    for prod in PRODUCTS:
        if prod["name"] in mech_names:
            build_product_slide(prs, prod)

    # Section: Maintenance
    build_section_divider(prs, "MAINTENANCE PLANS",
                          "ICE  |  EV Care  |  Competitive Makes", bg=HYUNDAI_SKY)
    maint_names = [
        "Maintenance Wrap [ICE]",
        "Pre-Paid Maintenance [ICE]",
        "Pre-Paid Maintenance — Competitive Makes [ICE]",
        "EV Care Maintenance",
    ]
    for prod in PRODUCTS:
        if prod["name"] in maint_names:
            build_product_slide(prs, prod)

    # Section: Damage Care
    build_section_divider(prs, "DAMAGE CARE",
                          "Lease-End  |  Multi-Coverage  |  Tire & Wheel  |  Dent  |  Windshield  |  Appearance",
                          bg=EV_TEAL)
    damage_names = [
        "Lease-End Protection",
        "Multi-Coverage Protection",
        "Tire & Wheel Protection",
        "Dent Protection",
        "Windshield Protection",
        "Appearance Protection",
    ]
    for prod in PRODUCTS:
        if prod["name"] in damage_names:
            build_product_slide(prs, prod)

    # Section: Loss & Theft
    build_section_divider(prs, "LOSS & THEFT PROTECTION",
                          "GAP  |  GAP Plus  |  Key Protection  |  Theft Protection",
                          bg=RGBColor(0x6A, 0x1A, 0x4D))
    loss_names = [
        "Guaranteed Asset Protection (GAP)",
        "Guaranteed Asset Protection Plus (GAP Plus)",
        "Key Protection",
        "Theft Protection",
    ]
    for prod in PRODUCTS:
        if prod["name"] in loss_names:
            build_product_slide(prs, prod)

    # Save
    out_path = "/Users/justin/Downloads/HPP Product Framework.pptx"
    prs.save(out_path)
    print(f"Saved: {out_path}")


if __name__ == "__main__":
    main()
