"""
HPP Product Framework — V2
Replicates the exact slide 3 formatting from PRES HPP.pptx:
  - Slide dims: 13.3333" x 7.5"
  - Title: placeholder-style, ~36pt
  - Banner bar: #002060, full width, T=1.33", H=0.54"
  - Overview textbox over banner: 12pt Aptos, bold label + regular body, white text
  - Section headers: 16pt bold Aptos (left-aligned, no fill)
  - Tables: 2-col, left col #EBF0F5 narrow, right col #FFFFFF, 9pt Calibri
  - Right column boxes (Exclusions, Transfers, Cancellations): AUTO_SHAPE fills
      Exclusions/Transfers: #EBF1F6   Cancellations: #FFEDEB
  - Corner code tag: 10pt Aptos, color=#1F497D, right-aligned
  - Circle callout box: distinct orange callout
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.oxml import parse_xml
from lxml import etree
import copy

# ── Exact colors from slide 3 ──────────────────────────────────────────────
C_NAVY          = RGBColor(0x00, 0x20, 0x60)   # #002060 banner
C_PLAN_LABEL    = RGBColor(0x00, 0x2C, 0x5F)   # #002C5F plan name text
C_PLAN_BODY     = RGBColor(0x1A, 0x1A, 0x1A)   # #1A1A1A plan detail text
C_CODE_TAG      = RGBColor(0x1F, 0x49, 0x7D)   # #1F497D corner tag
C_TBL_LEFT      = RGBColor(0xEB, 0xF0, 0xF5)   # #EBF0F5 table left col
C_TBL_RIGHT     = RGBColor(0xFF, 0xFF, 0xFF)   # #FFFFFF table right col
C_BOX_BLUE      = RGBColor(0xEB, 0xF1, 0xF6)   # #EBF1F6 exclusions/transfers
C_BOX_RED       = RGBColor(0xFF, 0xED, 0xEB)   # #FFEDEB cancellations
C_WHITE         = RGBColor(0xFF, 0xFF, 0xFF)
C_DARK          = RGBColor(0x1A, 0x1A, 0x1A)
C_CIRCLE_FILL   = RGBColor(0xE8, 0x6F, 0x1E)   # orange circle callout
C_CIRCLE_TXT    = RGBColor(0xFF, 0xFF, 0xFF)
C_CORE_GOLD     = RGBColor(0xC8, 0xA0, 0x32)
C_NON_CORE      = RGBColor(0x6D, 0x6D, 0x6D)
C_EV_TEAL       = RGBColor(0x00, 0x70, 0x70)
C_LIGHT_BG      = RGBColor(0xF4, 0xF6, 0xF9)
C_MED_GREY      = RGBColor(0x9E, 0x9E, 0x9E)

# Fonts
F_HEADER = "Aptos"
F_BODY   = "Calibri"

def I(x): return Inches(x)
def P(x): return Pt(x)

# ── Low-level helpers ──────────────────────────────────────────────────────

def set_cell_fill(cell, rgb):
    tc = cell._tc
    tcPr = tc.find(qn('a:tcPr'))
    if tcPr is None:
        tcPr = etree.SubElement(tc, qn('a:tcPr'))
    solidFill = tcPr.find(qn('a:solidFill'))
    if solidFill is not None:
        tcPr.remove(solidFill)
    solidFill = etree.SubElement(tcPr, qn('a:solidFill'))
    srgbClr = etree.SubElement(solidFill, qn('a:srgbClr'))
    srgbClr.set('val', str(rgb))

def set_cell_border_none(cell):
    """Remove all borders from a table cell."""
    tc = cell._tc
    tcPr = tc.find(qn('a:tcPr'))
    if tcPr is None:
        tcPr = etree.SubElement(tc, qn('a:tcPr'))
    for border_tag in ('a:lnL', 'a:lnR', 'a:lnT', 'a:lnB'):
        existing = tcPr.find(qn(border_tag))
        if existing is not None:
            tcPr.remove(existing)
        ln = etree.SubElement(tcPr, qn(border_tag))
        ln.set('w', '0')
        noFill = etree.SubElement(ln, qn('a:noFill'))

def add_cell_run(para, text, bold=False, size_pt=9, color=None, font_name=F_BODY):
    run = para.add_run()
    run.text = text
    run.font.size = P(size_pt)
    run.font.bold = bold
    run.font.name = font_name
    if color:
        run.font.color.rgb = color

def set_bg(slide, rgb=None):
    if rgb is None:
        return
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = rgb

def add_rect_shape(slide, left, top, width, height, fill_rgb=None, line=False):
    shape = slide.shapes.add_shape(1, I(left), I(top), I(width), I(height))
    if fill_rgb:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_rgb
    else:
        shape.fill.background()
    if not line:
        shape.line.fill.background()
    return shape

def add_textbox(slide, left, top, width, height):
    txBox = slide.shapes.add_textbox(I(left), I(top), I(width), I(height))
    txBox.word_wrap = True
    return txBox

def clear_para(para):
    for run in para.runs:
        r = run._r
        r.getparent().remove(r)

def set_para_align(para, align):
    para.alignment = align

# ── Section-header textbox (16pt bold Aptos) ───────────────────────────────
def add_section_header(slide, text, left, top, width, height=0.37,
                        color=None, align=PP_ALIGN.LEFT):
    tb = add_textbox(slide, left, top, width, height)
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = P(16)
    r.font.bold = True
    r.font.name = F_HEADER
    if color:
        r.font.color.rgb = color
    return tb

# ── Two-column table helper ────────────────────────────────────────────────
def add_two_col_table(slide, rows_data, left, top, width, height,
                       left_col_frac=0.25):
    """
    rows_data: list of (left_text, right_parts)
      left_text: str (shown in left column, small, colored)
      right_parts: list of (text, bold) tuples
    left_col_frac: fraction of width for left col (narrow)
    """
    n_rows = len(rows_data)
    if n_rows == 0:
        return None
    col_widths = [I(width * left_col_frac), I(width * (1 - left_col_frac))]
    tbl = slide.shapes.add_table(n_rows, 2, I(left), I(top), I(width), I(height)).table
    tbl.columns[0].width = col_widths[0]
    tbl.columns[1].width = col_widths[1]

    for ri, (left_text, right_parts) in enumerate(rows_data):
        lc = tbl.cell(ri, 0)
        rc = tbl.cell(ri, 1)

        # Alternating fill on left col
        if ri % 2 == 0:
            set_cell_fill(lc, C_TBL_LEFT)
        else:
            set_cell_fill(lc, C_TBL_LEFT)  # always same in source
        set_cell_fill(rc, C_TBL_RIGHT)

        # Left col text (plan name style)
        if left_text:
            lp = lc.text_frame.paragraphs[0]
            lp.alignment = PP_ALIGN.RIGHT
            parts = left_text.split('\n')
            first = True
            for pt in parts:
                if first:
                    cur_p = lp
                    first = False
                else:
                    cur_p = lc.text_frame.add_paragraph()
                    cur_p.alignment = PP_ALIGN.RIGHT
                r = cur_p.add_run()
                r.text = pt
                r.font.size = P(8.5)
                r.font.bold = True
                r.font.name = F_BODY
                r.font.color.rgb = C_PLAN_LABEL

        # Right col text
        rp = rc.text_frame.paragraphs[0]
        rp.alignment = PP_ALIGN.LEFT
        for (text, bold) in right_parts:
            add_cell_run(rp, text, bold=bold, size_pt=9, font_name=F_BODY)

        # Remove borders
        set_cell_border_none(lc)
        set_cell_border_none(rc)

    return tbl

# ── Filled box with body text (Exclusions, Transfers, Cancellations style) ─
def add_body_box(slide, paras_data, left, top, width, height, fill_rgb):
    """
    paras_data: list of list of (text, bold) — one list per paragraph
    """
    shape = add_rect_shape(slide, left, top, width, height, fill_rgb=fill_rgb)
    tf = shape.text_frame
    tf.word_wrap = True
    first = True
    for para_runs in paras_data:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.space_before = P(2)
        p.space_after = P(2)
        p.alignment = PP_ALIGN.LEFT
        for (text, bold) in para_runs:
            add_cell_run(p, text, bold=bold, size_pt=9, font_name=F_BODY)
    return shape

# ── Circle callout box ─────────────────────────────────────────────────────
def add_circle_callout(slide, circle_data, left, top, width):
    """
    circle_data: {code, benefit_desc, restrictions: [str]}
    Returns bottom y position after callout.
    """
    # Estimate height
    line_count = 2 + len(circle_data.get('restrictions', []))
    h = max(0.7, 0.18 * (line_count + 2))

    # Orange header
    hdr = add_rect_shape(slide, left, top, width, 0.26, fill_rgb=C_CIRCLE_FILL)
    tf_h = hdr.text_frame
    tf_h.word_wrap = True
    p_h = tf_h.paragraphs[0]
    p_h.alignment = PP_ALIGN.LEFT
    r1 = p_h.add_run()
    r1.text = "● CIRCLE PROGRAM"
    r1.font.size = P(9)
    r1.font.bold = True
    r1.font.name = F_HEADER
    r1.font.color.rgb = C_WHITE
    # code on same line
    if circle_data.get('code'):
        r2 = p_h.add_run()
        r2.text = f"  •  Code: {circle_data['code']}"
        r2.font.size = P(8.5)
        r2.font.bold = False
        r2.font.name = F_HEADER
        r2.font.color.rgb = RGBColor(0xFF, 0xE0, 0xC0)

    # Body box
    body_top = top + 0.26
    body_h = h
    body = add_rect_shape(slide, left, body_top, width, body_h,
                           fill_rgb=RGBColor(0xFF, 0xF4, 0xEC))
    shape = body
    tf = shape.text_frame
    tf.word_wrap = True
    first = True

    # Benefit
    if circle_data.get('benefit'):
        p = tf.paragraphs[0]
        first = False
        p.alignment = PP_ALIGN.LEFT
        p.space_before = P(2)
        p.space_after = P(1)
        r = p.add_run()
        r.text = circle_data['benefit']
        r.font.size = P(8.5)
        r.font.bold = False
        r.font.name = F_BODY
        r.font.color.rgb = C_DARK

    # Restrictions
    for restr in circle_data.get('restrictions', []):
        p = tf.add_paragraph() if not first else tf.paragraphs[0]
        if first:
            first = False
        p.alignment = PP_ALIGN.LEFT
        p.space_before = P(1)
        p.space_after = P(1)
        r = p.add_run()
        r.text = f"• {restr}"
        r.font.size = P(8)
        r.font.bold = False
        r.font.name = F_BODY
        r.font.color.rgb = RGBColor(0x60, 0x30, 0x00)

    return top + 0.26 + body_h + 0.08

# ── Cover slide ────────────────────────────────────────────────────────────
def build_cover(prs):
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    set_bg(slide, C_NAVY)

    add_rect_shape(slide, 0, 0, 13.333, 0.07, fill_rgb=RGBColor(0x00, 0x70, 0xC0))

    tb = add_textbox(slide, 0.5, 1.8, 12.33, 0.8)
    p = tb.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = "HYUNDAI PROTECTION PLAN"
    r.font.size = P(40)
    r.font.bold = True
    r.font.name = F_HEADER
    r.font.color.rgb = C_WHITE

    tb2 = add_textbox(slide, 0.5, 2.65, 12.33, 0.5)
    p2 = tb2.text_frame.paragraphs[0]
    p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run()
    r2.text = "Product Framework Guide"
    r2.font.size = P(22)
    r2.font.bold = False
    r2.font.name = F_HEADER
    r2.font.color.rgb = RGBColor(0xB0, 0xC8, 0xE8)

    add_rect_shape(slide, 2.5, 3.3, 8.33, 0.04, fill_rgb=RGBColor(0x00, 0x70, 0xC0))

    tb3 = add_textbox(slide, 0.5, 3.45, 12.33, 0.4)
    p3 = tb3.text_frame.paragraphs[0]
    p3.alignment = PP_ALIGN.CENTER
    r3 = p3.add_run()
    r3.text = "Core & Non-Core Products  |  F&I Dealer Reference"
    r3.font.size = P(12)
    r3.font.name = F_HEADER
    r3.font.color.rgb = RGBColor(0xB0, 0xC8, 0xE8)

    cats = [
        ("MECHANICAL COVERAGE", 4.1),
        ("MAINTENANCE PLANS", 4.65),
        ("DAMAGE CARE", 5.2),
        ("LOSS & THEFT PROTECTION", 5.75),
    ]
    for label, y in cats:
        add_rect_shape(slide, 1.5, y, 10.33, 0.42, fill_rgb=RGBColor(0x00, 0x40, 0x80))
        tb = add_textbox(slide, 1.6, y + 0.05, 10.13, 0.32)
        p = tb.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = label
        r.font.size = P(14)
        r.font.bold = True
        r.font.name = F_HEADER
        r.font.color.rgb = C_WHITE

    # Legend
    for x, label, col in [(0.5, "● CORE", C_CORE_GOLD),
                            (1.8, "● NON-CORE", C_NON_CORE),
                            (3.3, "● CIRCLE PROGRAM", C_CIRCLE_FILL)]:
        add_rect_shape(slide, x, 6.7, 1.35, 0.25, fill_rgb=col)
        tb = add_textbox(slide, x + 0.05, 6.72, 1.25, 0.2)
        p = tb.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = label
        r.font.size = P(8.5)
        r.font.bold = True
        r.font.name = F_HEADER
        r.font.color.rgb = C_WHITE

    add_textbox_simple(slide, "For Internal Use Only", 0.5, 7.2, 12.33, 0.25,
                        size=9, color=RGBColor(0x80, 0xA0, 0xC0), align=PP_ALIGN.CENTER)

def add_textbox_simple(slide, text, left, top, width, height,
                        size=10, bold=False, italic=False,
                        color=C_DARK, align=PP_ALIGN.LEFT, font=F_HEADER):
    tb = add_textbox(slide, left, top, width, height)
    p = tb.text_frame.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = P(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.name = font
    r.font.color.rgb = color
    return tb

# ── Section divider ────────────────────────────────────────────────────────
def build_section_divider(prs, title, subtitle="", bg=C_NAVY):
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    set_bg(slide, bg)
    add_rect_shape(slide, 0, 0, 13.333, 0.07, fill_rgb=RGBColor(0x00, 0x70, 0xC0))
    add_textbox_simple(slide, "HYUNDAI PROTECTION PLAN",
                        0.5, 2.4, 12.33, 0.5, size=14,
                        color=RGBColor(0xB0, 0xC8, 0xE8), align=PP_ALIGN.CENTER)
    add_textbox_simple(slide, title,
                        0.5, 3.0, 12.33, 1.0, size=42, bold=True,
                        color=C_WHITE, align=PP_ALIGN.CENTER)
    if subtitle:
        add_rect_shape(slide, 2.5, 4.15, 8.33, 0.04, fill_rgb=RGBColor(0x00, 0x70, 0xC0))
        add_textbox_simple(slide, subtitle,
                            0.5, 4.28, 12.33, 0.4, size=12,
                            color=RGBColor(0xB0, 0xC8, 0xE8), align=PP_ALIGN.CENTER)
    add_textbox_simple(slide, "For Internal Use Only",
                        0.5, 7.1, 12.33, 0.3, size=9,
                        color=RGBColor(0x60, 0x80, 0xA0), align=PP_ALIGN.CENTER)

# ── TOC ────────────────────────────────────────────────────────────────────
def build_toc(prs):
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)

    add_rect_shape(slide, 0, 0, 13.333, 1.0, fill_rgb=C_NAVY)
    add_textbox_simple(slide, "TABLE OF CONTENTS",
                        0.3, 0.25, 12.7, 0.65, size=28, bold=True,
                        color=C_WHITE, align=PP_ALIGN.LEFT)

    sections = [
        ("MECHANICAL COVERAGE", C_NAVY, [
            ("Vehicle Service Protection [ICE]", True),
            ("VSP — Competitive Makes", False),
            ("VSP Original Owner [ICE]", True),
            ("CUV Wrap [ICE]", True),
            ("High Mileage VSP [ICE]", False),
            ("EV Care VSP", True),
            ("EV Care VSP — Competitive Makes", False),
            ("EV Care VSP — Original Owner", True),
            ("EV Care CUV Wrap", True),
            ("VSP Livery Coverage [ICE & EV]", True),
            ("Wear Protection [ICE & EV]", True),
        ]),
        ("MAINTENANCE PLANS", RGBColor(0x00, 0x60, 0xA0), [
            ("Maintenance Basic Wrap [ICE]", True),
            ("Maintenance Scheduled Wrap [ICE]", True),
            ("Pre-Paid Maintenance [ICE]", True),
            ("Pre-Paid Maintenance — Competitive Makes", False),
            ("EV Care Maintenance", True),
        ]),
        ("DAMAGE CARE", C_EV_TEAL, [
            ("Lease-End Protection", True),
            ("Multi-Coverage Protection", False),
            ("Appearance Protection", False),
        ]),
        ("LOSS & THEFT PROTECTION", RGBColor(0x6A, 0x1A, 0x4D), [
            ("Guaranteed Asset Protection (GAP)", True),
            ("GAP Plus", True),
            ("Theft Protection", False),
        ]),
    ]

    col_x = [0.25, 6.9]
    col_y = [1.08, 1.08]
    col = 0
    COL_W = 6.45

    for sec_title, sec_color, items in sections:
        x = col_x[col]
        y = col_y[col]
        add_rect_shape(slide, x, y, COL_W, 0.28, fill_rgb=sec_color)
        add_textbox_simple(slide, sec_title, x + 0.1, y + 0.04, COL_W - 0.2, 0.22,
                            size=9.5, bold=True, color=C_WHITE)
        y += 0.3
        for (item_name, is_core) in items:
            add_textbox_simple(slide, f"  {item_name}", x, y, COL_W - 1.0, 0.21,
                                size=8, color=C_DARK, font=F_BODY)
            bc = C_CORE_GOLD if is_core else C_NON_CORE
            bl = "CORE" if is_core else "NON-CORE"
            add_rect_shape(slide, x + COL_W - 0.92, y + 0.01, 0.88, 0.19, fill_rgb=bc)
            tb = add_textbox(slide, x + COL_W - 0.92, y + 0.02, 0.88, 0.17)
            p = tb.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            r = p.add_run()
            r.text = bl
            r.font.size = P(7)
            r.font.bold = True
            r.font.name = F_HEADER
            r.font.color.rgb = C_WHITE
            y += 0.21
        col_y[col] = y + 0.14
        col = 1 - col

    add_textbox_simple(slide,
        "CORE = HPP Branded Products   |   NON-CORE = White-Label Products   |   ● = Circle Program Available",
        0.3, 7.28, 12.7, 0.22, size=8, italic=True,
        color=C_MED_GREY, align=PP_ALIGN.CENTER, font=F_BODY)

# ── MAIN PRODUCT SLIDE — replicates slide 3 layout exactly ────────────────

def build_product_slide(prs, d):
    """
    d keys:
      name, code, core (bool), ev_care (bool), circle (dict or None),
      overview,
      plans: [{name, sub, desc}]   sub = optional subscript like "(Stated Component)"
      term_coverage: [str]         rows for Term Coverage table
      eligibility: [(label, [(text,bold)])]
      reimbursement: [(label, [(text,bold)])]
      additional_benefits: [(label, [(text,bold)])]
      exclusions: [str]
      transfer: str
      cancel: [(str, bool)]   list of (text, bold) for cancel box
      circle: {code, benefit, restrictions:[str]}  or None
    """
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)

    is_ev = d.get("ev_care", False)
    banner_color = C_NAVY  # always navy per source

    # ── Title ────────────────────────────────────────────────────────────
    tb_title = add_textbox(slide, 0.667, 0.30, 12.0, 0.89)
    pt = tb_title.text_frame.paragraphs[0]
    r = pt.add_run()
    r.text = d["name"]
    r.font.size = P(36)
    r.font.name = F_HEADER
    # color inherits (white background reads dark automatically) — match source

    # ── Banner bar ───────────────────────────────────────────────────────
    banner = add_rect_shape(slide, -0.002, 1.331, 13.338, 0.541, fill_rgb=banner_color)

    # EV Care sub-label on banner if applicable
    if is_ev:
        add_textbox_simple(slide, "EV Care", 0.3, 1.35, 2.0, 0.35,
                            size=10, bold=True, color=RGBColor(0x70, 0xFF, 0xFF))

    # ── Overview text (over banner) ───────────────────────────────────────
    tb_ov = add_textbox(slide, 0.829, 1.331, 9.5, 0.505)
    tb_ov.text_frame.word_wrap = True
    p_ov = tb_ov.text_frame.paragraphs[0]
    r1 = p_ov.add_run()
    r1.text = "Product Overview: "
    r1.font.size = P(12)
    r1.font.bold = True
    r1.font.name = F_HEADER
    r1.font.color.rgb = C_WHITE
    r2 = p_ov.add_run()
    r2.text = d.get("overview", "")
    r2.font.size = P(12)
    r2.font.bold = False
    r2.font.name = F_HEADER
    r2.font.color.rgb = C_WHITE

    # ── Code / CORE tag (top right, over banner) ──────────────────────────
    core_label = f"{d.get('code','')}  |  {'CORE' if d.get('core', True) else 'NON-CORE'}"
    core_color = C_CORE_GOLD if d.get('core', True) else C_NON_CORE
    tb_tag = add_textbox(slide, 10.786, 1.378, 2.4, 0.4)
    p_tag = tb_tag.text_frame.paragraphs[0]
    p_tag.alignment = PP_ALIGN.RIGHT
    r_tag = p_tag.add_run()
    r_tag.text = core_label
    r_tag.font.size = P(10)
    r_tag.font.name = F_HEADER
    r_tag.font.color.rgb = C_WHITE

    # ── Layout constants (matching slide 3 positions) ─────────────────────
    # Left block  (Plans + Term Coverage)
    L1_X  = 0.496
    L1_W  = 4.418
    # Middle block (Eligibility + Reimbursement + Additional Benefits)
    L2_X  = 5.207
    L2_W  = 4.501
    # Right block  (Exclusions + Transfers + Cancellations + Circle)
    L3_X  = 9.958
    L3_W  = 2.846
    # Rows start below banner
    ROW_START = 2.157

    # Left col fraction inside the 2-col tables
    LEFT_FRAC = 0.18  # ~18% for label, 82% for content (matches source narrow left col)

    # ── LEFT COLUMN: Plans ───────────────────────────────────────────────
    y_l = ROW_START
    add_section_header(slide, "Plans", L1_X, y_l, L1_W)
    y_l += 0.37 + 0.02

    # Build plans rows
    plan_rows = []
    for pl in d.get("plans", []):
        label = pl["name"]
        if pl.get("sub"):
            label += "\n" + pl["sub"]
        right_parts = [(pl.get("desc", ""), False)]
        plan_rows.append((label, right_parts))

    plan_h = max(0.38 * len(plan_rows), 0.5)
    if plan_rows:
        add_two_col_table(slide, plan_rows, L1_X, y_l, L1_W, plan_h,
                          left_col_frac=LEFT_FRAC)
    y_l += plan_h + 0.12

    # ── LEFT COLUMN: Term Coverage ───────────────────────────────────────
    add_section_header(slide, "Term Coverage", L1_X, y_l, L1_W)
    y_l += 0.37 + 0.02

    term_rows = [(("", [(" " + t, False)]) if True else ("", [(t, False)]))
                 for t in d.get("term_coverage", [])]
    # Simplify: blank left col, right = text
    term_rows2 = [("", [(t, False)]) for t in d.get("term_coverage", [])]
    term_h = max(0.28 * len(term_rows2), 0.35)
    if term_rows2:
        add_two_col_table(slide, term_rows2, L1_X, y_l, L1_W, term_h,
                          left_col_frac=LEFT_FRAC)
    y_l += term_h + 0.08

    # ── MIDDLE COLUMN: Eligibility ───────────────────────────────────────
    y_m = ROW_START
    elig_label = d.get("elig_label", "Eligibility")
    add_section_header(slide, elig_label, L2_X, y_m, L2_W)
    y_m += 0.37 + 0.02

    elig_rows = [("", [(t, False)]) for t in d.get("eligibility", [])]
    elig_h = max(0.28 * len(elig_rows), 0.35)
    if elig_rows:
        add_two_col_table(slide, elig_rows, L2_X, y_m, L2_W, elig_h,
                          left_col_frac=LEFT_FRAC)
    y_m += elig_h + 0.12

    # ── MIDDLE COLUMN: Reimbursement ─────────────────────────────────────
    add_section_header(slide, "Reimbursement", L2_X, y_m, L2_W)
    y_m += 0.37 + 0.02

    reimb_rows = []
    for (lbl, parts) in d.get("reimbursement", []):
        reimb_rows.append(("", parts))
    if not reimb_rows:
        reimb_rows = [("", [("MSRP, replacement parts may be new or remanufactured like kind and quality (administrator discretion)", False)]),
                      ("", [("Labor: ", False), ("dealer door rate", True)])]
    reimb_h = max(0.28 * len(reimb_rows), 0.35)
    add_two_col_table(slide, reimb_rows, L2_X, y_m, L2_W, reimb_h,
                      left_col_frac=LEFT_FRAC)
    y_m += reimb_h + 0.12

    # ── MIDDLE COLUMN: Additional Benefits ───────────────────────────────
    add_section_header(slide, "Additional Benefits", L2_X, y_m, L2_W)
    y_m += 0.37 + 0.02

    ab = d.get("additional_benefits", [])
    if not ab:
        ab = [
            [("24 Hour Roadside Assistance", True)],
            [("Rental car: up to ", False), ("$55/day", True), (", 10 days max", False)],
            [("Trip interruption up to ", False), ("$300/day", True), (" x 5 Days", False)],
            [("Diagnostics: ", False), ("Paid in conjunction with covered repair, capped to labor guide time", False)],
            [("Replacement fluids covered in conjunction with a covered repair", False)],
        ]
    ab_rows = [("", row_parts) for row_parts in ab]
    ab_h = max(0.26 * len(ab_rows), 0.35)
    add_two_col_table(slide, ab_rows, L2_X, y_m, L2_W, ab_h,
                      left_col_frac=LEFT_FRAC)
    y_m += ab_h + 0.08

    # ── RIGHT COLUMN: Exclusions ─────────────────────────────────────────
    y_r = ROW_START
    add_section_header(slide, "Exclusions", L3_X, y_r, L3_W)
    y_r += 0.37 + 0.02

    excl_paras = [[(e, False)] for e in d.get("exclusions", [])]
    excl_h = max(0.22 * len(excl_paras), 0.35)
    if excl_paras:
        add_body_box(slide, excl_paras, L3_X, y_r, L3_W, excl_h, C_BOX_BLUE)
    y_r += excl_h + 0.37 + 0.04   # gap for Transfers header

    # ── RIGHT COLUMN: Transfers ──────────────────────────────────────────
    add_section_header(slide, "Transfers", L3_X, y_r - 0.37 - 0.02, L3_W)
    transfer_text = d.get("transfer", "$75 (private party sale or lease assumption only, not transferable if dealer is party to resale/assumption)")
    # Parse bold for $75
    tf_parts = []
    if "$75" in transfer_text:
        idx = transfer_text.index("$75")
        tf_parts = [(transfer_text[:idx], False), ("$75", True), (transfer_text[idx+3:], False)]
    elif "$50" in transfer_text:
        idx = transfer_text.index("$50")
        tf_parts = [(transfer_text[:idx], False), ("$50", True), (transfer_text[idx+3:], False)]
    else:
        tf_parts = [(transfer_text, False)]

    xfer_h = 0.72
    add_body_box(slide, [tf_parts], L3_X, y_r, L3_W, xfer_h, C_BOX_BLUE)
    y_r += xfer_h + 0.37 + 0.06

    # ── RIGHT COLUMN: Cancellations ──────────────────────────────────────
    add_section_header(slide, "Cancellations", L3_X, y_r - 0.37 - 0.04, L3_W)
    cancel_data = d.get("cancel", [
        [("Cancel within 30 days: full refund less claims paid (state exceptions may apply)", False)],
        [("Cancel after 30 days: pro-rata less claims, $75 processing fee (state exceptions may apply)", False)],
    ])
    cancel_h = max(0.28 * len(cancel_data), 0.5)
    add_body_box(slide, cancel_data, L3_X, y_r, L3_W, cancel_h, C_BOX_RED)
    y_r += cancel_h + 0.1

    # ── RIGHT COLUMN: Circle callout (if applicable) ──────────────────────
    circle = d.get("circle")
    if circle:
        # Check if we have room; if not, overlay below cancellations
        add_section_header(slide, "Circle Program", L3_X, y_r, L3_W)
        y_r += 0.37 + 0.02
        add_circle_callout(slide, circle, L3_X, y_r, L3_W)

    # ── Footer ───────────────────────────────────────────────────────────
    add_rect_shape(slide, 0, 7.3, 13.333, 0.2, fill_rgb=C_LIGHT_BG)
    add_textbox_simple(slide,
        "Some terms vary by state and/or contract revision date. For illustrative purposes only. See applicable product agreement for details.",
        0.2, 7.32, 13.0, 0.18, size=6.5, italic=True, color=C_MED_GREY,
        align=PP_ALIGN.LEFT, font=F_BODY)


# ── ALL PRODUCT DATA ───────────────────────────────────────────────────────

STD_REIMB = [
    ("", [("MSRP, replacement parts may be new or remanufactured like kind and quality (administrator discretion)", False)]),
    ("", [("Labor: ", False), ("dealer door rate", True)]),
]
STD_AB = [
    [("24 Hour Roadside Assistance", True)],
    [("Rental car: up to ", False), ("$55/day", True), (", 10 days max", False)],
    [("Trip interruption up to ", False), ("$300/day", True), (" x 5 Days", False)],
    [("Diagnostics: ", False), ("Paid in conjunction with a covered repair, capped to labor guide time", False)],
    [("Replacement fluids covered in conjunction with a covered repair", False)],
]
STD_CANCEL_75 = [
    [("Cancel within 30 days: full refund less claims paid (state exceptions may apply)", False)],
    [("Cancel after 30 days: pro-rata less claims, ", False), ("$75", True), (" processing fee (state exceptions may apply)", False)],
]
STD_CANCEL_50 = [
    [("Cancel within 30 days: full refund less claims paid (state exceptions may apply)", False)],
    [("Cancel after 30 days: pro-rata less claims, ", False), ("$50", True), (" processing fee (fee cannot exceed refund)", False)],
]

PRODUCTS = [
    # ════════════════ MECHANICAL ════════════════
    {
        "name": "Vehicle Service Protection (ICE)",
        "code": "HFVI", "core": True, "ev_care": False,
        "overview": "Covers eligible mechanical breakdown repairs for covered parts and labor, with term and mileage options up to 10 years/150,000 miles.",
        "plans": [
            {"name": "Powertrain", "sub": "(Stated Component)", "desc": "Engine, Transmission, Drive Axle incl. CV joints"},
            {"name": "Gold", "sub": "(Stated Cmpnt.)", "desc": "Powertrain plus front/rear suspension (incl shocks), A/C, fuel system, electrical system"},
            {"name": "High Technology", "sub": "(Stated Cmpnt.)", "desc": "Factory-installed ADAS, audio/visual, and accessories (USB ports, 120V ports, wireless charging, WiFi/hotspot, mobile interfaces)"},
            {"name": "Platinum", "sub": "(Exclusionary)", "desc": "Covers all Covered Parts except items listed in Exclusions"},
            {"name": "Platinum Plus", "sub": "(Florida Only)", "desc": "Platinum plus Term Protection items, headlamps, belts and hoses, electrical"},
        ],
        "term_coverage": [
            "Max options up to 10 years or 150,000 miles (from contract sale date)",
            "Additive coverage: begins day contract is sold, mileage based on odometer at sale",
            "Deductible options: $0 or $100",
        ],
        "elig_label": "VSP Eligibility",
        "eligibility": [
            "New or pre-owned ICE only (not Hybrid or EV)",
            "9 model years or less, under 140,001 miles at time of purchase",
            "Must have at least 1 month and 1,000 miles of manufacturer full warranty remaining if purchased after vehicle sale",
            "UCI required at 50,001+ miles or over 4.5 years old",
        ],
        "reimbursement": STD_REIMB,
        "additional_benefits": STD_AB,
        "exclusions": [
            "Total loss, salvaged, junk, buy-back vehicles not eligible",
            "Non-U.S. specification vehicles not eligible (commonly Canada-import vehicles)",
            "Commercial and fleet excluded unless it qualifies as a Permitted Commercial Purpose",
            "See agreement for full exclusion details",
        ],
        "transfer": "Transferable: $75 (private party sale or lease assumption only, not transferable if dealer is party to resale/assumption)",
        "cancel": STD_CANCEL_75,
        "circle": {
            "code": "HFVI",
            "benefit": "Available for employees, family members, and affiliates.",
            "restrictions": [
                "Post-sale purchase requires 1 month warranty remaining and 1,000 miles full warranty remaining",
                "UCI required for any pre-owned",
                "Not available in FL",
                "Max markup: $100 over dealer cost",
            ],
        },
    },
    {
        "name": "Vehicle Service Protection (ICE) — Competitive Makes",
        "code": "WFVI", "core": False, "ev_care": False,
        "overview": "PowerProtect mechanical coverage for pre-owned ICE competitive make (non-Hyundai) vehicles, parts & labor, up to 10 years or 120,000 miles.",
        "plans": [
            {"name": "Powertrain", "sub": "(Stated Component)", "desc": "Engine, Transmission, Drive Axle incl. CV joints"},
            {"name": "Gold", "sub": "(Stated Cmpnt.)", "desc": "Powertrain plus front/rear suspension (incl shocks), A/C, fuel system, electrical system"},
            {"name": "Platinum", "sub": "(Exclusionary)", "desc": "Covers all Covered Parts except items listed in Exclusions"},
            {"name": "Platinum Plus", "sub": "(Florida Only)", "desc": "Platinum plus headlamps, belts and hoses, electrical"},
        ],
        "term_coverage": [
            "Max options up to 10 years or 120,000 miles",
            "Additive coverage: begins day contract is sold",
            "Deductible options: $0, $100, or $250",
        ],
        "elig_label": "VSP Eligibility",
        "eligibility": [
            "Pre-owned ICE competitive make (non-Hyundai) vehicles only",
            "9 model years or less, under 120,000 miles at time of purchase",
            "UCI required for any pre-owned",
        ],
        "reimbursement": STD_REIMB,
        "additional_benefits": STD_AB,
        "exclusions": [
            "Hyundai vehicles — use standard HPP VSP",
            "Total loss, salvaged, junk, buy-back vehicles not eligible",
            "Non-U.S. specification vehicles not eligible",
            "High Technology coverage NOT available for Competitive Makes",
        ],
        "transfer": "Transferable: $75 (private party sale or lease assumption only)",
        "cancel": STD_CANCEL_75,
        "circle": None,
    },
    {
        "name": "Vehicle Service Protection Original Owner (ICE)",
        "code": "HFOI", "core": True, "ev_care": False,
        "overview": "Built for lessees buying their vehicle and original Hyundai owners with vehicles under 59,000 miles. Coverage measured from original in-service date.",
        "plans": [
            {"name": "Platinum", "sub": "(Exclusionary)", "desc": "Covers all Covered Parts except items listed in Exclusions"},
            {"name": "Platinum Plus", "sub": "(Florida Only)", "desc": "Platinum plus headlamps, belts and hoses, electrical"},
        ],
        "term_coverage": [
            "Coverage available up to 10 years/150,000 miles from original in-service date",
            "Deductible: $0 or $100",
        ],
        "elig_label": "VSP Eligibility",
        "eligibility": [
            "New or pre-owned ICE only (not Hybrid or EV)",
            "Vehicle mileage at sale: 10,001 to 58,999",
            "Original Owner requirement: currently owned by original owner or purchased by original lessee",
            "UCI required for any pre-owned",
        ],
        "reimbursement": STD_REIMB,
        "additional_benefits": STD_AB,
        "exclusions": [
            "Total loss, salvaged, junk, buy-back vehicles not eligible",
            "Non-U.S. specification vehicles not eligible",
            "Commercial and fleet excluded unless it qualifies as a Permitted Commercial Purpose",
            "See agreement for full exclusion details",
        ],
        "transfer": "Transferable: $75 (private party sale or lease assumption only, not transferable if dealer is party to resale/assumption)",
        "cancel": STD_CANCEL_75,
        "circle": None,
    },
    {
        "name": "Certified Used Vehicle (CUV) Wrap",
        "code": "HFCI", "core": True, "ev_care": False,
        "overview": "Wrap coverage that enhances the Hyundai CUV powertrain warranty on eligible Hyundai Certified Used (CUV) vehicles, covering mechanical repairs, parts, and labor.",
        "plans": [
            {"name": "Gold", "sub": "(Stated Cmpnt.)", "desc": "Front/rear suspension (incl shocks), air conditioning, fuel system, electrical system"},
            {"name": "Platinum", "sub": "(Exclusionary)", "desc": "All Covered Parts except those listed as exclusions"},
            {"name": "Platinum Plus", "sub": "(Florida Only)", "desc": "Platinum plus headlamp, belts and hoses, electrical"},
        ],
        "term_coverage": [
            "Max term and mileage: up to 10 years or 150,000 miles from original in-service date",
            "Coverage begins from original in-service date and zero miles",
            "Deductible options: $0, $50 (disappearing at selling dealer), or $100",
        ],
        "elig_label": "CPOW Eligibility",
        "eligibility": [
            "Hyundai CUV Vehicle Only with less than 6 years from the original in-service date and less than 80,000 miles",
            "Customer must have an active CPO / CUV certification to qualify",
            "UCI required for any pre-owned",
            "Vehicle cannot exceed 80,000 miles",
        ],
        "reimbursement": STD_REIMB,
        "additional_benefits": STD_AB,
        "exclusions": [
            "Vehicle cannot exceed 80,000 miles",
            "Commercial and fleet excluded unless it qualifies as a Permitted Commercial Purpose",
            "Platinum Plus cannot be sold with Wear Protection",
            "See agreement for full exclusion details",
        ],
        "transfer": "Transferable: $75 (private party sale or lease assumption only, not transferable if dealer is party to resale/assumption)",
        "cancel": STD_CANCEL_75,
        "circle": {
            "code": "HFCI",
            "benefit": "Requires Circle Certificate. Customer must also have an active CPO/CUV certification.",
            "restrictions": [
                "Retail rate restriction: maximum markup is $100 to dealer cost",
                "Requires Circle Certificate (and still must have active CPO / CUV certification)",
            ],
        },
    },
    {
        "name": "High Mileage Vehicle Service Protection (ICE)",
        "code": "WFHM", "core": False, "ev_care": False,
        "overview": "PowerProtect Powertrain-only coverage for used high-mileage ICE Hyundai and competitive make vehicles. Odometer between 60,001 and 175,000 miles. Available at time of sale only.",
        "plans": [
            {"name": "Powertrain", "sub": "(Stated Component)", "desc": "Engine, Transmission, Drive Axle incl. CV joints"},
        ],
        "term_coverage": [
            "Max options up to 5 years or 60,000 miles",
            "Additive coverage: begins day contract is sold",
            "Deductible options: $250 or $500",
        ],
        "elig_label": "Eligibility",
        "eligibility": [
            "Used ICE Hyundai vehicles AND competitive makes sold at Hyundai dealerships",
            "9 model years or less; odometer 60,001 to 175,000 miles",
            "Available at time of vehicle sale ONLY",
        ],
        "reimbursement": STD_REIMB,
        "additional_benefits": STD_AB,
        "exclusions": [
            "Vehicles under 60,001 or over 175,000 miles not eligible",
            "Post-sale purchase not permitted",
            "Total loss, salvaged, junk, buy-back vehicles not eligible",
            "See agreement for full exclusion details",
        ],
        "transfer": "Transferable: $75 (private party sale or lease assumption only)",
        "cancel": STD_CANCEL_75,
        "circle": None,
    },
    {
        "name": "EV Care Vehicle Service Protection",
        "code": "HFVE", "core": True, "ev_care": True,
        "overview": "Mechanical repair coverage for eligible Hyundai EV and Hybrid vehicles, parts and labor, up to 12 years or 200,000 miles.",
        "plans": [
            {"name": "Battery", "sub": "(Stated Cmpnt.)", "desc": "Lithium-Ion Battery Pack, Battery Module, Battery Management System, Battery Degradation (below 70%), High Voltage Pre Charger, Traction Motor (incl housing case), Onboard Charger, Inverter, Converter"},
            {"name": "High Technology", "sub": "(Stated Cmpnt.)", "desc": "Factory-installed ADAS, audio/visual systems, accessories (USB ports, 120V, wireless charging, WiFi/hotspot, phone interfaces)"},
            {"name": "Platinum", "sub": "(Exclusionary)", "desc": "All Covered Parts except those listed as exclusions"},
            {"name": "Platinum Plus", "sub": "(Florida Only)", "desc": "Platinum plus Term Protection items, headlamps, belts and hoses, electrical"},
        ],
        "term_coverage": [
            "Max options up to 12 years or 200,000 miles",
            "Additive coverage: begins day contract is sold, mileage based on odometer at sale",
            "Deductible options: $0 or $100",
        ],
        "elig_label": "EV VSP Eligibility",
        "eligibility": [
            "Vehicle type: Available only on EV and Hybrid, not available on ICE",
            "9 model years or less, under 140,001 miles at time of purchase",
            "In and Out of Warranty vehicles eligible at time of vehicle purchase",
            "After time of vehicle purchase: vehicle must have at least 1 month and 1,000 miles of manufacturer warranty remaining",
            "UCI required for any pre-owned",
        ],
        "reimbursement": STD_REIMB,
        "additional_benefits": STD_AB,
        "exclusions": [
            "Total loss, salvaged, junk, buy-back vehicles not eligible",
            "Non-U.S. specification vehicles not eligible",
            "Commercial and fleet excluded unless it qualifies as a Permitted Commercial Purpose",
            "See agreement for full exclusion details",
        ],
        "transfer": "Transferable: $75 (private party sale or lease assumption only, not transferable if dealer is party to resale/assumption)",
        "cancel": STD_CANCEL_75,
        "circle": {
            "code": "HFVE",
            "benefit": "Available for employees, family members, and affiliates.",
            "restrictions": [
                "Post-sale purchase requires 1 month warranty remaining and 1,000 miles full warranty remaining",
                "UCI required for any pre-owned",
                "Not available in FL",
                "Max markup: $100 over dealer cost",
            ],
        },
    },
    {
        "name": "EV Care VSP — Competitive Makes",
        "code": "WFVE", "core": False, "ev_care": True,
        "overview": "PowerProtect EV mechanical coverage for pre-owned competitive make EVs, parts and labor, up to 12 years or 120,000 miles.",
        "plans": [
            {"name": "Platinum", "sub": "(Exclusionary)", "desc": "All Covered Parts except those listed as exclusions"},
            {"name": "Platinum Plus", "sub": "(Florida Only)", "desc": "Platinum plus headlamps, belts and hoses, electrical"},
        ],
        "term_coverage": [
            "Max options up to 12 years or 120,000 miles",
            "Additive coverage: begins day contract is sold",
            "Deductible options: $0, $100, or $250",
        ],
        "elig_label": "EV VSP Eligibility",
        "eligibility": [
            "Pre-owned competitive make EVs only",
            "9 model years or less, under 108,001 miles at time of purchase",
            "UCI required for any pre-owned",
        ],
        "reimbursement": STD_REIMB,
        "additional_benefits": STD_AB,
        "exclusions": [
            "Hyundai EVs — use EV Care VSP (HFVE)",
            "Total loss, salvaged, junk, buy-back vehicles not eligible",
            "Non-U.S. specification vehicles not eligible",
            "See agreement for full exclusion details",
        ],
        "transfer": "Transferable: $75 (private party sale or lease assumption only)",
        "cancel": STD_CANCEL_75,
        "circle": None,
    },
    {
        "name": "EV Care VSP — Original Owner",
        "code": "HFOE", "core": True, "ev_care": True,
        "overview": "EV mechanical breakdown coverage for eligible Hyundai EVs and Hybrids, with coverage options up to 12 years/200,000 miles, measured from the original in-service date.",
        "plans": [
            {"name": "Battery", "sub": "(Stated Cmpnt.)", "desc": "Lithium-Ion Battery Pack, Battery Module, BMS, Battery Degradation (below 70%), HV Pre Charger, Traction Motor (incl housing case), Onboard Charger, Inverter, Converter"},
            {"name": "Platinum", "sub": "(Exclusionary)", "desc": "All Covered Parts except those listed as exclusions"},
            {"name": "Platinum Plus", "sub": "(Florida Only)", "desc": "Platinum plus Term Protection items, headlamps, belts and hoses, electrical"},
        ],
        "term_coverage": [
            "Coverage available up to 12 years or 200,000 miles from original in-service date",
            "In-service date coverage (measured from original in-service date)",
            "Vehicle mileage at sale: 10,001 to 58,999",
            "Deductible options: $0 or $100",
        ],
        "elig_label": "EV VSP Eligibility",
        "eligibility": [
            "Available only on EV and Hybrid, not available on ICE",
            "Original Owner requirement: currently owned by original owner or purchased by original lessee",
            "Vehicle mileage at sale: 10,001 to 58,999",
            "UCI required for any pre-owned",
        ],
        "reimbursement": STD_REIMB,
        "additional_benefits": STD_AB,
        "exclusions": [
            "Total loss, salvaged, junk, buy-back vehicles not eligible",
            "Non-U.S. specification vehicles not eligible",
            "Commercial and fleet excluded unless it qualifies as a Permitted Commercial Purpose",
            "See agreement for full exclusion details",
        ],
        "transfer": "Transferable: $75 (private party sale or lease assumption only, not transferable if dealer is party to resale/assumption)",
        "cancel": STD_CANCEL_75,
        "circle": None,
    },
    {
        "name": "EV Care Certified Used Vehicle (CUV) Wrap",
        "code": "HFCE", "core": True, "ev_care": True,
        "overview": "Wrap coverage for eligible Hyundai CUV certified EV vehicles, covering mechanical repairs, parts, and labor for the program term.",
        "plans": [
            {"name": "Battery", "sub": "(Stated Cmpnt.)", "desc": "Lithium-Ion Battery Pack, BMS, Battery degradation below 70%, HV Pre Charger, Traction Motor (incl housing case), Onboard Charger, Inverter, Converter"},
            {"name": "Gold", "sub": "(Stated Cmpnt.)", "desc": "Front/rear suspension (incl shocks), air conditioning, fuel system, electrical system"},
            {"name": "Platinum", "sub": "(Exclusionary)", "desc": "All Covered Parts except items listed as exclusions"},
            {"name": "Platinum Plus", "sub": "(Florida Only)", "desc": "Platinum plus headlamp, belts and hoses, electrical"},
        ],
        "term_coverage": [
            "Max term and mileage: up to 12 years or 200,000 miles from original in-service date",
            "Coverage begins from original in-service date and zero miles",
            "Deductible options: $0, $50 (disappearing at selling dealer), or $100",
        ],
        "elig_label": "EV CPOW Eligibility",
        "eligibility": [
            "Hyundai CUV Vehicle / EV Only",
            "Less than 6 years from the original in-service date and less than 80,000 miles",
            "Customer must have an active CPO / CUV certification to qualify",
            "UCI required for any pre-owned",
        ],
        "reimbursement": STD_REIMB,
        "additional_benefits": STD_AB,
        "exclusions": [
            "Vehicle cannot exceed 80,000 miles",
            "Commercial and fleet excluded unless it qualifies as a Permitted Commercial Purpose",
            "Platinum Plus cannot be sold with Wear Protection",
            "See agreement for full exclusion details",
        ],
        "transfer": "Transferable: $75 (private party sale or lease assumption only, not transferable if dealer is party to resale/assumption)",
        "cancel": STD_CANCEL_75,
        "circle": {
            "code": "HFCE",
            "benefit": "Requires Circle Certificate. Customer must also have an active CPO/CUV certification.",
            "restrictions": [
                "Max markup: $100 over dealer cost",
                "Requires Circle Certificate (and still must have active CPO / CUV certification)",
            ],
        },
    },
    {
        "name": "VSP Livery Coverage (ICE & EV)",
        "code": "HCVL", "core": True, "ev_care": False,
        "overview": "Mechanical breakdown coverage designed for livery use (limousines, chauffeured vehicles, private hire, shuttles). Available through select dealers only — contact your District Manager of Insurance.",
        "plans": [
            {"name": "ICE — Powertrain", "sub": "(Stated Component)", "desc": "Engine, Transmission, Drive Axle incl. CV joints"},
            {"name": "ICE — Gold", "sub": "(Stated Cmpnt.)", "desc": "Powertrain plus front/rear suspension (incl shocks), A/C, fuel system, electrical system"},
            {"name": "ICE — Platinum", "sub": "(Exclusionary)", "desc": "All Covered Parts except items listed in Exclusions"},
            {"name": "EV — Battery", "sub": "(Stated Cmpnt.)", "desc": "Lithium-Ion Battery Pack, BMS, Battery Degradation (<70%), HV Pre Charger, Traction Motor, Onboard Charger, Inverter, Converter"},
            {"name": "EV — Platinum", "sub": "(Exclusionary)", "desc": "All Covered Parts except items listed in Exclusions"},
        ],
        "term_coverage": [
            "Coverage available from 5 years/100,000 miles up to 5 years/300,000 miles",
            "Deductible: $0",
        ],
        "elig_label": "VSP Eligibility",
        "eligibility": [
            "New ICE and EV Hyundai vehicles only",
            "Vehicle miles cannot exceed 10,000 miles at sale",
            "Available at time of vehicle sale ONLY",
            "Only available to certain dealers — contact District Manager of Insurance",
            "Livery Use: limousine, chauffeured vehicle, private hire vehicle, or shuttle services",
        ],
        "reimbursement": STD_REIMB,
        "additional_benefits": STD_AB,
        "exclusions": [
            "Pre-owned vehicles not eligible",
            "Vehicles over 10,000 miles at purchase not eligible",
            "Total loss, salvaged, junk, buy-back vehicles not eligible",
            "Non-livery commercial use not permitted",
        ],
        "transfer": "Transferable: $75 (private party sale or lease assumption only)",
        "cancel": STD_CANCEL_75,
        "circle": None,
    },
    {
        "name": "Wear Protection (ICE & EV)",
        "code": "HFWP", "core": True, "ev_care": False,
        "overview": "Additive coverage plan for new Hyundai vehicles that pays for common wear items — brakes, 12V battery, alignment, lighting, belts, hoses, wipers — during the selected term.",
        "plans": [
            {"name": "Standard", "sub": "", "desc": "One set front & rear brake pads/shoes, 12V battery replacement (1), wheel alignment (1), unlimited headlamps (non-impact), unlimited engine belts & hoses, one set wiper blades, unlimited fuses & light bulbs (non-impact)"},
            {"name": "+ Brake Rotor", "sub": "(Optional Add-on)", "desc": "Adds one replacement set of front and rear brake rotors during the term"},
        ],
        "term_coverage": [
            "Max term and mileage: Up to 60 months / 60,000 miles",
            "Deductible: $0",
            "Available for purchase at time of sale only",
        ],
        "elig_label": "Eligibility",
        "eligibility": [
            "New Hyundai ICE, EV, Hybrid, and Hydrogen Fuel Cell vehicles",
            "Pre-owned ICE and EV competitive makes (max 10,000 miles at purchase)",
            "Vehicle mileage cannot exceed 10,000 at time of purchase",
            "Available at time of sale ONLY",
        ],
        "reimbursement": STD_REIMB,
        "additional_benefits": STD_AB,
        "exclusions": [
            "Cannot be sold with Hyundai Maintenance Wrap Plus (HCM)",
            "Hybrid & EV batteries excluded from 12V battery benefit",
            "Impact-damaged headlamps, bulbs not covered",
            "Platinum Plus (FL) cannot be sold with Wear Protection",
            "Competitive makes on new vehicles not eligible",
        ],
        "transfer": "Transferable: $75 (private party sale or lease assumption only, not transferable if dealer is party to resale/assumption)",
        "cancel": STD_CANCEL_75,
        "circle": {
            "code": "HFEW",
            "benefit": "Available for employees, family members, and affiliates.",
            "restrictions": [
                "Dealer must verify customer Circle status",
                "Max markup: $100 over dealer cost",
            ],
        },
    },

    # ════════════════ MAINTENANCE ════════════════
    {
        "name": "Maintenance Basic Wrap (ICE)",
        "code": "HFBI", "core": True, "ev_care": False,
        "overview": "Extends/wraps Hyundai Complimentary Maintenance (HCM) on MY2025 and older ICE vehicles. Mirrors HCM services: oil & oil filter, tire rotation, and multi-point inspection. Option to upgrade to Severe Usage intervals.",
        "plans": [
            {"name": "Basic Wrap", "sub": "(HFBI)", "desc": "Oil & oil filter change, tire rotation, multi-point inspection per interval. Severe Usage upgrade available. Extends beyond standard 3 yrs/36,000 mi HCM coverage."},
        ],
        "term_coverage": [
            "Coverage for up to 8 years/96,000 miles",
            "MY2025 and older ICE Hyundai vehicles with HCM only",
            "Deductible: $0",
        ],
        "elig_label": "Eligibility",
        "eligibility": [
            "MY2025 and older ICE Hyundai vehicles WITH HCM only",
            "Available at time of sale or before 6,000 miles on odometer",
            "Pre-Paid Maintenance NOT eligible on any vehicle with HCM",
            "Electric vehicles excluded",
        ],
        "reimbursement": [("", [("Reimbursed per Maintenance Interval sheets", False)])],
        "additional_benefits": [[("Non-Transferable", True)]],
        "exclusions": [
            "EV vehicles",
            "MY2026+ vehicles — use Pre-Paid Maintenance (HFPI)",
            "Services not performed within 5,000 miles or 6 months of scheduled interval",
            "Mechanical breakdown repairs not covered",
        ],
        "transfer": "Non-Transferable to subsequent owner/lessee or another vehicle",
        "cancel": [
            [("Within 30 days: full refund if no covered services provided", False)],
            [("Next 35 months: full refund less ", False), ("$50", True), (" cancellation fee if no covered services provided", False)],
            [("After initial look period or if covered service provided: pro-rata less cost of services + ", False), ("$50", True), (" fee", False)],
        ],
        "circle": None,
    },
    {
        "name": "Maintenance Scheduled Wrap (ICE)",
        "code": "HFSI", "core": True, "ev_care": False,
        "overview": "Extends/wraps Hyundai Complimentary Maintenance (HCM) on MY2025 and older ICE vehicles. Includes all Basic Wrap services plus owner's manual-scheduled items: air filters, brake fluid, spark plugs, transmission fluid, and other fluids.",
        "plans": [
            {"name": "Scheduled Wrap", "sub": "(HFSI)", "desc": "All Basic Wrap services (oil & filter, tire rotation, multi-point inspection) plus cabin air filter, engine air filter, brake fluid, spark plugs, transmission fluid, differential/case oil per owner's manual schedule. Severe Usage upgrade available."},
        ],
        "term_coverage": [
            "Coverage for up to 8 years/96,000 miles",
            "MY2025 and older ICE Hyundai vehicles with HCM only",
            "Deductible: $0",
        ],
        "elig_label": "Eligibility",
        "eligibility": [
            "MY2025 and older ICE Hyundai vehicles WITH HCM only",
            "Available at time of sale or before 6,000 miles on odometer",
            "Pre-Paid Maintenance NOT eligible on any vehicle with HCM",
            "Electric vehicles excluded",
        ],
        "reimbursement": [("", [("Reimbursed per Maintenance Interval sheets", False)])],
        "additional_benefits": [[("Non-Transferable", True)]],
        "exclusions": [
            "EV vehicles",
            "MY2026+ vehicles — use Pre-Paid Maintenance (HFSI/Scheduled PPM)",
            "Services not performed within 5,000 miles or 6 months of scheduled interval",
            "Mechanical breakdown repairs not covered",
        ],
        "transfer": "Non-Transferable to subsequent owner/lessee or another vehicle",
        "cancel": [
            [("Within 30 days: full refund if no covered services provided", False)],
            [("Next 35 months: full refund less ", False), ("$50", True), (" cancellation fee if no covered services provided", False)],
            [("After initial look period or if covered service provided: pro-rata less cost of services + ", False), ("$50", True), (" fee", False)],
        ],
        "circle": None,
    },
    {
        "name": "Pre-Paid Maintenance (ICE)",
        "code": "HFPI / HFSI", "core": True, "ev_care": False,
        "overview": "Prepaid maintenance services at set intervals, reimbursed per interval schedule, designed to keep customers on factory-recommended service cadence. Basic available anytime; Scheduled on new MY2026+ at time of sale.",
        "plans": [
            {"name": "Basic Maintenance", "sub": "(HFPI)", "desc": "Oil & oil filter services at selected intervals; tire rotations and multi-point inspection. Normal or Severe usage options. All model year Hyundai vehicles."},
            {"name": "Scheduled Maintenance", "sub": "(HFSI)", "desc": "Includes Basic services plus scheduled items: cabin filter, engine air filter, brake fluid, spark plugs, transmission fluid, differential/case oil. New MY2026+ only; at time of sale; ≤6,000 miles."},
        ],
        "term_coverage": [
            "Coverage up to 8 years (96 months), 96,000 miles",
            "Deductible: $0",
            "Basic available anytime; Scheduled at time of sale only on new MY2026+ vehicles",
        ],
        "elig_label": "PPM Eligibility",
        "eligibility": [
            "Only available on ICE & Hybrid Vehicles",
            "Basic: any vehicle without HCM, any time",
            "Scheduled: new MY2026+ vehicles only; max 6,000 miles odometer; at time of sale only",
            "NOT available for any vehicle with HCM",
        ],
        "reimbursement": [("", [("Reimbursed per Maintenance Interval sheets", False)])],
        "additional_benefits": [
            [("Locks in today's service prices; may be amortized within retail installment contract or lease", False)],
            [("Helps increase customer retention and loyalty", False)],
        ],
        "exclusions": [
            "Vehicles with HCM — use Maintenance Wrap instead",
            "EV vehicles — use EV Care Maintenance",
            "Mechanical breakdown repairs not covered",
            "Services not completed within interval requirements",
        ],
        "transfer": "Transferable: $50 (private party or lease assumption only, within 30 days of transfer event)",
        "cancel": STD_CANCEL_75,
        "circle": {
            "code": "HFEP",
            "benefit": "Friends and Family certificate required. Covers employees, family members, and affiliates.",
            "restrictions": [
                "Max markup: $100 over dealer cost",
            ],
        },
    },
    {
        "name": "Pre-Paid Maintenance — Competitive Makes (ICE)",
        "code": "WFPI", "core": False, "ev_care": False,
        "overview": "Prepaid maintenance for competitive make ICE vehicles sold at Hyundai dealerships. Three plans: Basic (oil only), Plus (oil + rotation), Synthetic Plus (synthetic oil + rotation). Sold at any time.",
        "plans": [
            {"name": "Basic", "sub": "", "desc": "Conventional oil & oil filter changes at selected intervals"},
            {"name": "Plus", "sub": "", "desc": "Conventional oil changes with tire rotations for the covered vehicle"},
            {"name": "Synthetic Plus", "sub": "", "desc": "Synthetic oil changes with tire rotations for the covered vehicle"},
        ],
        "term_coverage": [
            "Coverage up to 7 years / 105,000 miles",
            "Deductible: $0",
            "Can be sold at any time (no time-of-sale restriction)",
        ],
        "elig_label": "Eligibility",
        "eligibility": [
            "ICE competitive make (non-Hyundai) vehicles sold at Hyundai dealerships",
            "Can be sold at any time",
            "Commercial/fleet excluded unless permitted professional use",
        ],
        "reimbursement": [("", [("Reimbursed per Maintenance Interval sheets", False)])],
        "additional_benefits": [
            [("Locks in today's service prices; may be amortized within retail installment contract or lease", False)],
        ],
        "exclusions": [
            "Hyundai vehicles — use HPP Pre-Paid Maintenance",
            "EV / Hybrid competitive makes not eligible",
        ],
        "transfer": "Transferable: $50 (private party or lease assumption only)",
        "cancel": STD_CANCEL_75,
        "circle": None,
    },
    {
        "name": "EV Care Maintenance",
        "code": "HFEM", "core": True, "ev_care": True,
        "overview": "Prepaid EV maintenance program designed to keep the customer's EV running smoothly, reimbursed per maintenance interval schedule. Basic anytime; Maintenance Plus on new MY2026+ at time of sale.",
        "plans": [
            {"name": "Basic", "sub": "(New + Used)", "desc": "Tire rotations and multi-point maintenance inspection every interval; cabin air filter replacement every other interval"},
            {"name": "Maintenance Plus", "sub": "(New MY2026+ Only)", "desc": "All Basic items + all owner's manual scheduled services + wiper blades (every other interval), wheel alignment at 32K & 72K mi, 12V battery at 72K mi, brake pads at 72K mi, brake fluid at 48K & 96K mi"},
        ],
        "term_coverage": [
            "Up to 8 years (96 months), 96,000 miles",
            "Deductible: $0",
            "Basic: sold at any time; Maintenance Plus: new MY2026+ EV only at time of sale; ≤6,000 miles",
        ],
        "elig_label": "EV Care Maintenance Eligibility",
        "eligibility": [
            "Basic: all model year Hyundai EV and Fuel Cell vehicles",
            "Maintenance Plus: new MY2026+ Hyundai EV only; ≤3 years old; ≤6,000 miles",
            "Not available on vehicles with HCM (Maintenance Plus only)",
            "Basic can be sold at any time; Maintenance Plus at time of sale only",
        ],
        "reimbursement": [("", [("Reimbursed per Maintenance Interval sheets", False)])],
        "additional_benefits": [
            [("Plug-in Hyundai Hybrid vehicles are NOT eligible for EV Care at this time", False)],
        ],
        "exclusions": [
            "Plug-in Hybrid vehicles not eligible",
            "ICE vehicles — use Pre-Paid Maintenance",
            "Maintenance Plus not for vehicles with HCM",
        ],
        "transfer": "Transferable: $50 (private party or lease assumption only, within 30 days of transfer event)",
        "cancel": STD_CANCEL_75,
        "circle": {
            "code": "HFEM",
            "benefit": "Friends and Family certificate required. Covers employees, family members, and affiliates.",
            "restrictions": [
                "Max markup: $100 over dealer cost",
            ],
        },
    },

    # ════════════════ DAMAGE CARE ════════════════
    {
        "name": "Lease-End Protection",
        "code": "HFLP", "core": True, "ev_care": False,
        "overview": "Helps cover excess wear and use charges assessed at lease end, reducing the customer's out-of-pocket exposure. Available only on leases assigned to Hyundai Motor Finance / Hyundai Lease Titling Trust.",
        "plans": [
            {"name": "Single Plan", "sub": "", "desc": "Covered excess wear and use charges waived by Hyundai. Turn-in window: any time prior to, or within 1 year after, the original scheduled termination date. Unlimited waiver benefit; single event limit $1,000."},
        ],
        "term_coverage": [
            "Coverage ranging from 12 to 72 months, up to 150,000 miles",
            "Lease term MUST match the protection term",
            "Deductible: $0",
        ],
        "elig_label": "Lease Protection Eligibility",
        "eligibility": [
            "New Hyundai vehicle only (ICE, EV, hybrid, hydrogen)",
            "ICE, EV, hybrid, hydrogen eligible if within mileage restriction",
            "Leased vehicles only",
            "Must be assigned to HMF / Hyundai Lease Titling Trust",
            "Available at time of lease only",
        ],
        "reimbursement": [
            ("", [("Charges waived directly by Hyundai Lease Titling Trust / Hyundai Motor Finance", False)]),
            ("", [("Missing parts: up to ", False), ("$150 total", True), (" ($100 with Circle Certificate)", False)]),
            ("", [("Excess miles: up to 200 miles waived at ", False), ("$0.20 per mile", True)]),
        ],
        "additional_benefits": [
            [("Paint damage & exterior dents", False)],
            [("Glass chips/breakage, lights & lenses", False)],
            [("Interior rips/tears/stains, carpet damage", False)],
            [("Wheels/wheel covers, mirrors, trim/moldings, bumper covers", False)],
            [("Worn tires; Missing parts up to $150 total", False)],
        ],
        "exclusions": [
            "Not available in NY or MD",
            "Product is only for leases assigned to HMF / Hyundai Lease Titling Trust",
            "Damage exceeding benefit limits and $1,000 single event cap",
            "Alteration charges or pre-turn-in completed repairs",
            "No benefit if vehicle is purchased by customer or dealer",
        ],
        "transfer": "Transferable: $75 (to subsequent lessee only)",
        "cancel": [
            [("Within 30 days: full refund less claims paid", False)],
            [("After 30 days: pro-rata less claims, ", False), ("$75", True), (" processing fee", False)],
            [("Once waiver benefit is paid: contract deemed fully earned — no cancellation", False)],
        ],
        "circle": {
            "code": "HFEL",
            "benefit": "Friends and Family certificate required. Missing parts limit changes to $100 total (base is $150).",
            "restrictions": [
                "Max markup: $100 over dealer cost",
            ],
        },
    },
    {
        "name": "Multi-Coverage Protection",
        "code": "HFMC", "core": False, "ev_care": False,
        "overview": "Bundled coverage package combining Tire & Wheel, Dent, Windshield, and Key Protection into one contract. Products can be sold as a bundle or as stand-alone on the Multi-Coverage form.",
        "plans": [
            {"name": "Platinum", "sub": "(Complete)", "desc": "Tire & Wheel + Dent + Key + Windshield"},
            {"name": "Gold", "sub": "(No Key)", "desc": "Tire & Wheel + Dent + Windshield"},
            {"name": "Silver", "sub": "(FL Bundle Only)", "desc": "Tire & Wheel + Dent + Key"},
            {"name": "Bronze", "sub": "(Entry)", "desc": "Tire & Wheel + Dent"},
        ],
        "term_coverage": [
            "Up to 7 years (84 months), unlimited miles",
            "Deductible: $0",
            "No limit on occurrences (except hail damage & cosmetic wheel repairs)",
            "Available at time of vehicle purchase only",
        ],
        "elig_label": "PVP Eligibility",
        "eligibility": [
            "Available on New, Used, and CUV Hyundai vehicles",
            "Also available on Used Competitive Make vehicles (incl. Kia, Genesis)",
            "No mileage or year limitations",
            "Available at time of vehicle purchase only",
        ],
        "reimbursement": [
            ("", [("Reimbursement follows individual product agreement caps", False)]),
        ],
        "additional_benefits": [
            [("Windshield Protection NOT available stand-alone or in bundle in FL", False)],
            [("Cancellation of one product cancels ALL products on the contract", False)],
        ],
        "exclusions": [
            "Individual product exclusions apply (see underlying product agreements)",
            "Windshield stand-alone not available in FL",
            "Hail damage limited to $1,000 per occurrence",
        ],
        "transfer": "Transferable: $50 (private party or lease assumption only; transfer request within 30 days)",
        "cancel": STD_CANCEL_50,
        "circle": None,
    },
    {
        "name": "Appearance Protection",
        "code": "HFAP", "core": False, "ev_care": False,
        "overview": "Permanently applied interior and exterior protection system covering paint, fabric, vinyl, and leather from stains, finish damage, and environmental conditions. NON-CANCELLABLE.",
        "plans": [
            {"name": "Exterior Paint", "sub": "", "desc": "Protects factory painted surfaces from UV oxidation/fading, corrosion (sea/road salt, acid rain, hard water, insects, de-icing), bird droppings, tree sap, spray paint/overspray"},
            {"name": "Interior Fabric", "sub": "", "desc": "Covers fabric seating and interior surfaces from stains: food/drink (soda, coffee, gum, baby food), pet waste, makeup, dye transfer, ink, crayons"},
            {"name": "Interior Vinyl & Leather", "sub": "", "desc": "Covers stains or damage to vinyl/leather surfaces — benefit is professional cleaning, repair, dye, or replacement of affected area"},
        ],
        "term_coverage": [
            "Max term: 120 months (10 years) / Unlimited miles",
            "Deductible: $0",
            "Permanently applied at time of sale",
        ],
        "elig_label": "Appearance Eligibility",
        "eligibility": [
            "New and pre-owned Hyundai vehicles (ICE, EV, hybrid, hydrogen)",
            "Competitive makes sold at HPP dealership",
            "Pre-owned vehicles: 5 model years old or newer",
            "Available at time of sale only",
        ],
        "reimbursement": [
            ("", [("Repair, repaint, or refinish of affected area (per agreement limits)", False)]),
            ("", [("Professional cleaning, repair, or replacement for interior (per agreement limits)", False)]),
        ],
        "additional_benefits": [
            [("Rental Car Benefit: up to ", False), ("$50/day", True), (", 10 days max while vehicle is in for a covered repair", False)],
        ],
        "exclusions": [
            "Pre-owned vehicles older than 5 model years not eligible",
            "Commercial and fleet vehicles excluded (permitted professional use allowed)",
        ],
        "transfer": "Transferable: $50 (private party; within 30 days with required documents)",
        "cancel": [[("NON-CANCELLABLE", True)]],
        "circle": None,
    },

    # ════════════════ LOSS & THEFT ════════════════
    {
        "name": "Guaranteed Asset Protection (GAP)",
        "code": "HFGP / WFGP", "core": True, "ev_care": False,
        "overview": "Waives the difference between the vehicle's cash value (per primary insurer) and the remaining balance on the customer's finance agreement in the event of a total loss. Finance purchases only.",
        "plans": [
            {"name": "GAP", "sub": "(HFGP — Hyundai | WFGP — Off-Make)", "desc": "Waives the gap between vehicle cash value and remaining finance balance. Also waives up to $1,000 of customer's primary insurance deductible (not in all states)."},
        ],
        "term_coverage": [
            "Max term: up to 96 months",
            "No mileage limitations",
            "No LTV limit; No MSRP limitations",
            "Deductible: $0",
        ],
        "elig_label": "GAP Eligibility",
        "eligibility": [
            "Finance purchases only (NOT available on leases)",
            "New, Used, CUV Hyundai vehicles and Used Competitive Makes",
            "Available for purchase at time of sale only",
            "NOT available in CA, NY, TX",
        ],
        "reimbursement": [
            ("", [("Waives gap between vehicle cash value and remaining finance balance", False)]),
            ("", [("Also waives up to ", False), ("$1,000", True), (" of customer's primary insurance deductible (not available in all states)", False)]),
        ],
        "additional_benefits": [
            [("Competitive makes eligible (Off-Make product code: WFGP — NON-CORE)", False)],
            [("Deductible waiver: up to $1,000 (not available in all states)", False)],
        ],
        "exclusions": [
            "Leases not eligible",
            "Amounts owed due to missed payments or extraordinary events not waived",
            "NOT available in CA, NY, TX",
            "GAP deductible waiver not available in all states",
        ],
        "transfer": "Transferable: $75 (only upon assumption of retail installment contract)",
        "cancel": [
            [("Within 30 days: full refund less claims paid", False)],
            [("After 30 days: pro-rata less claims, ", False), ("$75", True), (" cancellation fee", False)],
            [("Once waiver benefit is paid: addendum deemed fully earned — no cancellation", False)],
        ],
        "circle": {
            "code": "HFEG",
            "benefit": "Customer cost fixed at $100 over dealer cost. Requires Circle Certificate verification.",
            "restrictions": [
                "Max markup: $100 over dealer cost (fixed)",
                "Circle code changes to HFEG",
            ],
        },
    },
    {
        "name": "Guaranteed Asset Protection Plus (GAP Plus)",
        "code": "HFGS / WFGS", "core": True, "ev_care": False,
        "overview": "All GAP benefits plus a $2,000 credit toward a replacement vehicle purchase at the original selling dealer (or nearest Hyundai dealer if customer moves >50 miles). Finance only.",
        "plans": [
            {"name": "GAP Plus", "sub": "(HFGS — Hyundai | WFGS — Off-Make)", "desc": "All GAP benefits + $2,000 credit toward purchase or lease of a replacement vehicle. Available only at the selling dealer unless customer moves more than 50 miles away."},
        ],
        "term_coverage": [
            "Max term: up to 96 months",
            "No mileage or LTV limitations",
            "Deductible: $0",
            "NOT available in AK, KS, NE, NY, TN, TX, WV, CA",
        ],
        "elig_label": "GAP Plus Eligibility",
        "eligibility": [
            "Finance purchases only (NOT available on leases)",
            "New, Used, CUV Hyundai vehicles and Used Competitive Makes",
            "Available for purchase at time of sale only",
            "NOT available in AK, KS, NE, NY, TN, TX, WV, CA",
        ],
        "reimbursement": [
            ("", [("All standard GAP benefits (gap waiver + up to $1,000 deductible waiver)", False)]),
            ("", [("GAP Plus: ", False), ("$2,000", True), (" replacement vehicle credit at original selling dealer or nearest Hyundai dealer", False)]),
        ],
        "additional_benefits": [
            [("$2,000 replacement credit available at original selling dealer or nearest Hyundai dealer (if customer moves >50 miles)", False)],
            [("Competitive makes eligible (Off-Make product code: WFGS — NON-CORE)", False)],
        ],
        "exclusions": [
            "NOT available in AK, KS, NE, NY, TN, TX, WV, CA",
            "Leases not eligible",
            "GAP Plus credit only applicable at a Hyundai dealership",
        ],
        "transfer": "Transferable: $75 (only upon assumption of retail installment contract)",
        "cancel": [
            [("Within 30 days: full refund less claims paid", False)],
            [("After 30 days: pro-rata less claims, ", False), ("$75", True), (" cancellation fee", False)],
            [("Once waiver benefit is paid: addendum deemed fully earned — no cancellation", False)],
        ],
        "circle": {
            "code": "HFEG",
            "benefit": "Customer cost fixed at $100 over dealer cost. Requires Circle Certificate verification.",
            "restrictions": [
                "Max markup: $100 over dealer cost (fixed)",
                "Circle code changes to HFEG for GAP Plus as well",
            ],
        },
    },
    {
        "name": "Theft Protection",
        "code": "HFTP", "core": False, "ev_care": False,
        "overview": "Permanently marks the vehicle with VIN identification codes and anti-theft window decals to deter theft. Provides a financial warranty benefit up to $5,000 if the vehicle is stolen and declared a total loss. NON-CANCELLABLE.",
        "plans": [
            {"name": "Theft Protection", "sub": "", "desc": "Permanent VIN markings + anti-theft decals affixed to vehicle windows. No activation required. Warranty benefit up to $5,000 if vehicle is stolen and not recovered or declared a total loss."},
        ],
        "term_coverage": [
            "Max term: 60 months (5 years), unlimited miles",
            "Deductible: $0",
            "No mileage limitations",
        ],
        "elig_label": "Theft Eligibility",
        "eligibility": [
            "New & pre-owned Hyundai vehicles (ICE, EV, Hybrid, Hydrogen)",
            "Competitive makes including Kia & Genesis sold at Hyundai dealerships",
            "Available at time of vehicle sale or lease only",
            "CT: dealer must have written consent before VIN etching",
        ],
        "reimbursement": [
            ("", [("Warranty benefit up to ", False), ("$5,000", True), (" — checks payable to customer, sent to dealership", False)]),
            ("", [("Benefit covers: insurance deductible, sales tax, tag & title fees, negative equity, down payment for replacement vehicle", False)]),
        ],
        "additional_benefits": [
            [("No activation required — vehicle is permanently marked at time of sale", False)],
        ],
        "exclusions": [
            "Commercial and fleet vehicles excluded (permitted professional use allowed)",
            "FL: flat warranty benefit capped at $5,000 (no check-off box benefit amounts)",
            "CT: written consent required before etching",
        ],
        "transfer": "Transferable: $50 (private party; within 30 days with required documents)",
        "cancel": [[("NON-CANCELLABLE", True)]],
        "circle": None,
    },
]


def main():
    prs = Presentation()
    prs.slide_width  = Inches(13.3333)
    prs.slide_height = Inches(7.5)

    # Cover
    build_cover(prs)

    # TOC
    build_toc(prs)

    # ── MECHANICAL ────────────────────────────────────────────────────────
    build_section_divider(prs, "MECHANICAL COVERAGE",
        "Vehicle Service Protection  |  EV Care  |  Wear Protection")
    mech = [
        "Vehicle Service Protection (ICE)",
        "Vehicle Service Protection (ICE) — Competitive Makes",
        "Vehicle Service Protection Original Owner (ICE)",
        "Certified Used Vehicle (CUV) Wrap",
        "High Mileage Vehicle Service Protection (ICE)",
        "EV Care Vehicle Service Protection",
        "EV Care VSP — Competitive Makes",
        "EV Care VSP — Original Owner",
        "EV Care Certified Used Vehicle (CUV) Wrap",
        "VSP Livery Coverage (ICE & EV)",
        "Wear Protection (ICE & EV)",
    ]
    for d in PRODUCTS:
        if d["name"] in mech:
            build_product_slide(prs, d)

    # ── MAINTENANCE ───────────────────────────────────────────────────────
    build_section_divider(prs, "MAINTENANCE PLANS",
        "ICE  |  EV Care  |  Competitive Makes",
        bg=RGBColor(0x00, 0x50, 0x90))
    maint = [
        "Maintenance Basic Wrap (ICE)",
        "Maintenance Scheduled Wrap (ICE)",
        "Pre-Paid Maintenance (ICE)",
        "Pre-Paid Maintenance — Competitive Makes (ICE)",
        "EV Care Maintenance",
    ]
    for d in PRODUCTS:
        if d["name"] in maint:
            build_product_slide(prs, d)

    # ── DAMAGE CARE ───────────────────────────────────────────────────────
    build_section_divider(prs, "DAMAGE CARE",
        "Lease-End  |  Multi-Coverage Protection  |  Appearance Protection",
        bg=C_EV_TEAL)
    damage = [
        "Lease-End Protection",
        "Multi-Coverage Protection",
        "Appearance Protection",
    ]
    for d in PRODUCTS:
        if d["name"] in damage:
            build_product_slide(prs, d)

    # ── LOSS & THEFT ──────────────────────────────────────────────────────
    build_section_divider(prs, "LOSS & THEFT PROTECTION",
        "GAP  |  GAP Plus  |  Theft Protection",
        bg=RGBColor(0x6A, 0x1A, 0x4D))
    loss = [
        "Guaranteed Asset Protection (GAP)",
        "Guaranteed Asset Protection Plus (GAP Plus)",
        "Theft Protection",
    ]
    for d in PRODUCTS:
        if d["name"] in loss:
            build_product_slide(prs, d)

    out = "/Users/justin/Downloads/HPP Product Framework V2.pptx"
    prs.save(out)
    print(f"Saved: {out}")
    print(f"Total slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
